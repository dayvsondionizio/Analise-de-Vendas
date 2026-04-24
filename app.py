"""
BI Fiscal — Análise Estratégica de Vendas (Panificação)
Lê os Excels exportados do conversor de XML e gera relatório completo.
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from itertools import combinations
from collections import Counter
import io
import re
import warnings
from pathlib import Path

warnings.filterwarnings("ignore")

st.set_page_config(
    page_title="Análise de Vendas CP",
    page_icon="",
    layout="wide",
    initial_sidebar_state="expanded",
)

# 
# CATEGORIAS — mapeamento por palavra-chave no xProd
# Ordem importa: primeiro match ganha
# 
# ──────────────────────────────────────────────────────────────────
# CATEGORIAS — baseadas em NCM (código fiscal oficial)
# NCM é sempre primário; keywords são fallback só para alimentos
# preparados na própria padaria (salgados, bolos, etc.) que dividem
# capítulo NCM 19/21 com produtos industrializados.
# ──────────────────────────────────────────────────────────────────

# Subcapítulos 4 dígitos têm prioridade sobre capítulos 2 dígitos
MAPA_NCM = {
    # ── subcapítulos 4 dígitos ──────────────────────────────────
    "1901": "Mercearia",       # misturas p/ bolos, preparações infantis
    "1902": "Massas",          # massas alimentícias (macarrão, lamen…)
    "1903": "Mercearia",       # tapioca granulada
    "1904": "Mercearia",       # cereais pré-cozidos (granola, cornflakes)
    "1905": "Panificação",     # pão, biscoito, wafer, torrada, croissant
    # ── capítulos 2 dígitos ─────────────────────────────────────
    "01": "Outros",
    "02": "Frios/Carnes",      # carnes bovinas, suínas, aves
    "03": "Frios/Carnes",      # peixes e frutos do mar
    "04": "Laticínios",        # leite, queijo, manteiga, ovos, iogurte
    "07": "Hortifruti",        # legumes e hortaliças
    "08": "Hortifruti",        # frutas
    "09": "Café/Chá",          # café, chá, mate, especiarias
    "10": "Mercearia",         # cereais (trigo, arroz, milho)
    "11": "Mercearia",         # farinhas e amidos
    "12": "Mercearia",         # sementes oleaginosas
    "13": "Mercearia",         # gomas e resinas vegetais
    "15": "Mercearia",         # gorduras e óleos vegetais/animais
    "16": "Frios/Carnes",      # conservas de carne e peixe
    "17": "Doces",             # açúcar, mel, balas, confeitos
    "18": "Doces",             # cacau e chocolate
    "19": "Panificação",       # demais preparações de cereais
    "20": "Mercearia",         # preparações de legumes/frutas (polpa, geleia)
    "21": "Mercearia",         # preparações alimentícias (tempero, caldo)
    "22": "Bebidas",           # bebidas (água, suco, refri, cerveja, vinho)
    "23": "Mercearia",         # resíduos da indústria alimentar
    "33": "Higiene/Limpeza",   # perfumaria, cosméticos
    "34": "Higiene/Limpeza",   # sabões, detergentes, velas
    "35": "Mercearia",         # albuminas e gelatinas
    "39": "Embalagens",        # plásticos e embalagens plásticas
    "48": "Embalagens",        # papel e embalagens de papel/papelão
    "63": "Outros",            # artigos têxteis
    "94": "Outros",            # móveis e iluminação
}

# Keywords: fallback APENAS para alimentos preparados na padaria
# (estes têm NCM 19/21 igual aos industrializados — NCM não distingue)
MAPA_CATEGORIAS_FALLBACK = [
    ("Salgados",   ["COXINHA", "PASTEL", "SALGADO", "ESFIHA", "KIBE",
                    "RISOLIS", "RISOLE", "EMPADA", "BOLINHA", "DOGUINHO",
                    "QUICHE", "CROQUET"]),
    ("Bolos",      ["BOLO", "CHEESECAKE", "BROWNIE", "CUPCAKE", "ROCAMBOLE",
                    "PANETONE", "CHOCOTONE", "PUDIM", "MOUSSE", "SEMIFRIO"]),
    ("Tapiocas",   ["TAPIOCA"]),
    ("Refeições",  ["ALMOCO", "ALMOÇO", "PRATO FEITO", "MARMITA",
                    "SELF SERVICE", "SELF-SERVICE", "BUFFET"]),
    ("Lanchonete", ["LANCHONETE", "SANDUICHE", "SANDUÍCHE", "WRAP"]),
    ("Sopas/Caldos", ["SOPAO", "SOPÃO", "MINGAU"]),
]

DAYS_MAP = {
    "Monday": "Segunda", "Tuesday": "Terça", "Wednesday": "Quarta",
    "Thursday": "Quinta", "Friday": "Sexta",
    "Saturday": "Sábado", "Sunday": "Domingo",
}

CORES_CATEGORIA = {
    "Panificação":      "#E67E22",
    "Massas":           "#E8A838",
    "Laticínios":       "#5DADE2",
    "Frios/Carnes":     "#3498DB",
    "Bolos":            "#9B59B6",
    "Salgados":         "#F1C40F",
    "Tapiocas":         "#1ABC9C",
    "Café/Chá":         "#795548",
    "Refeições":        "#27AE60",
    "Sopas/Caldos":     "#5D6D7E",
    "Bebidas":          "#2980B9",
    "Doces":            "#E91E63",
    "Hortifruti":       "#52BE80",
    "Mercearia":        "#F39C12",
    "Embalagens":       "#AAB7B8",
    "Higiene/Limpeza":  "#85C1E9",
    "Lanchonete":       "#E74C3C",
    "Outros":           "#95A5A6",
}

# CFOPs que compõem "compras de comercialização" para fins do Simples Nacional
# Inclui: revenda de mercadorias (com e sem ST)
# Exclui: imobilizado (14xx), uso e consumo (15xx), energia (12xx), industrialização (11xx)
CFOPS_COMERCIALIZACAO = {
    # Compra de mercadoria para comercialização (revenda)
    "1102", "2102", "3102",
    # Compra de mercadoria para comercialização — substituição tributária
    "1403", "2403", "3403",
    # Compra de mercadoria de produtor rural para comercialização
    "1104", "2104", "3104",
    # Compra por conta e ordem (comercialização)
    "1117", "2117", "3117",
    # Compra de mercadoria para comercialização em operação com DI
    "3102",
    # Retorno simbólico de mercadoria vendida c/ substituição tributária
    "1410", "2410",
}


def parse_entradas_xml(arquivos) -> pd.DataFrame:
    """Parseia XMLs de nota fiscal de ENTRADA (compras) e retorna DataFrame
    com os itens, filtrando apenas CFOPs de comercialização."""
    import xml.etree.ElementTree as ET
    import zipfile, io as _io

    NS = "{http://www.portalfiscal.inf.br/nfe}"

    def t(name):
        return f"{NS}{name}"

    def gettxt(parent, child):
        el = parent.find(t(child))
        return el.text if el is not None and el.text else ""

    def getfloat(parent, child):
        txt = gettxt(parent, child)
        try:
            return float(txt) if txt else 0.0
        except ValueError:
            return 0.0

    rows = []

    def _process_xml_bytes(xml_data: bytes, source_name: str = ""):
        try:
            root = ET.fromstring(xml_data)
            root_local = root.tag.split("}")[-1] if "}" in root.tag else root.tag
            if root_local not in ("nfeProc", "NFe"):
                return
            nfe_el = root.find(t("NFe")) if root_local == "nfeProc" else root
            prot_el = root.find(t("protNFe")) if root_local == "nfeProc" else None
            if nfe_el is None:
                return
            infNFe = nfe_el.find(t("infNFe"))
            if infNFe is None:
                return

            # Verifica autorização
            if prot_el is not None:
                infProt = prot_el.find(t("infProt"))
                if infProt is not None:
                    cStat = gettxt(infProt, "cStat")
                    if cStat not in ("100", "150"):
                        return

            inf_id = infNFe.get("Id", "")
            chave = inf_id[3:] if inf_id.startswith("NFe") else inf_id

            ide = infNFe.find(t("ide"))
            if ide is None:
                return
            dhEmi = gettxt(ide, "dhEmi") or None
            nNF = gettxt(ide, "nNF")

            emit_el = infNFe.find(t("emit"))
            emitente = gettxt(emit_el, "xNome") if emit_el is not None else ""
            cnpj_emit = gettxt(emit_el, "CNPJ") if emit_el is not None else ""

            vNF = 0.0
            total_el = infNFe.find(t("total"))
            if total_el is not None:
                icms = total_el.find(t("ICMSTot"))
                if icms is not None:
                    vNF = getfloat(icms, "vNF")

            for det in infNFe.findall(t("det")):
                prod = det.find(t("prod"))
                if prod is None:
                    continue
                # Normaliza CFOP: remove pontos/traços e mantém só 4 dígitos
                # "1.102", "1102001", "1102.001" → "1102"
                import re as _re_cfop
                cfop_raw = gettxt(prod, "CFOP")
                cfop = _re_cfop.sub(r"[^\d]", "", cfop_raw)[:4]
                rows.append({
                    "chave":     chave,
                    "nNF":       nNF,
                    "dhEmi":     dhEmi,
                    "emitente":  emitente,
                    "cnpj_emit": cnpj_emit,
                    "CFOP":      cfop,
                    "xProd":     gettxt(prod, "xProd"),
                    "NCM":       gettxt(prod, "NCM"),
                    "qCom":      getfloat(prod, "qCom"),
                    "vUnCom":    getfloat(prod, "vUnCom"),
                    "vProd":     getfloat(prod, "vProd"),
                    "vNF":       vNF,
                    "fonte":     source_name,
                })
        except Exception:
            pass

    for arq in (arquivos or []):
        try:
            raw = arq.read() if hasattr(arq, "read") else arq
            name = getattr(arq, "name", "")
            if name.lower().endswith(".zip"):
                with zipfile.ZipFile(_io.BytesIO(raw)) as zf:
                    for zi in zf.namelist():
                        if zi.lower().endswith(".xml"):
                            _process_xml_bytes(zf.read(zi), zi)
            else:
                _process_xml_bytes(raw, name)
        except Exception:
            continue

    if not rows:
        return pd.DataFrame()

    df = pd.DataFrame(rows)
    if "dhEmi" in df.columns:
        df["dhEmi"] = pd.to_datetime(df["dhEmi"], utc=True, errors="coerce").dt.tz_localize(None)
    return df


def calc_simples_nacional(df_entradas: pd.DataFrame, faturamento_total: float) -> dict:
    """Calcula métricas para verificação da regra dos 80% do Simples Nacional.

    Retorna dict com:
        total_compras_comercializacao: float
        pct_faturamento: float (0-100)
        status: "OK" | "ALERTA" | "EXCEDIDO"
        df_por_cfop: DataFrame com breakdown por CFOP
        df_por_fornecedor: DataFrame com breakdown por fornecedor
        df_entradas_filtradas: DataFrame somente com linhas de comercialização
    """
    if df_entradas is None or df_entradas.empty:
        return {
            "total_compras_comercializacao": 0.0,
            "pct_faturamento": 0.0,
            "status": "SEM_DADOS",
            "df_por_cfop": pd.DataFrame(),
            "df_por_fornecedor": pd.DataFrame(),
            "df_entradas_filtradas": pd.DataFrame(),
        }

    # CFOP já vem normalizado (4 dígitos limpos) do parse_entradas_xml
    df_com = df_entradas[df_entradas["CFOP"].isin(CFOPS_COMERCIALIZACAO)].copy()
    total = df_com["vProd"].sum()
    pct = (total / faturamento_total * 100) if faturamento_total else 0.0

    if pct > 80:
        status = "EXCEDIDO"
    elif pct > 70:
        status = "ALERTA"
    else:
        status = "OK"

    df_por_cfop = (
        df_com.groupby("CFOP")
        .agg(
            total_compras=("vProd", "sum"),
            notas=("chave", "nunique"),
            itens=("vProd", "count"),
        )
        .reset_index()
        .sort_values("total_compras", ascending=False)
    )

    df_por_fornecedor = (
        df_com.groupby("emitente")
        .agg(
            total_compras=("vProd", "sum"),
            notas=("chave", "nunique"),
        )
        .reset_index()
        .sort_values("total_compras", ascending=False)
        .head(20)
    )

    return {
        "total_compras_comercializacao": total,
        "pct_faturamento": pct,
        "status": status,
        "df_por_cfop": df_por_cfop,
        "df_por_fornecedor": df_por_fornecedor,
        "df_entradas_filtradas": df_com,
    }


def categorizar_serie(s: pd.Series, ncm: pd.Series = None) -> pd.Series:
    """
    Classifica produtos pelo NCM fiscal (primário).
    Fallback por keyword apenas para alimentos preparados na padaria
    que compartilham capítulo NCM com industrializados.
    """
    result = pd.Series("Outros", index=s.index, dtype="object")

    # 1) NCM como classificação primária
    if ncm is not None:
        ncm_str = ncm.astype(str).str.zfill(8)
        sub4 = ncm_str.str[:4]
        cap2 = ncm_str.str[:2]
        ncm2 = {k: v for k, v in MAPA_NCM.items() if len(k) == 2}
        ncm4 = {k: v for k, v in MAPA_NCM.items() if len(k) == 4}
        # Aplica 2 dígitos primeiro, depois 4 dígitos sobrescrevem onde necessário
        for cod, cat in ncm2.items():
            result[cap2 == cod] = cat
        for cod, cat in ncm4.items():
            result[sub4 == cod] = cat

    # 2) Keyword como fallback só para quem ficou "Outros" ou "Panificação"/"Mercearia"
    #    (alimentos preparados na padaria que o NCM não distingue)
    s_up = s.str.upper().fillna("")
    for cat, palavras in reversed(MAPA_CATEGORIAS_FALLBACK):
        mask = pd.Series(False, index=s.index)
        for p in palavras:
            mask |= s_up.str.contains(p, regex=False, na=False)
        result[mask] = cat

    return result


def categorizar(nome: str) -> str:
    """Versão escalar — mantida para compatibilidade."""
    nome_up = str(nome).upper()
    for cat, palavras in MAPA_CATEGORIAS_FALLBACK:
        for p in palavras:
            if p in nome_up:
                return cat
    return "Outros"


# 
# CARREGAMENTO DOS DADOS
# 
@st.cache_data(show_spinner=False)
def carregar_nfce(file_bytes: bytes) -> pd.DataFrame:
    df = pd.read_excel(io.BytesIO(file_bytes))

    # Renomeia para nomes padronizados (tolerante a variações)
    col_map = {}
    for c in df.columns:
        cl = c.lower().strip()
        if "chave" in cl:             col_map[c] = "chave"
        elif cl == "nnf":             col_map[c] = "nNF"
        elif cl == "numitem":         col_map[c] = "numItem"
        elif cl == "cprod":           col_map[c] = "cProd"
        elif cl == "xprod":           col_map[c] = "xProd"
        elif cl == "ncm":             col_map[c] = "NCM"
        elif cl == "cfop":            col_map[c] = "CFOP"
        elif cl == "qcom":            col_map[c] = "qCom"
        elif cl == "vuncom":          col_map[c] = "vUnCom"
        elif cl == "vprod":           col_map[c] = "vProd"
        elif "icmstot_vnf" in cl:     col_map[c] = "vNF"
        elif "dhemi" in cl:           col_map[c] = "dhEmi"
    df = df.rename(columns=col_map)

    # Colunas numéricas
    for col in ["qCom", "vUnCom", "vProd", "vNF"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)

    # Data/hora — se existir dhEmi
    if "dhEmi" in df.columns:
        df["dhEmi"]     = pd.to_datetime(df["dhEmi"], errors="coerce")
        df["hora"]      = df["dhEmi"].dt.hour
        df["dia_semana"] = df["dhEmi"].dt.day_name()
        df["turno"]     = df["hora"].apply(
            lambda h: "Manhã" if 5 <= h < 12 else ("Tarde" if 12 <= h < 18 else "Noite")
            if pd.notna(h) else None
        )

    # Se não tem dhEmi, tenta extrair mês/ano da chave (pos 2-5 = AAMM)
    if "dhEmi" not in df.columns and "chave" in df.columns:
        def chave_para_data(ch):
            ch = str(ch)
            if len(ch) >= 6:
                try:
                    aa, mm = int(ch[2:4]), int(ch[4:6])
                    return pd.Timestamp(f"20{aa:02d}-{mm:02d}-01")
                except Exception:
                    pass
            return pd.NaT
        df["dhEmi"] = df["chave"].apply(chave_para_data)

    # Categoria
    if "xProd" in df.columns:
        df["categoria"] = categorizar_serie(df["xProd"], df.get("NCM"))

    df["fonte"] = "NFC-e"
    return df


@st.cache_data(show_spinner=False)
def carregar_nfe(file_bytes: bytes) -> pd.DataFrame:
    df = pd.read_excel(io.BytesIO(file_bytes))

    col_map = {}
    for c in df.columns:
        cl = c.lower().strip()
        if "chave" in cl:             col_map[c] = "chave"
        elif cl == "nnf":             col_map[c] = "nNF"
        elif cl == "numitem":         col_map[c] = "numItem"
        elif cl == "cprod":           col_map[c] = "cProd"
        elif cl == "xprod":           col_map[c] = "xProd"
        elif cl == "ncm":             col_map[c] = "NCM"
        elif cl == "cfop":            col_map[c] = "CFOP"
        elif cl == "qcom":            col_map[c] = "qCom"
        elif cl == "vuncom":          col_map[c] = "vUnCom"
        elif cl == "vprod":           col_map[c] = "vProd"
        elif "icmstot_vnf" in cl:     col_map[c] = "vNF"
        elif "dhemi" in cl:           col_map[c] = "dhEmi"
        elif "xnome" in cl and ".1" in c: col_map[c] = "destinatario"
        elif "situacao" in cl or "situa" in cl: col_map[c] = "situacao"
    df = df.rename(columns=col_map)

    for col in ["qCom", "vUnCom", "vProd", "vNF"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)

    if "dhEmi" in df.columns:
        df["dhEmi"] = pd.to_datetime(df["dhEmi"], errors="coerce")

    if "situacao" in df.columns:
        df = df[df["situacao"].astype(str).str.lower().str.contains("autoriza", na=False)]

    if "xProd" in df.columns:
        df["categoria"] = categorizar_serie(df["xProd"], df.get("NCM"))

    df["fonte"] = "NF-e"
    return df


@st.cache_data(show_spinner=False)
def carregar_zip(file_bytes: bytes):
    """
    Lê um arquivo ZIP contendo XMLs de NF-e/NFC-e.
    Separa automaticamente pelo campo <mod>: 65=NFC-e, 55=NF-e.
    Retorna (df_nfce, df_nfe) com o mesmo schema dos carregadores Excel.
    """
    import zipfile
    import xml.etree.ElementTree as ET

    NS = "{http://www.portalfiscal.inf.br/nfe}"

    def t(name):
        return f"{NS}{name}"

    def gettxt(parent, child):
        el = parent.find(t(child))
        return el.text if el is not None and el.text else ""

    def getfloat(parent, child):
        txt = gettxt(parent, child)
        try:
            return float(txt) if txt else 0.0
        except ValueError:
            return 0.0

    def _parse_entry(xml_data):
        try:
            root = ET.fromstring(xml_data)
            root_local = root.tag.split("}")[-1] if "}" in root.tag else root.tag
            if root_local not in ("nfeProc", "NFe"):
                return [], [], 1

            nfe_el  = root.find(t("NFe")) if root_local == "nfeProc" else root
            prot_el = root.find(t("protNFe")) if root_local == "nfeProc" else None
            if nfe_el is None:
                return [], [], 0
            infNFe = nfe_el.find(t("infNFe"))
            if infNFe is None:
                return [], [], 0

            inf_id = infNFe.get("Id", "")
            chave  = inf_id[3:] if inf_id.startswith("NFe") else inf_id
            ide    = infNFe.find(t("ide"))
            if ide is None:
                return [], [], 0

            mod   = gettxt(ide, "mod")
            nNF   = gettxt(ide, "nNF")
            dhEmi = gettxt(ide, "dhEmi") or None

            situacao = "Autorizada"
            if prot_el is not None:
                infProt = prot_el.find(t("infProt"))
                if infProt is not None:
                    cStat    = gettxt(infProt, "cStat")
                    situacao = "Autorizada" if cStat in ("100", "150") else f"cStat:{cStat}"

            vNF      = 0.0
            total_el = infNFe.find(t("total"))
            if total_el is not None:
                icms = total_el.find(t("ICMSTot"))
                if icms is not None:
                    vNF = getfloat(icms, "vNF")

            dest_el      = infNFe.find(t("dest"))
            destinatario = gettxt(dest_el, "xNome") if dest_el is not None else ""
            emit_el      = infNFe.find(t("emit"))
            emitente     = gettxt(emit_el, "xNome") if emit_el is not None else ""
            cnpj_emit    = gettxt(emit_el, "CNPJ")  if emit_el is not None else ""

            r_nfce, r_nfe = [], []
            for det in infNFe.findall(t("det")):
                numItem = det.get("nItem", "")
                prod    = det.find(t("prod"))
                if prod is None:
                    continue
                # CFOP 5929/001/002 = doc complementar de ECF, não é receita nova — ignora em NF-e
                if mod == "55" and gettxt(prod, "CFOP").startswith("5929"):
                    continue
                row = {
                    "chave": chave, "nNF": nNF, "numItem": numItem,
                    "cProd": gettxt(prod, "cProd"), "xProd": gettxt(prod, "xProd"),
                    "NCM":   gettxt(prod, "NCM"),   "CFOP":  gettxt(prod, "CFOP"),
                    "qCom":  getfloat(prod, "qCom"), "vUnCom": getfloat(prod, "vUnCom"),
                    "vProd": getfloat(prod, "vProd"), "vNF": vNF,
                    "dhEmi": dhEmi, "destinatario": destinatario, "emitente": emitente, "cnpj_emit": cnpj_emit, "situacao": situacao,
                }
                if mod == "65":
                    r_nfce.append(row)
                elif mod == "55":
                    r_nfe.append(row)
            return r_nfce, r_nfe, 0
        except Exception:
            return [], [], 1

    # Lê todos os bytes do ZIP de uma vez (sequencial — ZipFile não é thread-safe)
    with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
        xml_names  = [n for n in zf.namelist() if n.lower().endswith(".xml")]
        all_bytes  = [zf.read(entry) for entry in xml_names]

    # Parseia em paralelo
    from concurrent.futures import ThreadPoolExecutor, as_completed
    rows_nfce = []
    rows_nfe  = []
    skipped   = 0
    n_workers = min(16, max(4, len(all_bytes) // 500 + 1))
    with ThreadPoolExecutor(max_workers=n_workers) as pool:
        for r_nfce, r_nfe, sk in pool.map(_parse_entry, all_bytes):
            rows_nfce.extend(r_nfce)
            rows_nfe.extend(r_nfe)
            skipped += sk

    def montar_df(rows, fonte):
        if not rows:
            return pd.DataFrame()
        df = pd.DataFrame(rows)
        for col in ["qCom", "vUnCom", "vProd", "vNF"]:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
        df["numItem"] = pd.to_numeric(df["numItem"], errors="coerce").fillna(0).astype(int)
        if "dhEmi" in df.columns:
            df["dhEmi"] = pd.to_datetime(df["dhEmi"], errors="coerce", utc=False)
            # Se timezone-aware (ex: "2026-02-01T08:30:00-03:00"), converte para naive local
            if not df["dhEmi"].empty and hasattr(df["dhEmi"].dt, "tz") and df["dhEmi"].dt.tz is not None:
                df["dhEmi"] = df["dhEmi"].dt.tz_convert("America/Sao_Paulo").dt.tz_localize(None)
        if "xProd" in df.columns:
            df["categoria"] = categorizar_serie(df["xProd"], df.get("NCM"))
        df["fonte"] = fonte
        return df

    df_nfce_z = montar_df(rows_nfce, "NFC-e")
    df_nfe_z  = montar_df(rows_nfe,  "NF-e")

    # Pós-processamento NFC-e (igual ao carregar_nfce)
    if not df_nfce_z.empty and "dhEmi" in df_nfce_z.columns:
        df_nfce_z["hora"]       = df_nfce_z["dhEmi"].dt.hour
        df_nfce_z["dia_semana"] = df_nfce_z["dhEmi"].dt.day_name()
        df_nfce_z["turno"]      = df_nfce_z["hora"].apply(
            lambda h: "Manhã" if 5 <= h < 12 else ("Tarde" if 12 <= h < 18 else "Noite")
            if pd.notna(h) else None
        )

    # NF-e: filtra apenas autorizadas
    if not df_nfe_z.empty and "situacao" in df_nfe_z.columns:
        df_nfe_z = df_nfe_z[df_nfe_z["situacao"] == "Autorizada"].reset_index(drop=True)

    return df_nfce_z, df_nfe_z, len(xml_names), skipped


@st.cache_data(show_spinner=False)
def carregar_xmls_multi(arquivos: tuple):
    """
    Recebe uma tupla de (nome, bytes) de XMLs selecionados pelo usuário.
    Separa NFC-e (mod=65) e NF-e (mod=55) automaticamente.
    Retorna (df_nfce, df_nfe, total, ignorados).
    """
    import xml.etree.ElementTree as ET

    NS = "{http://www.portalfiscal.inf.br/nfe}"

    def t(name):
        return f"{NS}{name}"

    def gettxt(parent, child):
        el = parent.find(t(child))
        return el.text if el is not None and el.text else ""

    def getfloat(parent, child):
        txt = gettxt(parent, child)
        try:
            return float(txt) if txt else 0.0
        except ValueError:
            return 0.0

    rows_nfce = []
    rows_nfe  = []
    skipped   = 0

    for _nome, xml_data in arquivos:
        try:
            root       = ET.fromstring(xml_data)
            root_local = root.tag.split("}")[-1] if "}" in root.tag else root.tag

            if root_local not in ("nfeProc", "NFe"):
                skipped += 1
                continue

            nfe_el  = root.find(t("NFe")) if root_local == "nfeProc" else root
            prot_el = root.find(t("protNFe")) if root_local == "nfeProc" else None

            if nfe_el is None:
                continue

            infNFe = nfe_el.find(t("infNFe"))
            if infNFe is None:
                continue

            inf_id = infNFe.get("Id", "")
            chave  = inf_id[3:] if inf_id.startswith("NFe") else inf_id

            ide = infNFe.find(t("ide"))
            if ide is None:
                continue

            mod   = gettxt(ide, "mod")
            nNF   = gettxt(ide, "nNF")
            dhEmi = gettxt(ide, "dhEmi") or None

            situacao = "Autorizada"
            if prot_el is not None:
                infProt = prot_el.find(t("infProt"))
                if infProt is not None:
                    cStat    = gettxt(infProt, "cStat")
                    situacao = "Autorizada" if cStat in ("100", "150") else f"cStat:{cStat}"

            vNF      = 0.0
            total_el = infNFe.find(t("total"))
            if total_el is not None:
                icms = total_el.find(t("ICMSTot"))
                if icms is not None:
                    vNF = getfloat(icms, "vNF")

            dest_el      = infNFe.find(t("dest"))
            destinatario = gettxt(dest_el, "xNome") if dest_el is not None else ""
            emit_el      = infNFe.find(t("emit"))
            emitente     = gettxt(emit_el, "xNome") if emit_el is not None else ""
            cnpj_emit    = gettxt(emit_el, "CNPJ")  if emit_el is not None else ""

            for det in infNFe.findall(t("det")):
                numItem = det.get("nItem", "")
                prod    = det.find(t("prod"))
                if prod is None:
                    continue
                if mod == "55" and gettxt(prod, "CFOP").startswith("5929"):
                    continue

                row = {
                    "chave": chave, "nNF": nNF, "numItem": numItem,
                    "cProd": gettxt(prod, "cProd"), "xProd": gettxt(prod, "xProd"),
                    "NCM":   gettxt(prod, "NCM"),   "CFOP":  gettxt(prod, "CFOP"),
                    "qCom":  getfloat(prod, "qCom"), "vUnCom": getfloat(prod, "vUnCom"),
                    "vProd": getfloat(prod, "vProd"), "vNF": vNF,
                    "dhEmi": dhEmi, "destinatario": destinatario, "emitente": emitente, "cnpj_emit": cnpj_emit, "situacao": situacao,
                }

                if mod == "65":
                    rows_nfce.append(row)
                elif mod == "55":
                    rows_nfe.append(row)

        except Exception:
            skipped += 1
            continue

    def montar_df(rows, fonte):
        if not rows:
            return pd.DataFrame()
        df = pd.DataFrame(rows)
        for col in ["qCom", "vUnCom", "vProd", "vNF"]:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
        df["numItem"] = pd.to_numeric(df["numItem"], errors="coerce").fillna(0).astype(int)
        if "dhEmi" in df.columns:
            df["dhEmi"] = pd.to_datetime(df["dhEmi"], errors="coerce", utc=False)
            if not df["dhEmi"].empty and hasattr(df["dhEmi"].dt, "tz") and df["dhEmi"].dt.tz is not None:
                df["dhEmi"] = df["dhEmi"].dt.tz_convert("America/Sao_Paulo").dt.tz_localize(None)
        if "xProd" in df.columns:
            df["categoria"] = categorizar_serie(df["xProd"], df.get("NCM"))
        df["fonte"] = fonte
        return df

    df_nfce_m = montar_df(rows_nfce, "NFC-e")
    df_nfe_m  = montar_df(rows_nfe,  "NF-e")

    if not df_nfce_m.empty and "dhEmi" in df_nfce_m.columns:
        df_nfce_m["hora"]       = df_nfce_m["dhEmi"].dt.hour
        df_nfce_m["dia_semana"] = df_nfce_m["dhEmi"].dt.day_name()
        df_nfce_m["turno"]      = df_nfce_m["hora"].apply(
            lambda h: "Manhã" if 5 <= h < 12 else ("Tarde" if 12 <= h < 18 else "Noite")
            if pd.notna(h) else None
        )

    if not df_nfe_m.empty and "situacao" in df_nfe_m.columns:
        df_nfe_m = df_nfe_m[df_nfe_m["situacao"] == "Autorizada"].reset_index(drop=True)

    return df_nfce_m, df_nfe_m, len(arquivos), skipped


@st.cache_data(show_spinner=False)
def carregar_pasta(caminho: str):
    """
    Lê todos os XMLs de uma pasta local (recursivamente), sem precisar zipar.
    Separa NFC-e (mod=65) e NF-e (mod=55) automaticamente.
    Retorna (df_nfce, df_nfe, total_xmls, skipped).
    """
    from pathlib import Path
    import xml.etree.ElementTree as ET

    pasta = Path(caminho.strip())
    if not pasta.exists() or not pasta.is_dir():
        return pd.DataFrame(), pd.DataFrame(), 0, 0

    # Busca recursiva — case-insensitive nos sistemas que suportam
    xml_files = sorted(set(pasta.rglob("*.xml")) | set(pasta.rglob("*.XML")))

    NS = "{http://www.portalfiscal.inf.br/nfe}"

    def t(name):
        return f"{NS}{name}"

    def gettxt(parent, child):
        el = parent.find(t(child))
        return el.text if el is not None and el.text else ""

    def getfloat(parent, child):
        txt = gettxt(parent, child)
        try:
            return float(txt) if txt else 0.0
        except ValueError:
            return 0.0

    rows_nfce = []
    rows_nfe  = []
    skipped   = 0

    for xml_path in xml_files:
        try:
            xml_data = xml_path.read_bytes()
            root = ET.fromstring(xml_data)
            root_local = root.tag.split("}")[-1] if "}" in root.tag else root.tag

            if root_local not in ("nfeProc", "NFe"):
                skipped += 1
                continue

            nfe_el  = root.find(t("NFe")) if root_local == "nfeProc" else root
            prot_el = root.find(t("protNFe")) if root_local == "nfeProc" else None

            if nfe_el is None:
                continue

            infNFe = nfe_el.find(t("infNFe"))
            if infNFe is None:
                continue

            inf_id = infNFe.get("Id", "")
            chave  = inf_id[3:] if inf_id.startswith("NFe") else inf_id

            ide = infNFe.find(t("ide"))
            if ide is None:
                continue

            mod   = gettxt(ide, "mod")
            nNF   = gettxt(ide, "nNF")
            dhEmi = gettxt(ide, "dhEmi") or None

            situacao = "Autorizada"
            if prot_el is not None:
                infProt = prot_el.find(t("infProt"))
                if infProt is not None:
                    cStat    = gettxt(infProt, "cStat")
                    situacao = "Autorizada" if cStat in ("100", "150") else f"cStat:{cStat}"

            vNF = 0.0
            total_el = infNFe.find(t("total"))
            if total_el is not None:
                icms = total_el.find(t("ICMSTot"))
                if icms is not None:
                    vNF = getfloat(icms, "vNF")

            dest_el      = infNFe.find(t("dest"))
            destinatario = gettxt(dest_el, "xNome") if dest_el is not None else ""
            emit_el      = infNFe.find(t("emit"))
            emitente     = gettxt(emit_el, "xNome") if emit_el is not None else ""
            cnpj_emit    = gettxt(emit_el, "CNPJ")  if emit_el is not None else ""

            for det in infNFe.findall(t("det")):
                numItem = det.get("nItem", "")
                prod    = det.find(t("prod"))
                if prod is None:
                    continue
                if mod == "55" and gettxt(prod, "CFOP").startswith("5929"):
                    continue

                row = {
                    "chave":        chave,
                    "nNF":          nNF,
                    "numItem":      numItem,
                    "cProd":        gettxt(prod, "cProd"),
                    "xProd":        gettxt(prod, "xProd"),
                    "NCM":          gettxt(prod, "NCM"),
                    "CFOP":         gettxt(prod, "CFOP"),
                    "qCom":         getfloat(prod, "qCom"),
                    "vUnCom":       getfloat(prod, "vUnCom"),
                    "vProd":        getfloat(prod, "vProd"),
                    "vNF":          vNF,
                    "dhEmi":        dhEmi,
                    "destinatario": destinatario,
                    "emitente":     emitente,
                    "cnpj_emit":    cnpj_emit,
                    "situacao":     situacao,
                }

                if mod == "65":
                    rows_nfce.append(row)
                elif mod == "55":
                    rows_nfe.append(row)

        except Exception:
            skipped += 1
            continue

    def montar_df(rows, fonte):
        if not rows:
            return pd.DataFrame()
        df = pd.DataFrame(rows)
        for col in ["qCom", "vUnCom", "vProd", "vNF"]:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
        df["numItem"] = pd.to_numeric(df["numItem"], errors="coerce").fillna(0).astype(int)
        if "dhEmi" in df.columns:
            df["dhEmi"] = pd.to_datetime(df["dhEmi"], errors="coerce", utc=False)
            if not df["dhEmi"].empty and hasattr(df["dhEmi"].dt, "tz") and df["dhEmi"].dt.tz is not None:
                df["dhEmi"] = df["dhEmi"].dt.tz_convert("America/Sao_Paulo").dt.tz_localize(None)
        if "xProd" in df.columns:
            df["categoria"] = categorizar_serie(df["xProd"], df.get("NCM"))
        df["fonte"] = fonte
        return df

    df_nfce_p = montar_df(rows_nfce, "NFC-e")
    df_nfe_p  = montar_df(rows_nfe,  "NF-e")

    if not df_nfce_p.empty and "dhEmi" in df_nfce_p.columns:
        df_nfce_p["hora"]       = df_nfce_p["dhEmi"].dt.hour
        df_nfce_p["dia_semana"] = df_nfce_p["dhEmi"].dt.day_name()
        df_nfce_p["turno"]      = df_nfce_p["hora"].apply(
            lambda h: "Manhã" if 5 <= h < 12 else ("Tarde" if 12 <= h < 18 else "Noite")
            if pd.notna(h) else None
        )

    if not df_nfe_p.empty and "situacao" in df_nfe_p.columns:
        df_nfe_p = df_nfe_p[df_nfe_p["situacao"] == "Autorizada"].reset_index(drop=True)

    return df_nfce_p, df_nfe_p, len(xml_files), skipped


def _contar_xmls_pasta(caminho: str) -> int:
    """Conta XMLs em uma pasta sem carregar — para preview na sidebar."""
    try:
        from pathlib import Path
        p = Path(caminho.strip())
        if not p.exists() or not p.is_dir():
            return -1
        return len(list(p.rglob("*.xml")) + list(p.rglob("*.XML")))
    except Exception:
        return -1


@st.cache_data(show_spinner=False)
def carregar_pastas(caminhos: tuple):
    """
    Lê XMLs de múltiplas pastas combinadas.
    caminhos: tuple de strings com os caminhos das pastas.
    Retorna (df_nfce, df_nfe, total_xmls, ignorados).
    """
    from pathlib import Path
    import xml.etree.ElementTree as ET

    NS = "{http://www.portalfiscal.inf.br/nfe}"

    def t(name):
        return f"{NS}{name}"

    def gettxt(parent, child):
        el = parent.find(t(child))
        return el.text if el is not None and el.text else ""

    def getfloat(parent, child):
        txt = gettxt(parent, child)
        try:
            return float(txt) if txt else 0.0
        except ValueError:
            return 0.0

    # Coleta todos os XMLs de todas as pastas (sem duplicatas por path absoluto)
    xml_files = []
    vistos    = set()
    for caminho in caminhos:
        p = Path(caminho.strip())
        if not p.exists() or not p.is_dir():
            continue
        for xf in sorted(set(p.rglob("*.xml")) | set(p.rglob("*.XML"))):
            abs_p = str(xf.resolve())
            if abs_p not in vistos:
                vistos.add(abs_p)
                xml_files.append(xf)

    def _parse_xml(xml_path):
        """Parseia um XML e retorna (lista_rows_nfce, lista_rows_nfe, skipped_count)."""
        try:
            xml_data   = xml_path.read_bytes()
            root       = ET.fromstring(xml_data)
            root_local = root.tag.split("}")[-1] if "}" in root.tag else root.tag

            if root_local not in ("nfeProc", "NFe"):
                return [], [], 1

            nfe_el  = root.find(t("NFe")) if root_local == "nfeProc" else root
            prot_el = root.find(t("protNFe")) if root_local == "nfeProc" else None

            if nfe_el is None:
                return [], [], 0
            infNFe = nfe_el.find(t("infNFe"))
            if infNFe is None:
                return [], [], 0

            inf_id = infNFe.get("Id", "")
            chave  = inf_id[3:] if inf_id.startswith("NFe") else inf_id

            ide = infNFe.find(t("ide"))
            if ide is None:
                return [], [], 0

            mod   = gettxt(ide, "mod")
            nNF   = gettxt(ide, "nNF")
            dhEmi = gettxt(ide, "dhEmi") or None

            situacao = "Autorizada"
            if prot_el is not None:
                infProt = prot_el.find(t("infProt"))
                if infProt is not None:
                    cStat    = gettxt(infProt, "cStat")
                    situacao = "Autorizada" if cStat in ("100", "150") else f"cStat:{cStat}"

            vNF      = 0.0
            total_el = infNFe.find(t("total"))
            if total_el is not None:
                icms = total_el.find(t("ICMSTot"))
                if icms is not None:
                    vNF = getfloat(icms, "vNF")

            dest_el      = infNFe.find(t("dest"))
            destinatario = gettxt(dest_el, "xNome") if dest_el is not None else ""
            emit_el      = infNFe.find(t("emit"))
            emitente     = gettxt(emit_el, "xNome") if emit_el is not None else ""
            cnpj_emit    = gettxt(emit_el, "CNPJ")  if emit_el is not None else ""

            r_nfce, r_nfe = [], []
            for det in infNFe.findall(t("det")):
                numItem = det.get("nItem", "")
                prod    = det.find(t("prod"))
                if prod is None:
                    continue
                # CFOP 5929/001/002 = doc complementar de ECF, não é receita nova — ignora em NF-e
                if mod == "55" and gettxt(prod, "CFOP").startswith("5929"):
                    continue
                row = {
                    "chave": chave, "nNF": nNF, "numItem": numItem,
                    "cProd": gettxt(prod, "cProd"), "xProd": gettxt(prod, "xProd"),
                    "NCM":   gettxt(prod, "NCM"),   "CFOP":  gettxt(prod, "CFOP"),
                    "qCom":  getfloat(prod, "qCom"), "vUnCom": getfloat(prod, "vUnCom"),
                    "vProd": getfloat(prod, "vProd"), "vNF": vNF,
                    "dhEmi": dhEmi, "destinatario": destinatario, "emitente": emitente, "cnpj_emit": cnpj_emit, "situacao": situacao,
                }
                if mod == "65":
                    r_nfce.append(row)
                elif mod == "55":
                    r_nfe.append(row)
            return r_nfce, r_nfe, 0
        except Exception:
            return [], [], 1

    # Leitura paralela — I/O bound, ThreadPoolExecutor é ideal
    from concurrent.futures import ThreadPoolExecutor, as_completed
    rows_nfce = []
    rows_nfe  = []
    skipped   = 0

    n_workers = min(16, max(4, len(xml_files) // 500 + 1))
    with ThreadPoolExecutor(max_workers=n_workers) as pool:
        futures = {pool.submit(_parse_xml, xf): xf for xf in xml_files}
        for fut in as_completed(futures):
            r_nfce, r_nfe, sk = fut.result()
            rows_nfce.extend(r_nfce)
            rows_nfe.extend(r_nfe)
            skipped += sk

    def montar_df(rows, fonte):
        if not rows:
            return pd.DataFrame()
        df = pd.DataFrame(rows)
        for col in ["qCom", "vUnCom", "vProd", "vNF"]:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
        df["numItem"] = pd.to_numeric(df["numItem"], errors="coerce").fillna(0).astype(int)
        if "dhEmi" in df.columns:
            df["dhEmi"] = pd.to_datetime(df["dhEmi"], errors="coerce", utc=False)
            if not df["dhEmi"].empty and hasattr(df["dhEmi"].dt, "tz") and df["dhEmi"].dt.tz is not None:
                df["dhEmi"] = df["dhEmi"].dt.tz_convert("America/Sao_Paulo").dt.tz_localize(None)
        if "xProd" in df.columns:
            df["categoria"] = categorizar_serie(df["xProd"], df.get("NCM"))
        df["fonte"] = fonte
        return df

    df_nfce_p = montar_df(rows_nfce, "NFC-e")
    df_nfe_p  = montar_df(rows_nfe,  "NF-e")

    if not df_nfce_p.empty and "dhEmi" in df_nfce_p.columns:
        df_nfce_p["hora"]       = df_nfce_p["dhEmi"].dt.hour
        df_nfce_p["dia_semana"] = df_nfce_p["dhEmi"].dt.day_name()
        df_nfce_p["turno"]      = df_nfce_p["hora"].apply(
            lambda h: "Manhã" if 5 <= h < 12 else ("Tarde" if 12 <= h < 18 else "Noite")
            if pd.notna(h) else None
        )

    if not df_nfe_p.empty and "situacao" in df_nfe_p.columns:
        df_nfe_p = df_nfe_p[df_nfe_p["situacao"] == "Autorizada"].reset_index(drop=True)

    return df_nfce_p, df_nfe_p, len(xml_files), skipped


def _is_cloud() -> bool:
    """True quando rodando no Streamlit Cloud (sem display/tkinter disponível)."""
    import os
    if os.environ.get("STREAMLIT_SHARING_MODE") or os.environ.get("HOME") == "/home/appuser":
        return True
    try:
        import tkinter as _tk
        _r = _tk.Tk()
        _r.destroy()
        return False
    except Exception:
        return True


@st.cache_data(show_spinner=False)
def processar_fontes_universal(arquivos: tuple, pastas: tuple):
    """
    Processador unificado: recebe (nome, bytes) de qualquer arquivo
    (ZIP, RAR, 7z, XML) e caminhos de pastas locais.
    Vasculha tudo recursivamente e retorna (df_nfce, df_nfe, n_xml, n_skip).
    """
    import xml.etree.ElementTree as ET
    import zipfile

    # Palavras que — no NOME DO ARQUIVO (sem o caminho) — indicam notas não autorizadas
    _PALAVRAS_EXCLUIR = ("CANCELAD", "INUTILIZAD", "DENEGAD")

    def _arquivo_excluido(nome_completo: str) -> bool:
        """
        Verifica só o ÚLTIMO componente do path (nome do arquivo/zip em si),
        ignorando nomes de pastas ancestrais — evita falsos positivos quando a
        pasta pai tem palavras como 'Canceladas' no nome.
        """
        # Pega apenas o filename, descartando qualquer caminho de diretório
        nome_base = nome_completo.replace("\\", "/").split("/")[-1].upper()
        return any(p in nome_base for p in _PALAVRAS_EXCLUIR)

    def extrair_xml_bytes(data: bytes, nome: str, herdou_exclusao: bool = False) -> list:
        """
        Extrai lista de bytes de XMLs de qualquer arquivo.
        Arquivos/pastas cujo NOME (não o caminho) contenha CANCELADAS,
        INUTILIZADAS ou DENEGADAS são ignorados — e a exclusão propaga para
        todo o conteúdo interno.
        """
        ext = nome.lower().rsplit(".", 1)[-1] if "." in nome else ""
        excluir = herdou_exclusao or _arquivo_excluido(nome)
        resultado = []
        if ext == "xml":
            return [] if excluir else [data]
        if ext == "zip":
            try:
                with zipfile.ZipFile(io.BytesIO(data)) as zf:
                    for entry in zf.namelist():
                        eb = zf.read(entry)
                        resultado.extend(extrair_xml_bytes(eb, entry, excluir))
            except Exception:
                pass
        elif ext == "rar":
            try:
                import rarfile
                with rarfile.RarFile(io.BytesIO(data)) as rf:
                    for entry in rf.namelist():
                        eb = rf.read(entry)
                        resultado.extend(extrair_xml_bytes(eb, entry, excluir))
            except Exception:
                pass
        elif ext in ("7z", "7zip"):
            try:
                import py7zr
                with py7zr.SevenZipFile(io.BytesIO(data)) as szf:
                    for nome_arq, bio in szf.read().items():
                        resultado.extend(extrair_xml_bytes(bio.read(), nome_arq, excluir))
            except Exception:
                pass
        return resultado

    def parse_xml(xml_data: bytes):
        """
        Retorna (rows_nfce, rows_nfe, skipped, chave_cancelada).
        - rows_*        : lista de dicts de itens (pode ser None)
        - skipped       : 1 se parse falhou, 0 caso contrário
        - chave_cancelada: chave de 44 dígitos se este XML é um evento de
                           cancelamento aceito pela SEFAZ (procEventoNFe /
                           retEnvEvento com tpEvento=110111 e cStat=135).
                           None caso contrário.
        """
        try:
            root = ET.fromstring(xml_data)

            # ── Detecta namespace dinamicamente ──
            _ns = root.tag.split("}")[0] + "}" if "}" in root.tag else ""
            def _t(name): return f"{_ns}{name}"
            def _gettxt(p, c):
                el = p.find(_t(c)); return el.text if el is not None and el.text else ""
            def _getfloat(p, c):
                try: return float(_gettxt(p, c) or 0)
                except: return 0.0

            root_local = root.tag.split("}")[-1] if "}" in root.tag else root.tag

            # ── Evento de cancelamento (procEventoNFe ou retEnvEvento) ──
            if root_local in ("procEventoNFe", "retEnvEvento"):
                # Procura em qualquer profundidade por tpEvento=110111 + cStat=135
                # (cancelamento aceito pela SEFAZ)
                chave_canc = None
                for el in root.iter():
                    local = el.tag.split("}")[-1] if "}" in el.tag else el.tag
                    if local == "tpEvento" and el.text == "110111":
                        # Acha a chave da NF-e cancelada
                        parent = root
                        for candidate in root.iter():
                            c_local = candidate.tag.split("}")[-1] if "}" in candidate.tag else candidate.tag
                            if c_local == "chNFe" and candidate.text and len(candidate.text) == 44:
                                chave_canc = candidate.text.strip()
                                break
                        break
                if chave_canc:
                    # Só registra como cancelada se a SEFAZ aceitou (cStat=135)
                    cstat_aceito = False
                    for el in root.iter():
                        local = el.tag.split("}")[-1] if "}" in el.tag else el.tag
                        if local == "cStat" and el.text in ("135", "155", "101"):
                            cstat_aceito = True
                            break
                    if cstat_aceito:
                        return None, None, 0, chave_canc
                return None, None, 0, None

            if root_local not in ("nfeProc", "NFe"):
                return None, None, 0, None   # ignora silenciosamente outros tipos

            nfe_el  = root.find(_t("NFe")) if root_local == "nfeProc" else root
            prot_el = root.find(_t("protNFe")) if root_local == "nfeProc" else None
            if nfe_el is None: return None, None, 0, None
            infNFe = nfe_el.find(_t("infNFe"))
            if infNFe is None: return None, None, 0, None
            inf_id = infNFe.get("Id", "")
            chave  = inf_id[3:] if inf_id.startswith("NFe") else inf_id
            ide    = infNFe.find(_t("ide"))
            if ide is None: return None, None, 0, None
            mod   = _gettxt(ide, "mod")
            nNF   = _gettxt(ide, "nNF")
            dhEmi = _gettxt(ide, "dhEmi") or None
            situacao = "Autorizada"
            if prot_el is not None:
                infProt = prot_el.find(_t("infProt"))
                if infProt is not None:
                    cStat = _gettxt(infProt, "cStat")
                    situacao = "Autorizada" if cStat in ("100", "150") else f"cStat:{cStat}"
            # vNF: tenta ICMSTot/vNF, fallback para soma dos vProd dos itens
            vNF = 0.0
            tot = infNFe.find(_t("total"))
            if tot is not None:
                icms = tot.find(_t("ICMSTot"))
                if icms is not None:
                    vNF = _getfloat(icms, "vNF")
            dest_el = infNFe.find(_t("dest"))
            destinatario = _gettxt(dest_el, "xNome") if dest_el is not None else ""
            emit_el = infNFe.find(_t("emit"))
            emitente = _gettxt(emit_el, "xNome") if emit_el is not None else ""
            cnpj_emit = _gettxt(emit_el, "CNPJ") if emit_el is not None else ""
            rows_n, rows_e = [], []
            soma_vprod = 0.0
            for det in infNFe.findall(_t("det")):
                numItem = det.get("nItem", "")
                prod = det.find(_t("prod"))
                if prod is None: continue
                if mod == "55" and _gettxt(prod, "CFOP").startswith("5929"):
                    continue
                vp = _getfloat(prod, "vProd")
                soma_vprod += vp
                row = {
                    "chave": chave, "nNF": nNF, "numItem": numItem,
                    "cProd": _gettxt(prod, "cProd"), "xProd": _gettxt(prod, "xProd"),
                    "NCM": _gettxt(prod, "NCM"), "CFOP": _gettxt(prod, "CFOP"),
                    "qCom": _getfloat(prod, "qCom"), "vUnCom": _getfloat(prod, "vUnCom"),
                    "vProd": vp, "vNF": vNF,
                    "dhEmi": dhEmi, "destinatario": destinatario, "emitente": emitente, "cnpj_emit": cnpj_emit, "situacao": situacao,
                }
                if mod == "65": rows_n.append(row)
                elif mod == "55": rows_e.append(row)
            # Se vNF ficou 0 mas há itens, usa soma dos vProd como fallback
            if vNF == 0.0 and soma_vprod > 0:
                for r in rows_n: r["vNF"] = soma_vprod
                for r in rows_e: r["vNF"] = soma_vprod
            return rows_n, rows_e, 0, None
        except Exception:
            return None, None, 1, None

    # Coleta todos os bytes de XMLs
    all_xml_bytes = []
    for nome, data in arquivos:
        all_xml_bytes.extend(extrair_xml_bytes(data, nome))

    # Pastas locais
    if pastas:
        from pathlib import Path
        vistos = set()
        for caminho in pastas:
            p = Path(caminho.strip())
            if not p.exists(): continue
            for xf in sorted(p.rglob("*")):
                abs_p = str(xf.resolve())
                if abs_p in vistos: continue
                vistos.add(abs_p)
                ext = xf.suffix.lower().lstrip(".")
                if ext in ("xml",):
                    try: all_xml_bytes.append(xf.read_bytes())
                    except: pass
                elif ext in ("zip", "rar", "7z"):
                    try: all_xml_bytes.extend(extrair_xml_bytes(xf.read_bytes(), xf.name))
                    except: pass

    # Parseia em paralelo
    from concurrent.futures import ThreadPoolExecutor
    rows_nfce, rows_nfe, skipped = [], [], 0
    chaves_canceladas: set = set()
    with ThreadPoolExecutor(max_workers=min(16, max(4, len(all_xml_bytes) // 500 + 1))) as pool:
        for r_n, r_e, sk, chave_canc in pool.map(parse_xml, all_xml_bytes):
            if r_n: rows_nfce.extend(r_n)
            if r_e: rows_nfe.extend(r_e)
            skipped += sk
            if chave_canc: chaves_canceladas.add(chave_canc)

    # Remove notas cujo evento de cancelamento foi detectado
    if chaves_canceladas:
        rows_nfce = [r for r in rows_nfce if r["chave"] not in chaves_canceladas]
        rows_nfe  = [r for r in rows_nfe  if r["chave"] not in chaves_canceladas]
        skipped  += len(chaves_canceladas)

    def montar_df(rows, fonte):
        if not rows: return pd.DataFrame()
        df = pd.DataFrame(rows)
        for col in ["qCom", "vUnCom", "vProd", "vNF"]:
            if col in df.columns:
                df[col] = pd.to_numeric(df[col], errors="coerce").fillna(0)
        df["numItem"] = pd.to_numeric(df["numItem"], errors="coerce").fillna(0).astype(int)
        if "dhEmi" in df.columns:
            df["dhEmi"] = pd.to_datetime(df["dhEmi"], errors="coerce", utc=False)
            if not df["dhEmi"].empty and hasattr(df["dhEmi"].dt, "tz") and df["dhEmi"].dt.tz is not None:
                df["dhEmi"] = df["dhEmi"].dt.tz_convert("America/Sao_Paulo").dt.tz_localize(None)
        if "xProd" in df.columns:
            df["categoria"] = categorizar_serie(df["xProd"], df.get("NCM"))
        df["fonte"] = fonte
        return df

    df_nfce = montar_df(rows_nfce, "NFC-e")
    df_nfe  = montar_df(rows_nfe,  "NF-e")

    if not df_nfce.empty and "dhEmi" in df_nfce.columns:
        df_nfce["hora"]       = df_nfce["dhEmi"].dt.hour
        df_nfce["dia_semana"] = df_nfce["dhEmi"].dt.day_name()
        df_nfce["turno"]      = df_nfce["hora"].apply(
            lambda h: "Manhã" if 5 <= h < 12 else ("Tarde" if 12 <= h < 18 else "Noite")
            if pd.notna(h) else None)

    # Conta rejeitadas antes de filtrar (por nota única, não por linha/item)
    n_rejeitadas_nfce = 0
    n_rejeitadas_nfe  = 0
    if not df_nfce.empty and "situacao" in df_nfce.columns:
        n_rejeitadas_nfce = df_nfce[df_nfce["situacao"] != "Autorizada"]["chave"].nunique()
        df_nfce = df_nfce[df_nfce["situacao"] == "Autorizada"].reset_index(drop=True)
    if not df_nfe.empty and "situacao" in df_nfe.columns:
        n_rejeitadas_nfe = df_nfe[df_nfe["situacao"] != "Autorizada"]["chave"].nunique()
        df_nfe = df_nfe[df_nfe["situacao"] == "Autorizada"].reset_index(drop=True)
    skipped += (n_rejeitadas_nfce + n_rejeitadas_nfe)

    return df_nfce, df_nfe, len(all_xml_bytes), skipped


def _abrir_seletor_pasta() -> str:
    """Abre o seletor nativo de pasta do Windows via tkinter."""
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk()
        root.withdraw()
        root.wm_attributes("-topmost", 1)
        pasta = filedialog.askdirectory(title="Selecione a pasta com os XMLs fiscais")
        root.destroy()
        return pasta or ""
    except Exception:
        return ""


def _logo_base64() -> str:
    """Retorna a logo em base64 para exibir em HTML, ou string vazia se não existir."""
    import base64
    from pathlib import Path
    logo_path = Path(__file__).parent / "LOGO S FUNDO 2.png"
    if not logo_path.exists():
        return ""
    try:
        return base64.b64encode(logo_path.read_bytes()).decode()
    except Exception:
        return ""


#
# ANALYTICS
#
def calc_kpis(df: pd.DataFrame) -> dict:
    notas        = df.drop_duplicates("chave")
    fat          = notas["vNF"].sum()
    n            = len(notas)
    tm           = fat / n if n else 0
    total_itens  = len(df)
    ipc          = total_itens / n if n else 0   # total linhas ÷ total pedidos
    return {
        "faturamento":  fat,
        "n_pedidos":    n,
        "ticket_medio": tm,
        "ipc":          ipc,
        "total_itens":  total_itens,
    }


def calc_categorias(df: pd.DataFrame) -> pd.DataFrame:
    notas_por_cat = (
        df.groupby(["chave", "categoria"])["vProd"]
        .sum()
        .reset_index()
        .groupby("categoria")
        .agg(receita=("vProd", "sum"), n_pedidos=("chave", "count"))
        .reset_index()
        .sort_values("receita", ascending=False)
    )
    total = notas_por_cat["receita"].sum()
    notas_por_cat["pct"] = notas_por_cat["receita"] / total * 100
    return notas_por_cat


def calc_basket_pares(df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    # Só produtos que aparecem em pelo menos 2 notas (elimina cauda longa inútil)
    freq_prod = df.groupby("xProd")["chave"].nunique()
    prods_validos = set(freq_prod[freq_prod >= 2].index)
    df_f = df[df["xProd"].isin(prods_validos)]

    # Agrupa por nota — lista de valores direto (sem iterrows)
    baskets = df_f.groupby("chave")["xProd"].agg(lambda x: list(set(x)))

    contagem = Counter()
    for prods in baskets:
        if len(prods) >= 2:
            for par in combinations(sorted(prods), 2):
                contagem[par] += 1

    rows = [{"Produto A": a, "Produto B": b, "Frequência": f}
            for (a, b), f in contagem.most_common(top_n)]
    return pd.DataFrame(rows)


def calc_basket_trios(df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    # Limita aos top 300 produtos mais frequentes para não explodir em O(n³)
    freq_prod = df.groupby("xProd")["chave"].nunique()
    top_prods = set(freq_prod.nlargest(300).index)
    df_f = df[df["xProd"].isin(top_prods)]

    baskets = df_f.groupby("chave")["xProd"].agg(lambda x: list(set(x)))

    contagem = Counter()
    for prods in baskets:
        prods_filt = sorted(set(prods) & top_prods)
        if len(prods_filt) >= 3:
            for trio in combinations(prods_filt, 3):
                contagem[trio] += 1

    nomes_combo = {
        0: "Combo Clássico", 1: "Combo Matinal", 2: "Combo Família",
        3: "Combo Sabor",    4: "Combo Tradicional", 5: "Combo Premium",
        6: "Combo Lanche",   7: "Combo Especial",    8: "Combo Favorito",
        9: "Combo Delícia",
    }
    rows = []
    for i, ((a, b, c), f) in enumerate(contagem.most_common(top_n)):
        rows.append({
            "Produto A": a, "Produto B": b, "Produto C": c,
            "Frequência": f,
            "Sugestão": nomes_combo.get(i, f"Combo {i+1}"),
        })
    return pd.DataFrame(rows)


def calc_cesta(df: pd.DataFrame) -> pd.DataFrame:
    ipc = df.groupby("chave")["numItem"].count().reset_index(name="itens")
    dist = ipc["itens"].value_counts().sort_index().reset_index()
    dist.columns = ["Itens/Pedido", "Nº Pedidos"]
    dist["% do Total"] = (dist["Nº Pedidos"] / dist["Nº Pedidos"].sum() * 100).round(1)
    return dist


def calc_bcg(df: pd.DataFrame) -> pd.DataFrame:
    prod = (
        df.groupby("xProd")
        .agg(
            frequencia=("chave", "nunique"),
            receita=("vProd", "sum"),
        )
        .reset_index()
    )
    med_f = prod["frequencia"].median()
    med_r = prod["receita"].median()

    alto_f = prod["frequencia"] >= med_f
    alta_r = prod["receita"]    >= med_r
    prod["BCG"] = np.select(
        [alto_f & alta_r, ~alto_f & alta_r, alto_f & ~alta_r],
        ["Estrela",        "Vaca",           "Interrogação"],
        default="Cão",
    )
    return prod.sort_values("receita", ascending=False)


def calc_curva_abc(df: pd.DataFrame) -> pd.DataFrame:
    """
    Classifica todos os produtos pela Curva ABC com base na receita acumulada.
      A: produtos que representam os primeiros 80 % da receita (itens vitais)
      B: 80–95 % acumulado (itens importantes)
      C: 95–100 % acumulado (itens triviais)
    """
    prod = (
        df.groupby("xProd")
        .agg(frequencia=("chave", "nunique"), receita=("vProd", "sum"))
        .reset_index()
        .sort_values("receita", ascending=False)
        .reset_index(drop=True)
    )
    total = prod["receita"].sum()
    prod["pct_receita"]  = prod["receita"] / total * 100
    prod["pct_acumulado"] = prod["pct_receita"].cumsum()

    def _grupo(acum):
        if acum <= 80:
            return "A"
        elif acum <= 95:
            return "B"
        return "C"

    prod["Curva"] = prod["pct_acumulado"].apply(_grupo)
    prod.insert(0, "Rank", range(1, len(prod) + 1))
    prod = prod.rename(columns={
        "xProd":        "Produto",
        "frequencia":   "Frequência",
        "receita":      "Receita (R$)",
        "pct_receita":  "% Receita",
        "pct_acumulado":"% Acumulado",
    })
    return prod[["Rank", "Produto", "Curva", "Frequência", "Receita (R$)", "% Receita", "% Acumulado"]]


def calc_crossell(df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    baskets_cat = df.groupby("chave")["categoria"].agg(lambda x: list(set(x)))
    contagem = Counter()
    for cats in baskets_cat:
        if len(cats) >= 2:
            for par in combinations(sorted(cats), 2):
                contagem[par] += 1
    rows = [{"Categoria A": a, "Categoria B": b, "Frequência": f}
            for (a, b), f in contagem.most_common(top_n)]
    return pd.DataFrame(rows)


def calc_remocao(df: pd.DataFrame, top_n: int = 15) -> pd.DataFrame:
    prod = (
        df.groupby("xProd")
        .agg(frequencia=("chave", "nunique"), receita=("vProd", "sum"))
        .reset_index()
        .sort_values("receita")
        .head(top_n)
    )
    return prod


def calc_turno(df: pd.DataFrame) -> pd.DataFrame | None:
    if "turno" not in df.columns:
        return None
    return (
        df.groupby(["turno", "xProd"])
        .agg(frequencia=("chave", "nunique"), receita=("vProd", "sum"))
        .reset_index()
        .sort_values(["turno", "receita"], ascending=[True, False])
    )


#  PROMPT 5 — Produtos âncora (pedido com 1 só item) 
def calc_solo_produtos(df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    contagem = df.groupby("chave")["numItem"].count()
    notas_solo = contagem[contagem == 1].index
    solo = df[df["chave"].isin(notas_solo)]
    return (
        solo.groupby("xProd")
        .agg(frequencia=("chave", "count"), receita=("vProd", "sum"))
        .reset_index()
        .sort_values("frequencia", ascending=False)
        .head(top_n)
    )


#  PROMPT 3 — Anti-pares (raramente juntos) 
def calc_anti_pares(df: pd.DataFrame, min_freq: int = 30, top_n: int = 10) -> pd.DataFrame:
    freq = df.groupby("xProd")["chave"].nunique()
    produtos_freq = freq[freq >= min_freq].nlargest(60).index.tolist()
    total = df["chave"].nunique()

    basket = (
        df[df["xProd"].isin(produtos_freq)]
        .groupby("chave")["xProd"]
        .apply(set)
    )
    par_obs = Counter()
    for prods in basket:
        for par in combinations(sorted(prods), 2):
            par_obs[par] += 1

    rows = []
    for (a, b) in combinations(sorted(produtos_freq), 2):
        obs = par_obs.get((min(a, b), max(a, b)), 0)
        exp = (freq[a] / total) * (freq[b] / total) * total
        lift = obs / exp if exp > 0 else 0
        if lift < 0.25 and obs < 5:
            rows.append({
                "Produto A": a, "Produto B": b,
                "Freq. A": int(freq[a]), "Freq. B": int(freq[b]),
                "Juntos": obs, "Lift": round(lift, 3),
            })

    return pd.DataFrame(sorted(rows, key=lambda x: x["Lift"])[:top_n])


#  PROMPT 7 — Por dia da semana 
def calc_por_dia_semana(df: pd.DataFrame):
    if "dia_semana" not in df.columns:
        return None, None

    dia_pt   = df["dia_semana"].map(DAYS_MAP)
    tipo_dia = df["dia_semana"].isin(["Saturday", "Sunday"]).map({True: "Final de Semana", False: "Dia Útil"})
    df2 = df.assign(dia_pt=dia_pt, tipo_dia=tipo_dia)

    por_tipo = (
        df2.groupby(["tipo_dia", "xProd"])
        .agg(frequencia=("chave", "nunique"), receita=("vProd", "sum"))
        .reset_index()
        .sort_values(["tipo_dia", "receita"], ascending=[True, False])
    )

    por_dia = (
        df2.groupby(["dia_pt", "xProd"])
        .agg(frequencia=("chave", "nunique"), receita=("vProd", "sum"))
        .reset_index()
        .sort_values(["dia_pt", "receita"], ascending=[True, False])
    )

    return por_tipo, por_dia


#  PROMPTS 8 & 9 — Ticket drivers 
def calc_ticket_drivers(df: pd.DataFrame, top_n: int = 10):
    ticket_nota = df.drop_duplicates("chave")[["chave", "vNF"]].set_index("chave")["vNF"]
    tm_geral    = ticket_nota.mean()

    # Vetorizado: junta ticket de cada nota com os produtos da mesma nota
    prod_chave = df[["xProd", "chave"]].drop_duplicates(["xProd", "chave"])
    merged = prod_chave.merge(ticket_nota.rename("ticket"), left_on="chave", right_index=True)
    agg = (
        merged.groupby("xProd")
        .agg(n_pedidos=("chave", "nunique"),
             tm_prod=("ticket", "mean"))
        .reset_index()
    )
    agg = agg[agg["n_pedidos"] >= 10].copy()
    if agg.empty:
        return pd.DataFrame(), pd.DataFrame()

    agg["Ticket Médio Geral"]       = round(tm_geral, 2)
    agg["Diferença R$"]             = (agg["tm_prod"] - tm_geral).round(2)
    agg["Diferença %"]              = ((agg["tm_prod"] / tm_geral - 1) * 100).round(1)
    agg["Ticket Médio c/ Produto"]  = agg["tm_prod"].round(2)
    agg = agg.rename(columns={"xProd": "Produto", "n_pedidos": "Nº Pedidos"})
    agg = agg.drop(columns="tm_prod")

    elevadores = agg.nlargest(top_n, "Diferença R$").reset_index(drop=True)
    redutores  = agg.nsmallest(top_n, "Diferença R$").reset_index(drop=True)
    return elevadores, redutores


#  PROMPT 16 — Simulação de preços 
def calc_simulacao_precos(df: pd.DataFrame, top_n: int = 15) -> pd.DataFrame:
    prod = (
        df.groupby("xProd")
        .agg(preco_medio=("vUnCom", "mean"), volume=("qCom", "sum"), receita=("vProd", "sum"))
        .reset_index()
        .sort_values("receita", ascending=False)
        .head(top_n)
    )
    for pct in [10, 15, 20]:
        prod[f"+{pct}% preço (-5% vol)"] = (prod["receita"] * (1 + pct/100) * 0.95).round(2)
        prod[f"Δ +{pct}%"] = (prod[f"+{pct}% preço (-5% vol)"] - prod["receita"]).round(2)
    return prod


#  PROMPT 41 — Simulação de crescimento de receita 
def calc_simulacao_receita(kpis: dict, df: pd.DataFrame) -> pd.DataFrame:
    fat = kpis["faturamento"]
    tm  = kpis["ticket_medio"]
    ipc = kpis["ipc"]
    n   = kpis["n_pedidos"]

    valor_por_item = tm / ipc if ipc else 0

    pot_dias = 0
    if "dhEmi" in df.columns and df["dhEmi"].notna().any():
        notas_dia = df.drop_duplicates("chave").copy()
        notas_dia["_data"] = notas_dia["dhEmi"].dt.date
        daily   = notas_dia.groupby("_data")["vNF"].sum()
        media_d = daily.mean()
        fracos  = daily[daily < media_d * 0.7]
        if len(fracos):
            pot_dias = (media_d - fracos.mean()) * len(fracos)

    return pd.DataFrame([
        {"Estratégia": "Ticket Médio +10%",
         "Como": "Upsell no caixa e combos sugeridos",
         "Impacto Mensal": fat * 0.10,
         "Impacto Anual": fat * 0.10 * 12,
         "Complexidade": "Baixa"},
        {"Estratégia": "Ticket Médio +20%",
         "Como": "Combos premium e cardápio reposicionado",
         "Impacto Mensal": fat * 0.20,
         "Impacto Anual": fat * 0.20 * 12,
         "Complexidade": "Média"},
        {"Estratégia": "+1 item por pedido (Média)",
         "Como": "Sugestão ativa no balcão para complementar",
         "Impacto Mensal": valor_por_item * n,
         "Impacto Anual": valor_por_item * n * 12,
         "Complexidade": "Baixa"},
    ])


#  PROMPT 48 — Precificação de combos 
def calc_combo_pricing(df_pares: pd.DataFrame, df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    if df_pares.empty:
        return pd.DataFrame()

    preco_medio = df.groupby("xProd")["vUnCom"].mean()
    rows = []
    for _, row in df_pares.head(top_n).iterrows():
        a, b = row["Produto A"], row["Produto B"]
        pa = preco_medio.get(a, 0)
        pb = preco_medio.get(b, 0)
        total = pa + pb
        if total <= 0:
            continue
        rows.append({
            "Combo":               f"{a} + {b}",
            "Preço A":             round(pa, 2),
            "Preço B":             round(pb, 2),
            "Total Individual":    round(total, 2),
            "Combo c/ 5% desc.":   round(total * 0.95, 2),
            "Combo c/ 10% desc.":  round(total * 0.90, 2),
            "Frequência":          int(row["Frequência"]),
        })
    return pd.DataFrame(rows)


#  PROMPT 46 — Metas mensais 
def calc_metas(df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    prod = (
        df.groupby("xProd")
        .agg(volume=("qCom", "sum"), receita=("vProd", "sum"),
             frequencia=("chave", "nunique"))
        .reset_index()
        .sort_values("receita", ascending=False)
        .head(top_n)
    )
    prod["Meta +10%"] = (prod["receita"] * 1.10).round(2)
    prod["Meta +20%"] = (prod["receita"] * 1.20).round(2)
    return prod


# ── Vendas por horário (transações + receita por hora/turno) ─
def calc_vendas_horario(df: pd.DataFrame):
    """Retorna agregação por hora e por turno para análise de fluxo."""
    if "hora" not in df.columns:
        return None, None

    notas = df.drop_duplicates("chave").copy()

    # Por hora
    por_hora = (
        notas.groupby("hora")
        .agg(transacoes=("chave", "count"), receita=("vNF", "sum"))
        .reset_index()
    )
    por_hora["ticket_medio"] = por_hora["receita"] / por_hora["transacoes"]
    por_hora["turno"] = por_hora["hora"].apply(
        lambda h: "Manhã" if 5 <= h < 12 else ("Tarde" if 12 <= h < 18 else "Noite")
    )
    media_geral = por_hora["transacoes"].mean()
    por_hora["pct_media"] = por_hora["transacoes"] / media_geral * 100

    # Por turno
    por_turno = (
        notas.groupby("turno")
        .agg(transacoes=("chave", "count"), receita=("vNF", "sum"))
        .reset_index()
    )
    por_turno["pct"] = por_turno["transacoes"] / por_turno["transacoes"].sum() * 100
    por_turno["ticket_medio"] = por_turno["receita"] / por_turno["transacoes"]
    ordem = {"Manhã": 0, "Tarde": 1, "Noite": 2}
    por_turno["_ord"] = por_turno["turno"].map(ordem)
    por_turno = por_turno.sort_values("_ord").drop(columns="_ord").reset_index(drop=True)

    return por_hora, por_turno


#  PROMPT 22 — Horas com oportunidade
def calc_horas_oportunidade(df: pd.DataFrame):
    if "hora" not in df.columns:
        return None
    hourly = (
        df.drop_duplicates("chave")
        .groupby("hora")
        .agg(notas=("chave", "count"), receita=("vNF", "sum"))
        .reset_index()
    )
    media = hourly["notas"].mean()
    hourly["Status"] = hourly["notas"].apply(
        lambda n: "Potencial inexplorado" if n < media * 0.5
        else ("Abaixo da média" if n < media * 0.8 else "Ativo")
    )
    return hourly


# 
# FORMATAÇÃO
# 
def brl(v: float) -> str:
    # Formato brasileiro: R$ 1.096.468,19
    s = f"{v:,.2f}" # "1,096,468.19" (padrão americano)
    s = s.replace(",", "X")   # "1X096X468.19"
    s = s.replace(".", ",")   # "1X096X468,19"
    s = s.replace("X", ".")   # "1.096.468,19"
    return f"R$ {s}"


def fmt_num(v: float) -> str:
    s = f"{int(v):,}".replace(",", ".")
    return s


# 
# GRÁFICOS
# 
def fig_categorias(df_cat: pd.DataFrame):
    fig = px.bar(
        df_cat.sort_values("receita"),
        x="receita", y="categoria", orientation="h",
        color="categoria",
        color_discrete_map=CORES_CATEGORIA,
        text=df_cat.sort_values("receita")["pct"].apply(lambda x: f"{x:.1f}".replace(".", ",") + "%"),
        labels={"receita": "Receita (R$)", "categoria": ""},
        title="Receita por Categoria",
    )
    fig.update_traces(textposition="outside")
    fig.update_layout(showlegend=False, height=420,
                      xaxis_tickformat=",.0f")
    return fig


def fig_bcg(df_bcg: pd.DataFrame):
    cores = {
        "Estrela":     "#F39C12",
        "Vaca":        "#27AE60",
        "Interrogação": "#3498DB",
        "Cão":         "#E74C3C",
    }
    fig = px.scatter(
        df_bcg, x="frequencia", y="receita",
        color="BCG", color_discrete_map=cores,
        hover_data={"xProd": True,
                    "frequencia": True, "receita": ":.2f"},
        title="Matriz BCG — Produtos",
        labels={"frequencia": "Frequência (pedidos)", "receita": "Receita (R$)"},
        size="receita", size_max=40,
    )
    med_f = df_bcg["frequencia"].median()
    med_r = df_bcg["receita"].median()
    fig.add_vline(x=med_f, line_dash="dash", line_color="gray", opacity=0.5)
    fig.add_hline(y=med_r, line_dash="dash", line_color="gray", opacity=0.5)
    max_f, max_r = df_bcg["frequencia"].max(), df_bcg["receita"].max()
    fig.add_annotation(x=max_f * 0.80, y=max_r * 0.95, text="ESTRELAS",
                       showarrow=False, font=dict(size=11, color="#F39C12"))
    fig.add_annotation(x=med_f * 0.1,  y=max_r * 0.95, text="VACAS",
                       showarrow=False, font=dict(size=11, color="#27AE60"))
    fig.add_annotation(x=max_f * 0.80, y=med_r * 0.15, text="INTERROGAÇÃO",
                       showarrow=False, font=dict(size=11, color="#3498DB"))
    fig.add_annotation(x=med_f * 0.1,  y=med_r * 0.15, text="CÃES",
                       showarrow=False, font=dict(size=11, color="#E74C3C"))
    fig.update_layout(height=500)
    return fig


def fig_cesta(df_cesta: pd.DataFrame):
    fig = px.bar(
        df_cesta.head(12), x="Itens/Pedido", y="Nº Pedidos",
        text="% do Total",
        color="Nº Pedidos", color_continuous_scale="Blues",
        title="Distribuição do Tamanho da Cesta de Compras",
        labels={"Itens/Pedido": "Itens por Pedido", "Nº Pedidos": "Nº Pedidos"},
    )
    # Formata porcentagem no padrão brasileiro (ex: 25,8%)
    df_cesta = df_cesta.copy()
    df_cesta["pct_br"] = df_cesta["% do Total"].apply(lambda x: f"{x:.1f}".replace(".", ",") + "%")
    fig.data[0].text = df_cesta.head(12)["pct_br"].tolist()
    fig.update_traces(textposition="outside")
    fig.update_layout(showlegend=False, height=380, coloraxis_showscale=False)
    return fig


def fig_crossell(df_cross: pd.DataFrame):
    fig = px.bar(
        df_cross.sort_values("Frequência"),
        x="Frequência", y=df_cross.sort_values("Frequência").apply(
            lambda r: f"{r['Categoria A']} + {r['Categoria B']}", axis=1),
        orientation="h", color="Frequência",
        color_continuous_scale="Oranges",
        title="Cross-sell entre Categorias",
        labels={"y": "", "Frequência": "Pedidos conjuntos"},
    )
    fig.update_layout(height=420, showlegend=False, coloraxis_showscale=False)
    return fig


# 
# EXPORT EXCEL
# 
def exportar_excel(kpis, df_pares, df_trios,
                   df_cesta, df_bcg, df_abc, df_remocao,
                   df_elev, df_redu, df_sim_preco, df_sim_rec,
                   df_combos, df_metas,
                   cliente: str, periodo: str,
                   sn_result=None) -> bytes:
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        # Resumo Geral
        resumo = pd.DataFrame({
            "Indicador": ["Faturamento Total", "Nº de Pedidos",
                          "Ticket Médio", "Itens por Pedido", "Total de Itens"],
            "Valor": [brl(kpis["faturamento"]), fmt_num(kpis["n_pedidos"]),
                      brl(kpis["ticket_medio"]), f"{kpis['ipc']:.2f}",
                      fmt_num(kpis["total_itens"])],
        })
        resumo.to_excel(writer, sheet_name="Resumo Geral", index=False)

        df_pares.to_excel(writer, sheet_name="Pares de Produtos", index=False)
        df_trios.to_excel(writer, sheet_name="Combos de 3", index=False)
        df_cesta.to_excel(writer, sheet_name="Distribuição Cesta", index=False)

        bcg_exp = df_bcg[["xProd", "BCG", "frequencia", "receita"]].copy()
        bcg_exp["receita"] = bcg_exp["receita"].apply(brl)
        bcg_exp.columns    = ["Produto", "Classificação", "Frequência", "Receita"]
        bcg_exp.to_excel(writer, sheet_name="Classificação Produtos", index=False)

        if df_abc is not None and not df_abc.empty:
            abc_exp = df_abc.copy()
            abc_exp["Receita (R$)"] = abc_exp["Receita (R$)"].apply(brl)
            abc_exp["% Receita"]    = abc_exp["% Receita"].round(2).astype(str) + "%"
            abc_exp["% Acumulado"]  = abc_exp["% Acumulado"].round(2).astype(str) + "%"
            abc_exp.to_excel(writer, sheet_name="Curva ABC", index=False)

        rem_exp = df_remocao.copy()
        rem_exp["receita"] = rem_exp["receita"].apply(brl)
        rem_exp.columns    = ["Produto", "Frequência", "Receita"]
        rem_exp.to_excel(writer, sheet_name="Candidatos Remoção", index=False)

        if not df_elev.empty:
            elev_exp = df_elev.copy()
            for col in ["Ticket Médio c/ Produto", "Ticket Médio Geral", "Diferença R$"]:
                if col in elev_exp.columns:
                    elev_exp[col] = elev_exp[col].apply(brl)
            elev_exp.to_excel(writer, sheet_name="Ticket Drivers Elevam", index=False)

        if not df_redu.empty:
            redu_exp = df_redu.copy()
            for col in ["Ticket Médio c/ Produto", "Ticket Médio Geral", "Diferença R$"]:
                if col in redu_exp.columns:
                    redu_exp[col] = redu_exp[col].apply(brl)
            redu_exp.to_excel(writer, sheet_name="Ticket Drivers Reduzem", index=False)

        if not df_sim_rec.empty:
            sim_rec_exp = df_sim_rec.copy()
            for col in ["Impacto Mensal", "Impacto Anual"]:
                sim_rec_exp[col] = sim_rec_exp[col].apply(brl)
            sim_rec_exp.to_excel(writer, sheet_name="Simulação Receita", index=False)

        if not df_sim_preco.empty:
            sp = df_sim_preco.copy()
            for col in ["receita", "+10% preço (-5% vol)", "+15% preço (-5% vol)", "+20% preço (-5% vol)"]:
                if col in sp.columns:
                    sp[col] = sp[col].apply(brl)
            sp.to_excel(writer, sheet_name="Simulação Preços", index=False)

        if not df_combos.empty:
            cb = df_combos.copy()
            for col in ["Preço A", "Preço B", "Total Individual", "Combo c/ 5% desc.", "Combo c/ 10% desc."]:
                if col in cb.columns:
                    cb[col] = cb[col].apply(brl)
            cb.to_excel(writer, sheet_name="Combos Precificados", index=False)

        if not df_metas.empty:
            mt = df_metas.copy()
            for col in ["receita", "Meta +10%", "Meta +20%"]:
                if col in mt.columns:
                    mt[col] = mt[col].apply(brl)
            mt.to_excel(writer, sheet_name="Metas por Produto", index=False)

        # ── Simples Nacional ──────────────────────────────────────────
        if sn_result is not None and sn_result.get("status") not in (None, "SEM_DADOS"):
            _sn_resumo = pd.DataFrame({
                "Indicador": [
                    "Faturamento Total",
                    "Compras de Comercialização",
                    "% do Faturamento",
                    "Limite Legal",
                    "Status",
                ],
                "Valor": [
                    brl(kpis["faturamento"]),
                    brl(sn_result["total_compras_comercializacao"]),
                    f"{sn_result['pct_faturamento']:.2f}%",
                    "80,00%",
                    sn_result["status"],
                ],
            })
            _sn_resumo.to_excel(writer, sheet_name="SN Resumo", index=False)

            _df_cfop = sn_result.get("df_por_cfop", pd.DataFrame())
            if not _df_cfop.empty:
                _sn_cfop = _df_cfop.copy()
                _sn_cfop["total_compras"] = _sn_cfop["total_compras"].apply(brl)
                _sn_cfop.columns = ["CFOP", "Total Compras", "Notas", "Itens"]
                _sn_cfop.to_excel(writer, sheet_name="SN Por CFOP", index=False)

            _df_forn = sn_result.get("df_por_fornecedor", pd.DataFrame())
            if not _df_forn.empty:
                _sn_forn = _df_forn.copy()
                _sn_forn["total_compras"] = _sn_forn["total_compras"].apply(brl)
                _sn_forn.columns = ["Fornecedor", "Total Compras", "Notas"]
                _sn_forn.to_excel(writer, sheet_name="SN Fornecedores", index=False)

            _df_items = sn_result.get("df_entradas_filtradas", pd.DataFrame())
            if not _df_items.empty:
                _sn_items = _df_items[["chave", "nNF", "dhEmi", "emitente", "CFOP",
                                       "xProd", "qCom", "vUnCom", "vProd"]].copy()
                _sn_items.columns = ["Chave NF-e", "Nº NF", "Emissão", "Fornecedor",
                                     "CFOP", "Produto", "Qtd", "Vl Unit", "Vl Total"]
                _sn_items.to_excel(writer, sheet_name="SN Itens Compra", index=False)

    return buf.getvalue()


def pptx_para_pdf(pptx_bytes: bytes):
    """Converte cada slide do PPTX em imagem PNG e monta PDF via reportlab.
    Funciona sem dependências de sistema (sem LibreOffice).
    """
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches
        from matplotlib.backends.backend_pdf import PdfPages
        from pptx import Presentation
        from pptx.util import Emu
        import io as _io

        prs = Presentation(_io.BytesIO(pptx_bytes))
        slide_w = prs.slide_width.inches   # normalmente 13.33
        slide_h = prs.slide_height.inches  # normalmente 7.5

        buf_pdf = _io.BytesIO()
        with PdfPages(buf_pdf) as pdf:
            for slide in prs.slides:
                # Renderiza o slide como imagem via thumbnail usando python-pptx + pillow
                from pptx.util import Inches
                dpi = 150
                px_w = int(slide_w * dpi)
                px_h = int(slide_h * dpi)

                # Gera imagem do slide usando matplotlib para montar cada shape
                fig = plt.figure(figsize=(slide_w, slide_h), facecolor="white")
                ax = fig.add_axes([0, 0, 1, 1])
                ax.set_xlim(0, slide_w)
                ax.set_ylim(0, slide_h)
                ax.set_aspect("equal")
                ax.axis("off")
                ax.invert_yaxis()

                for shape in slide.shapes:
                    left   = shape.left.inches   if shape.left   else 0
                    top    = shape.top.inches    if shape.top    else 0
                    width  = shape.width.inches  if shape.width  else 0
                    height = shape.height.inches if shape.height else 0

                    # Retângulo colorido
                    if shape.shape_type == 1:  # MSO_SHAPE_TYPE.AUTO_SHAPE
                        try:
                            fill = shape.fill
                            if fill.type == 1:  # SOLID
                                rgb = fill.fore_color.rgb
                                clr = (rgb[0]/255, rgb[1]/255, rgb[2]/255)
                                rect = plt.Rectangle((left, top), width, height,
                                                     facecolor=clr, edgecolor="none",
                                                     transform=ax.transData)
                                ax.add_patch(rect)
                        except Exception:
                            pass

                    # Texto
                    if shape.has_text_frame:
                        try:
                            tf = shape.text_frame
                            full_text = tf.text.strip()
                            if not full_text:
                                continue
                            # Cor do primeiro run
                            p0 = tf.paragraphs[0]
                            run0 = p0.runs[0] if p0.runs else None
                            if run0 and run0.font.color and run0.font.color.type:
                                rgb = run0.font.color.rgb
                                txt_clr = (rgb[0]/255, rgb[1]/255, rgb[2]/255)
                            else:
                                txt_clr = (0, 0, 0)
                            fs_pt = run0.font.size.pt if (run0 and run0.font.size) else 10
                            bold = run0.font.bold if run0 else False
                            # alinhamento
                            from pptx.enum.text import PP_ALIGN
                            align_map = {PP_ALIGN.CENTER: "center",
                                         PP_ALIGN.RIGHT: "right"}
                            ha = align_map.get(p0.alignment, "left")
                            cx = left + width / 2 if ha == "center" else (
                                 left + width     if ha == "right"  else left + 0.05)
                            cy = top + height / 2
                            ax.text(cx, cy, full_text,
                                    color=txt_clr,
                                    fontsize=fs_pt * 0.75,
                                    fontweight="bold" if bold else "normal",
                                    ha=ha, va="center",
                                    wrap=True,
                                    transform=ax.transData,
                                    clip_on=True)
                        except Exception:
                            pass

                    # Imagem (gráficos matplotlib embutidos)
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            from PIL import Image
                            img_bytes = _io.BytesIO(shape.image.blob)
                            img = Image.open(img_bytes)
                            ax.imshow(img,
                                      extent=[left, left + width,
                                              top + height, top],
                                      aspect="auto", zorder=5)
                        except Exception:
                            pass

                pdf.savefig(fig, bbox_inches="tight", dpi=dpi)
                plt.close(fig)

        buf_pdf.seek(0)
        return buf_pdf.getvalue()
    except Exception:
        return None


def exportar_pdf(kpis, kpis_nfce, df_pares, df_bcg,
                 df_remocao, df_elev, df_redu, df_sim_rec,
                 tem_nfe, df_nfe,
                 cliente: str, periodo: str, fonte_label: str) -> bytes:
    """Gera relatório em PDF usando matplotlib (sem dependências extras)."""
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.backends.backend_pdf import PdfPages
        import matplotlib.patches as mpatches
        from pathlib import Path
    except ImportError:
        return None

    # ── Paleta ──────────────────────────────────────────────────────
    AZUL   = "#1e3a5f"
    AZUL2  = "#2563eb"
    VERDE  = "#10b981"
    CINZA  = "#f8f9fa"
    TEXTO  = "#1f2937"
    SUBTXT = "#6b7280"

    # ── Logo ────────────────────────────────────────────────────────
    logo_path = Path(__file__).parent / "LOGO S FUNDO 2.png"
    logo_img = None
    if logo_path.exists():
        try:
            logo_img = plt.imread(str(logo_path))
        except Exception:
            logo_img = None

    buf = io.BytesIO()

    with PdfPages(buf) as pdf:

        # ── helpers ────────────────────────────────────────────────
        def new_slide(title_text=None, has_logo=True):
            fig = plt.figure(figsize=(13.33, 7.5), facecolor="white")
            # Barra de topo
            ax_bar = fig.add_axes([0, 0.91, 1, 0.09])
            ax_bar.set_facecolor(AZUL)
            ax_bar.axis("off")

            if title_text:
                ax_bar.text(0.02, 0.5, title_text, color="white",
                            fontsize=18, fontweight="bold", va="center",
                            transform=ax_bar.transAxes)
            if has_logo and logo_img is not None:
                ax_logo = fig.add_axes([0.88, 0.91, 0.11, 0.09])
                ax_logo.imshow(logo_img)
                ax_logo.axis("off")

            # Rodapé
            ax_ft = fig.add_axes([0, 0, 1, 0.04])
            ax_ft.set_facecolor(AZUL)
            ax_ft.axis("off")
            ax_ft.text(0.5, 0.5, f"{cliente}  ·  {periodo}  ·  Análise de Vendas CP",
                       color="white", fontsize=12, va="center", ha="center",
                       transform=ax_ft.transAxes)
            return fig

        def add_kpi_box(fig, x, y, w, h, label, value, sub="", color=AZUL2):
            ax = fig.add_axes([x, y, w, h])
            ax.set_facecolor(color)
            ax.axis("off")
            ax.text(0.5, 0.72, value, color="white", fontsize=18,
                    fontweight="bold", ha="center", va="center",
                    transform=ax.transAxes)
            ax.text(0.5, 0.35, label, color="white", fontsize=12,
                    ha="center", va="center", transform=ax.transAxes,
                    fontweight="bold")
            if sub:
                ax.text(0.5, 0.12, sub, color=(1, 1, 1, 0.75),
                        fontsize=11, ha="center", va="center",
                        transform=ax.transAxes)

        def draw_table(fig, rect, data, col_headers, col_widths=None,
                       title=None, row_colors=None):
            ax = fig.add_axes(rect)
            ax.axis("off")
            if title:
                ax.text(0, 1.04, title, fontsize=13, fontweight="bold",
                        color=AZUL, transform=ax.transAxes, va="bottom")
            n_rows = len(data)
            n_cols = len(col_headers)
            if col_widths is None:
                col_widths = [1/n_cols] * n_cols
            tbl = ax.table(
                cellText=data,
                colLabels=col_headers,
                cellLoc="center",
                loc="center",
                colWidths=col_widths,
            )
            tbl.auto_set_font_size(False)
            tbl.set_fontsize(12)
            tbl.scale(1, 1.4)
            for (r, c), cell in tbl.get_celld().items():
                if r == 0:
                    cell.set_facecolor(AZUL)
                    cell.set_text_props(color="white", fontweight="bold")
                    cell.set_edgecolor("white")
                else:
                    bg = "#eef2ff" if r % 2 == 0 else "white"
                    if row_colors and r-1 < len(row_colors):
                        bg = row_colors[r-1]
                    cell.set_facecolor(bg)
                    cell.set_edgecolor("#e5e7eb")

        # ══════════════════════════════════════════════════════════
        # Slide 1 — Capa
        # ══════════════════════════════════════════════════════════
        fig = plt.figure(figsize=(13.33, 7.5), facecolor=AZUL)
        if logo_img is not None:
            ax_logo = fig.add_axes([0.38, 0.58, 0.24, 0.22])
            ax_logo.imshow(logo_img)
            ax_logo.axis("off")
        ax_c = fig.add_axes([0, 0.25, 1, 0.35])
        ax_c.set_facecolor(AZUL)
        ax_c.axis("off")
        ax_c.text(0.5, 0.75, "ANÁLISE ESTRATÉGICA DE VENDAS",
                  color="white", fontsize=24, fontweight="bold",
                  ha="center", va="center", transform=ax_c.transAxes)
        ax_c.text(0.5, 0.48, cliente.upper(),
                  color="#60a5fa", fontsize=20, fontweight="bold",
                  ha="center", va="center", transform=ax_c.transAxes)
        ax_c.text(0.5, 0.24, periodo,
                  color="white", fontsize=18, ha="center", va="center",
                  transform=ax_c.transAxes, alpha=0.8)
        ax_ft2 = fig.add_axes([0, 0, 1, 0.06])
        ax_ft2.set_facecolor("#0f1f3a")
        ax_ft2.axis("off")
        ax_ft2.text(0.5, 0.5, "Análise de Vendas CP  ·  Dados Fiscais NFC-e",
                    color="white", fontsize=13, ha="center", va="center",
                    transform=ax_ft2.transAxes, alpha=0.7)
        pdf.savefig(fig, bbox_inches="tight")
        plt.close(fig)

        # ══════════════════════════════════════════════════════════
        # Slide 2 — KPIs principais
        # ══════════════════════════════════════════════════════════
        fig = new_slide("INDICADORES GERAIS")
        fat_nfce_v = kpis_nfce["faturamento"]
        fat_total  = kpis["faturamento"]
        fat_nfe_v  = fat_total - fat_nfce_v if tem_nfe else 0

        kpi_items = [
            (brl(fat_nfce_v),              "FATURAMENTO NFC-e",    fonte_label, AZUL),
            (fmt_num(kpis["n_pedidos"]),   "PEDIDOS",              "",          AZUL2),
            (brl(kpis["ticket_medio"]),    "TICKET MÉDIO",         "",          "#0369a1"),
            (f"{kpis['ipc']:.2f}".replace(".", ","), "ITENS / PEDIDO", "", "#065f46"),
        ]
        if tem_nfe and fat_nfe_v > 0:
            kpi_items[0] = (brl(fat_total), "FATURAMENTO TOTAL", "NFC-e + NF-e", AZUL)

        n = len(kpi_items)
        box_w = 0.22
        gap   = (1 - n * box_w) / (n + 1)
        for i, (val, lbl, sub, clr) in enumerate(kpi_items):
            x = gap + i * (box_w + gap)
            add_kpi_box(fig, x, 0.50, box_w, 0.32, lbl, val, sub, clr)

        # Top produtos no slide de resumo
        if df_bcg is not None and not df_bcg.empty:
            top5_r = df_bcg.nlargest(5, "receita")
            ax_b = fig.add_axes([0.05, 0.10, 0.52, 0.36])
            ax_b.set_facecolor("white")
            bars = ax_b.barh(top5_r["xProd"].str[:22][::-1], top5_r["receita"][::-1],
                             color=AZUL2, edgecolor="none", height=0.6)
            ax_b.set_xlabel("Receita (R$)", fontsize=12, color=SUBTXT)
            ax_b.tick_params(labelsize=9)
            ax_b.spines[["top","right","left"]].set_visible(False)
            ax_b.set_title("Top 5 Produtos", fontsize=12, fontweight="bold",
                           color=AZUL, pad=6)
            for bar in bars:
                w = bar.get_width()
                ax_b.text(w * 1.01, bar.get_y() + bar.get_height()/2,
                          brl(w), va="center", fontsize=11, color=TEXTO)

        pdf.savefig(fig, bbox_inches="tight")
        plt.close(fig)

        # ══════════════════════════════════════════════════════════
        # Slide 3 — Top Produtos
        # ══════════════════════════════════════════════════════════
        fig = new_slide("TOP PRODUTOS")
        # Top 20 produtos por receita
        from collections import defaultdict
        if df_bcg is not None and not df_bcg.empty:
            top_prod = df_bcg.nlargest(20, "receita")[["xProd", "BCG", "frequencia", "receita"]].copy()
            tbl_data = [
                [str(i+1),
                 row["xProd"][:50] if len(row["xProd"]) > 50 else row["xProd"],
                 row["BCG"],
                 fmt_num(int(row["frequencia"])),
                 brl(row["receita"])]
                for i, (_, row) in enumerate(top_prod.iterrows())
            ]
            draw_table(fig, [0.03, 0.07, 0.94, 0.80], tbl_data,
                       ["#", "Produto", "BCG", "Frequência", "Receita"],
                       col_widths=[0.04, 0.54, 0.14, 0.14, 0.14],
                       title="Top 20 Produtos por Receita")

        pdf.savefig(fig, bbox_inches="tight")
        plt.close(fig)

        # ══════════════════════════════════════════════════════════
        # Slide 4 — Market Basket
        # ══════════════════════════════════════════════════════════
        _n_p = len(df_pares) if df_pares is not None else 0
        fig = new_slide(f"MARKET BASKET — TOP {_n_p} PARES MAIS COMPRADOS JUNTOS")
        if df_pares is not None and not df_pares.empty:
            half  = (_n_p + 1) // 2  # split point for two columns
            cols_pares = list(df_pares.columns)
            if _n_p <= 10:
                tbl_data = [[str(i+1)] + [str(v) for v in row]
                            for i, row in enumerate(df_pares.values.tolist())]
                draw_table(fig, [0.03, 0.07, 0.94, 0.80], tbl_data,
                           ["#"] + cols_pares,
                           title=f"Top {_n_p} Pares de Produtos")
            else:
                left_df  = df_pares.iloc[:half]
                right_df = df_pares.iloc[half:]
                tbl_l = [[str(i+1)]       + [str(v) for v in row] for i, row in enumerate(left_df.values.tolist())]
                tbl_r = [[str(i+1+half)]  + [str(v) for v in row] for i, row in enumerate(right_df.values.tolist())]
                draw_table(fig, [0.02, 0.07, 0.47, 0.80], tbl_l,
                           ["#"] + cols_pares, title=f"Pares 1–{half}")
                draw_table(fig, [0.51, 0.07, 0.47, 0.80], tbl_r,
                           ["#"] + cols_pares, title=f"Pares {half+1}–{_n_p}")

        pdf.savefig(fig, bbox_inches="tight")
        plt.close(fig)

        # ══════════════════════════════════════════════════════════
        # Slide 5 — Ticket Drivers + Simulações
        # ══════════════════════════════════════════════════════════
        fig = new_slide("TICKET DRIVERS E SIMULAÇÕES")

        if df_elev is not None and not df_elev.empty:
            top_elev = df_elev.head(12)
            tbl_data = [[str(i+1)] + [str(v) for v in row]
                        for i, row in enumerate(top_elev.values.tolist())]
            draw_table(fig, [0.03, 0.50, 0.55, 0.37], tbl_data,
                       ["#"] + list(top_elev.columns),
                       title="Produtos que Elevam o Ticket")

        if df_sim_rec is not None and not df_sim_rec.empty:
            tbl_data = [[str(v) for v in row]
                        for row in df_sim_rec.values.tolist()]
            draw_table(fig, [0.62, 0.50, 0.35, 0.37], tbl_data,
                       list(df_sim_rec.columns),
                       title="Simulação de Crescimento")

        if df_remocao is not None and not df_remocao.empty:
            top_rem = df_remocao.head(12)
            rem_data = [
                [str(i+1),
                 row.xProd[:45] if len(str(row.xProd)) > 45 else str(row.xProd),
                 fmt_num(int(row.frequencia)), brl(float(row.receita))]
                for i, row in enumerate(top_rem.itertuples(index=False))
            ]
            draw_table(fig, [0.03, 0.07, 0.55, 0.37], rem_data,
                       ["#", "Produto", "Freq.", "Receita"],
                       col_widths=[0.05, 0.60, 0.17, 0.18],
                       title="Candidatos a Remoção do Cardápio")

        pdf.savefig(fig, bbox_inches="tight")
        plt.close(fig)

    buf.seek(0)
    return buf.getvalue()


#
# EXPORT PPTX
#
def exportar_pptx(kpis, df_pares, df_trios,
                  df_cesta, df_turno, df_bcg,
                  df_elev, df_redu, df_sim_rec, df_sim_preco,
                  df_combos, df_metas, df_horas,
                  df_solo, df_remocao, df_dia_tipo,
                  df_nfe, kpis_nfce,
                  fonte_label: str, cliente: str, periodo: str,
                  df_abc=None,
                  df_por_hora=None, df_por_turno=None,
                  sn_result=None) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import matplotlib.patches as mpatches
    except ImportError:
        return None

    AZUL_ESC  = RGBColor(0x1E, 0x3A, 0x5F)
    AZUL_MED  = RGBColor(0x25, 0x63, 0xEB)
    BRANCO    = RGBColor(0xFF, 0xFF, 0xFF)
    CINZA_CLR = RGBColor(0xF8, 0xF9, 0xFA)
    LARANJA   = RGBColor(0xE6, 0x7E, 0x22)
    TEXTO     = RGBColor(0x1F, 0x29, 0x37)

    def pct_br(v: float) -> str:
        """Percentual no padrão brasileiro: 25,8%"""
        return f"{v:.1f}".replace(".", ",") + "%"

    def brl_k(v: float) -> str:
        """Valor abreviado em mil: R$ 264,6 mil"""
        if v >= 1_000_000:
            return f"R$ {v/1_000_000:.1f}".replace(".", ",") + " mi"
        elif v >= 1_000:
            return f"R$ {v/1_000:.1f}".replace(".", ",") + " mil"
        return brl(v)

    W = Inches(13.33)
    H = Inches(7.5)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]  # Blank

    def add_rect(slide, left, top, width, height, fill_color, alpha=None):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        return shape

    def add_text(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=None,
                 align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tb.word_wrap = wrap
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return tb

    def plt_to_pptx_image(fig_mpl, slide, left, top, width, height):
        buf2 = io.BytesIO()
        fig_mpl.savefig(buf2, format="png", bbox_inches="tight",
                        dpi=150, facecolor="white")
        buf2.seek(0)
        slide.shapes.add_picture(buf2, left, top, width, height)
        plt.close(fig_mpl)

    #  SLIDE 1: CAPA 
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, H, AZUL_ESC)
    add_rect(sl, 0, Inches(5.2), W, Inches(2.3), AZUL_MED)

    # Caixas de texto com largura total do slide — centralização garantida
    add_text(sl, cliente.upper(),
             Inches(0), Inches(1.5), W, Inches(1.4),
             font_size=42, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
    add_text(sl, "ANÁLISE ESTRATÉGICA DE VENDAS",
             Inches(0), Inches(3.0), W, Inches(0.7),
             font_size=24, bold=False, color=RGBColor(0xCA, 0xDC, 0xFC),
             align=PP_ALIGN.CENTER)
    add_text(sl, periodo.upper() + f"  |  {fonte_label.upper()}",
             Inches(0), Inches(3.75), W, Inches(0.55),
             font_size=16, color=RGBColor(0xCA, 0xDC, 0xFC),
             align=PP_ALIGN.CENTER)
    add_text(sl, "Análise de dados orientada ao crescimento: clareza para decidir, estratégia para vender mais.",
             Inches(0), Inches(5.5), W, Inches(0.75),
             font_size=14, color=BRANCO, align=PP_ALIGN.CENTER)

    # Logo centralizada na parte inferior da capa
    try:
        from pathlib import Path as _PL
        _logo_p = _PL(__file__).parent / "LOGO S FUNDO 2.png"
        if _logo_p.exists():
            _lw = Inches(2.8)
            _lh = Inches(1.1)
            sl.shapes.add_picture(
                str(_logo_p),
                (W - _lw) / 2,   # centralizada horizontalmente
                Inches(6.15),    # parte inferior do slide
                _lw, _lh,
            )
    except Exception:
        pass

    #  SLIDE 2: PANORAMA GERAL
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, Inches(1.2), AZUL_ESC)
    add_text(sl, f"PANORAMA GERAL  |  {periodo.upper()}",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
             font_size=22, bold=True, color=BRANCO)

    cards = [
        (brl(kpis["faturamento"]),  "FATURAMENTO TOTAL"),
        (fmt_num(kpis["n_pedidos"]), "PEDIDOS EMITIDOS"),
        (brl(kpis["ticket_medio"]), "TICKET MÉDIO"),
        (f"{kpis['ipc']:.2f}".replace(".", ","), "ITENS POR PEDIDO EM MÉDIA"),
    ]
    card_w = Inches(2.8)
    card_h = Inches(2.2)
    card_top = Inches(1.6)
    cores_cards = [AZUL_MED, RGBColor(0x16, 0xA0, 0x85),
                   RGBColor(0x8E, 0x44, 0xAD), RGBColor(0xD3, 0x54, 0x00)]
    for i, (val, lbl) in enumerate(cards):
        left = Inches(0.4 + i * 3.1)
        add_rect(sl, left, card_top, card_w, card_h, cores_cards[i])
        add_text(sl, val, left + Inches(0.1), card_top + Inches(0.3),
                 card_w - Inches(0.2), Inches(1),
                 font_size=28, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, left + Inches(0.1), card_top + Inches(1.3),
                 card_w - Inches(0.2), Inches(0.6),
                 font_size=15, color=BRANCO, align=PP_ALIGN.CENTER)

    sub_cards = [
        (fmt_num(kpis["total_itens"]), "TOTAL DE ITENS"),
        (fmt_num(len(df_bcg)) if df_bcg is not None else "—", "PRODUTOS ÚNICOS"),
        (fmt_num(df_bcg[df_bcg["BCG"]=="Estrela"].shape[0]) if df_bcg is not None else "—", "PRODUTOS ESTRELA"),
    ]
    for i, (val, lbl) in enumerate(sub_cards):
        left = Inches(1.0 + i * 3.8)
        add_rect(sl, left, Inches(4.3), Inches(3.2), Inches(1.6),
                 RGBColor(0xEC, 0xF0, 0xF1))
        add_text(sl, val, left + Inches(0.1), Inches(4.45),
                 Inches(3.0), Inches(0.7),
                 font_size=22, bold=True, color=AZUL_ESC, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, left + Inches(0.1), Inches(5.1),
                 Inches(3.0), Inches(0.5),
                 font_size=14, color=RGBColor(0x7F, 0x8C, 0x8D),
                 align=PP_ALIGN.CENTER)

    add_text(sl, f"Fonte: {fonte_label} — {periodo}",
             Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
             font_size=13, color=RGBColor(0x9E, 0x9E, 0x9E))

    #  SLIDE 3: PARES DE PRODUTOS
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)

    _n_pares = len(df_pares)
    add_text(sl, f"TOP {_n_pares} PARES DE PRODUTOS MAIS COMPRADOS JUNTOS",
             Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.8),
             font_size=22, bold=True, color=BRANCO)

    def _trunc(txt, n=30):
        return txt if len(txt) <= n else txt[:n].rstrip() + "…"

    # helper: draw one pares table at given x offset, items offset num_offset
    def _draw_pares_tbl(rows_df, left_x, y0, col_ws, num_offset=0):
        rh = Inches(0.41)
        hdrs_p = ["#", "Produto A", "Produto B", "Freq."]
        x = left_x
        for hdr, w in zip(hdrs_p, col_ws):
            add_rect(sl, x, y0, w, rh, AZUL_ESC)
            add_text(sl, hdr, x + Inches(0.04), y0 + Inches(0.06),
                     w - Inches(0.08), rh - Inches(0.1),
                     font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
            x += w
        for i, row in enumerate(rows_df.itertuples()):
            y  = y0 + rh * (i + 1)
            bg = CINZA_CLR if i % 2 == 0 else BRANCO
            x  = left_x
            vals = [str(i + 1 + num_offset),
                    str(row._1)[:38], str(row._2)[:38],
                    fmt_num(getattr(row, "Frequência", getattr(row, "_3", 0)))]
            for j, (val, w) in enumerate(zip(vals, col_ws)):
                add_rect(sl, x, y, w, rh, bg)
                al = PP_ALIGN.CENTER if j in (0, 3) else PP_ALIGN.LEFT
                add_text(sl, val, x + Inches(0.05), y + Inches(0.06),
                         w - Inches(0.10), rh - Inches(0.08),
                         font_size=12, color=TEXTO, align=al)
                x += w

    # Footer insights — empilhados verticalmente para evitar overflow
    insights = []
    if not df_pares.empty:
        top1 = df_pares.iloc[0]
        pa = _trunc(top1["Produto A"])
        pb = _trunc(top1["Produto B"])
        freq1 = fmt_num(top1["Frequência"])
        insights.append(f"{pa} + {pb}: par mais frequente com {freq1} ocorrencias")
    insights.append("Identifique o produto ancora — aquele que aparece na maioria dos pares")
    insights.append("Crie kits e combos baseados nos pares mais frequentes")

    _ins_row_h = Inches(0.30)
    footer_h   = Inches(0.28) + len(insights) * _ins_row_h + Inches(0.14)
    footer_y   = Inches(7.5) - footer_h - Inches(0.08)   # âncora no rodapé
    add_rect(sl, Inches(0.2), footer_y, Inches(13.0), footer_h,
             RGBColor(0xEB, 0xF5, 0xFF))
    add_text(sl, "INSIGHTS DE CROSS-SELL",
             Inches(0.35), footer_y + Inches(0.04), Inches(5.0), Inches(0.26),
             font_size=11, bold=True, color=AZUL_ESC)
    for k, ins in enumerate(insights):
        y_ins = footer_y + Inches(0.30) + k * _ins_row_h
        add_text(sl, f"•  {ins}",
                 Inches(0.40), y_ins, Inches(12.6), _ins_row_h,
                 font_size=11, color=TEXTO)

    # Table(s)
    y_tbl = Inches(1.15)
    if _n_pares <= 10:
        # Single wide table — full width
        cw1 = [Inches(0.5), Inches(4.5), Inches(4.5), Inches(1.4)]
        _draw_pares_tbl(df_pares, Inches(1.2), y_tbl, cw1, num_offset=0)
    else:
        # Two tables side by side — left: 1–10, right: 11–N
        cw2 = [Inches(0.45), Inches(2.35), Inches(2.35), Inches(0.75)]  # 5.9" each
        _draw_pares_tbl(df_pares.head(10),    Inches(0.3), y_tbl, cw2, num_offset=0)
        _draw_pares_tbl(df_pares.iloc[10:20], Inches(6.8), y_tbl, cw2, num_offset=10)

    #  SLIDE 5: COMBOS DE 3 
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
    add_text(sl, "TOP 10 COMBINAÇÕES DE 3 PRODUTOS",
             Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=BRANCO)

    rows_t = df_trios.head(10)
    row_h   = Inches(0.42)
    col_w2  = [Inches(0.5), Inches(2.8), Inches(2.8), Inches(2.8), Inches(1.0), Inches(2.0)]
    hdrs2   = ["#", "Produto A", "Produto B", "Produto C", "Freq.", "Sugestão"]
    x = Inches(0.3)
    for hdr, w in zip(hdrs2, col_w2):
        add_rect(sl, x, Inches(1.15), w, row_h, AZUL_ESC)
        add_text(sl, hdr, x + Inches(0.05), Inches(1.22),
                 w - Inches(0.1), row_h - Inches(0.1),
                 font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        x += w

    for i, row in enumerate(rows_t.itertuples()):
        y  = Inches(1.15) + row_h * (i + 1)
        bg = CINZA_CLR if i % 2 == 0 else BRANCO
        x  = Inches(0.3)
        vals = [str(i + 1), row._1, row._2, row._3,
                fmt_num(row.Frequência), row.Sugestão]
        for j, (val, w) in enumerate(zip(vals, col_w2)):
            add_rect(sl, x, y, w, row_h, bg)
            al = PP_ALIGN.CENTER if j in (0, 4) else PP_ALIGN.LEFT
            add_text(sl, val, x + Inches(0.05), y + Inches(0.07),
                     w - Inches(0.1), row_h - Inches(0.1),
                     font_size=12, color=TEXTO, align=al)
            x += w

    add_text(sl, "Estes combos ocorrem naturalmente — crie promoções formais para ampliar a recorrência e elevar o ticket médio",
             Inches(0.3), Inches(6.9), Inches(12), Inches(0.4),
             font_size=13, color=RGBColor(0x7F, 0x8C, 0x8D))

    #  SLIDE 6: DISTRIBUIÇÃO DA CESTA 
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
    add_text(sl, "DISTRIBUIÇÃO DO TAMANHO DA CESTA DE COMPRAS",
             Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=BRANCO)

    fig_m2, ax2 = plt.subplots(figsize=(7.5, 5.0))
    dc = df_cesta.head(10)
    bars2 = ax2.bar(dc["Itens/Pedido"].astype(str), dc["Nº Pedidos"],
                    color="#2563EB", edgecolor="white", width=0.7)
    max_n = max(dc["Nº Pedidos"])
    for bar, pct in zip(bars2, dc["% do Total"]):
        ax2.text(bar.get_x() + bar.get_width() / 2,
                 bar.get_height() + max_n * 0.015,
                 pct_br(pct), ha="center", fontsize=10, fontweight="bold", color="#1F2937")
    ax2.set_xlabel("Itens por Pedido", fontsize=11)
    ax2.set_ylabel("Nº de Pedidos", fontsize=11)
    ax2.tick_params(labelsize=10)
    ax2.set_ylim(0, max_n * 1.15)
    ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{int(x):,}".replace(",", ".")))
    ax2.spines[["top", "right"]].set_visible(False)
    fig_m2.tight_layout()
    plt_to_pptx_image(fig_m2, sl, Inches(0.3), Inches(1.1), Inches(7.8), Inches(6.2))

    # Painel de oportunidade
    cesta_1_2 = df_cesta[df_cesta["Itens/Pedido"] <= 2]["Nº Pedidos"].sum()
    pct_1_2   = cesta_1_2 / df_cesta["Nº Pedidos"].sum() * 100
    tm        = kpis["ticket_medio"]

    add_rect(sl, Inches(8.4), Inches(1.1), Inches(4.6), Inches(6.2),
             RGBColor(0xFF, 0xF8, 0xE1))
    add_text(sl, "OPORTUNIDADE",
             Inches(8.6), Inches(1.3), Inches(4.2), Inches(0.55),
             font_size=15, bold=True, color=LARANJA)
    add_text(sl, pct_br(pct_1_2) + " dos pedidos\ntêm apenas 1 ou 2 itens.",
             Inches(8.6), Inches(2.0), Inches(4.2), Inches(0.9),
             font_size=18, bold=True, color=AZUL_ESC)
    add_text(sl, "Enorme potencial de upsell.",
             Inches(8.6), Inches(3.0), Inches(4.2), Inches(0.5),
             font_size=14, color=TEXTO)
    add_rect(sl, Inches(8.4), Inches(3.65), Inches(4.6), Inches(1.0),
             RGBColor(0xFE, 0xF3, 0xC7))
    add_text(sl, f"Cada +1 item =",
             Inches(8.6), Inches(3.72), Inches(4.2), Inches(0.35),
             font_size=15, color=TEXTO)
    add_text(sl, f"+{brl(tm)} de ticket",
             Inches(8.6), Inches(4.07), Inches(4.2), Inches(0.45),
             font_size=16, bold=True, color=RGBColor(0x10, 0xB9, 0x81))
    add_text(sl, "→ Sugestão no balcão e\ncombos prontos são\nas principais alavancas.",
             Inches(8.6), Inches(4.85), Inches(4.2), Inches(1.2),
             font_size=14, color=TEXTO)

    #  SLIDE 7: TURNO (se disponível) 
    if df_turno is not None:
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
        add_text(sl, "PRODUTOS MAIS VENDIDOS POR TURNO",
                 Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
                 font_size=22, bold=True, color=BRANCO)

        turnos = ["Manhã", "Tarde", "Noite"]
        col_w_t = Inches(4.0)
        for t_i, turno in enumerate(turnos):
            left_t = Inches(0.3 + t_i * 4.3)
            sub = df_turno[df_turno["turno"] == turno].head(7)
            add_rect(sl, left_t, Inches(1.1), col_w_t, Inches(0.55),
                     AZUL_MED if turno == "Manhã" else
                     (RGBColor(0xD3, 0x54, 0x00) if turno == "Tarde" else
                      RGBColor(0x1A, 0x23, 0x4E)))
            icon = "" if turno == "Manhã" else ("" if turno == "Tarde" else "")
            add_text(sl, f"{icon} {turno.upper()}",
                     left_t + Inches(0.1), Inches(1.17),
                     col_w_t - Inches(0.2), Inches(0.4),
                     font_size=15, bold=True, color=BRANCO)
            for k, row in enumerate(sub.itertuples()):
                y_r = Inches(1.65 + k * 0.73)
                add_rect(sl, left_t, y_r, col_w_t, Inches(0.65),
                         CINZA_CLR if k % 2 == 0 else BRANCO)
                add_text(sl, f"{k+1}. {row.xProd}",
                         left_t + Inches(0.1), y_r + Inches(0.05),
                         col_w_t - Inches(0.2), Inches(0.3),
                         font_size=13, bold=True, color=TEXTO)
                add_text(sl, brl(row.receita),
                         left_t + Inches(0.1), y_r + Inches(0.33),
                         col_w_t - Inches(0.2), Inches(0.25),
                         font_size=13, color=LARANJA)

    #  SLIDE 8: PRODUTOS SOLO (ÂNCORA)
    if df_solo is not None and not df_solo.empty:
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
        add_text(sl, "PRODUTOS ÂNCORA — COMPRADOS COMO ÚNICO ITEM",
                 Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
                 font_size=22, bold=True, color=BRANCO)
        add_text(sl, "Clientes que vêm especificamente por esses produtos — máxima oportunidade de upsell",
                 Inches(0.5), Inches(1.05), Inches(12), Inches(0.35),
                 font_size=14, color=RGBColor(0x6B, 0x72, 0x80))

        hdrs_sl  = ["#", "Produto", "Pedidos Solo", "Receita Solo"]
        col_ws_sl = [Inches(0.5), Inches(7.5), Inches(2.0), Inches(2.8)]
        y_sl = Inches(1.5)
        rh   = Inches(0.46)
        x    = Inches(0.3)
        for hdr, ww in zip(hdrs_sl, col_ws_sl):
            add_rect(sl, x, y_sl, ww, rh, AZUL_ESC)
            add_text(sl, hdr, x + Inches(0.05), y_sl + Inches(0.08),
                     ww - Inches(0.1), rh - Inches(0.1),
                     font_size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
            x += ww
        # Máx 9 linhas para caber o insight embaixo
        for ri, (_, row) in enumerate(df_solo.head(9).iterrows()):
            y_r = y_sl + rh * (ri + 1)
            bg  = CINZA_CLR if ri % 2 == 0 else BRANCO
            x   = Inches(0.3)
            vals_sl = [str(ri+1), str(row["xProd"])[:65],
                       fmt_num(row["frequencia"]), brl(row["receita"])]
            for vi, (val, ww) in enumerate(zip(vals_sl, col_ws_sl)):
                add_rect(sl, x, y_r, ww, rh, bg)
                al = PP_ALIGN.LEFT if vi == 1 else PP_ALIGN.CENTER
                add_text(sl, val, x + Inches(0.05), y_r + Inches(0.08),
                         ww - Inches(0.1), rh - Inches(0.1),
                         font_size=13, color=TEXTO, align=al)
                x += ww

        # Insight abaixo da tabela (sem sobreposição)
        insight_top = y_sl + rh * 10 + Inches(0.1)   # logo após a última linha
        top_solo = df_solo.iloc[0]
        add_rect(sl, Inches(0.3), insight_top, Inches(12.7), Inches(0.8), RGBColor(0xEB, 0xF5, 0xFF))
        add_text(sl, f"⭐ {top_solo['xProd']} é o produto mais comprado sozinho — "
                     f"{fmt_num(top_solo['frequencia'])} pedidos. "
                     "Sugira um complemento natural no momento da compra para elevar o ticket.",
                 Inches(0.5), insight_top + Inches(0.08), Inches(12.3), Inches(0.65),
                 font_size=15, bold=True, color=AZUL_ESC)

    #  SLIDE 9: CANDIDATOS A REMOÇÃO 
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
    add_text(sl, "CANDIDATOS A REMOÇÃO DO CARDÁPIO",
             Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=BRANCO)
    add_text(sl, "Produtos com menor receita e baixa frequência — avalie custo antes de remover",
             Inches(0.5), Inches(1.05), Inches(12), Inches(0.35),
             font_size=14, color=RGBColor(0x6B, 0x72, 0x80))

    hdrs_rem  = ["#", "Produto", "Frequência", "Receita"]
    col_ws_rem = [Inches(0.5), Inches(7.8), Inches(1.8), Inches(2.8)]
    y_rm = Inches(1.5)
    rh   = Inches(0.44)
    x    = Inches(0.3)
    for hdr, ww in zip(hdrs_rem, col_ws_rem):
        add_rect(sl, x, y_rm, ww, rh, RGBColor(0xEF, 0x44, 0x44))
        add_text(sl, hdr, x + Inches(0.05), y_rm + Inches(0.07),
                 ww - Inches(0.1), rh - Inches(0.1),
                 font_size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
        x += ww
    # Máx 10 linhas para caber footer sem sobreposição
    for ri, (_, row) in enumerate(df_remocao.head(10).iterrows()):
        y_r = y_rm + rh * (ri + 1)
        bg  = RGBColor(0xFF, 0xF0, 0xF0) if ri % 2 == 0 else BRANCO
        x   = Inches(0.3)
        vals_rem = [str(ri+1), str(row["xProd"])[:70],
                    fmt_num(row["frequencia"]), brl(row["receita"])]
        for vi, (val, ww) in enumerate(zip(vals_rem, col_ws_rem)):
            add_rect(sl, x, y_r, ww, rh, bg)
            al = PP_ALIGN.LEFT if vi == 1 else PP_ALIGN.CENTER
            add_text(sl, val, x + Inches(0.05), y_r + Inches(0.07),
                     ww - Inches(0.1), rh - Inches(0.1),
                     font_size=13, color=TEXTO, align=al)
            x += ww

    # Footer sempre abaixo da última linha da tabela
    footer_top = y_rm + rh * 11 + Inches(0.08)
    add_rect(sl, Inches(0.3), footer_top, Inches(12.7), Inches(0.32), RGBColor(0xFF, 0xF8, 0xE1))
    add_text(sl, "⚠ Antes de remover: verifique custo de produção, sazonalidade e perfil de cliente fiel.",
             Inches(0.5), footer_top + Inches(0.04), Inches(12.3), Inches(0.26),
             font_size=13, color=RGBColor(0x92, 0x40, 0x0E))

    #  SLIDE CURVA ABC
    if df_abc is not None and not df_abc.empty:
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
        add_text(sl, "CURVA ABC — ANÁLISE DE RELEVÂNCIA DOS PRODUTOS",
                 Inches(0.3), Inches(0.1), W - Inches(0.6), Inches(0.8),
                 font_size=22, bold=True, color=BRANCO)

        _abc_cores = {"A": RGBColor(0x05, 0x96, 0x69),   # verde
                      "B": RGBColor(0xD9, 0x77, 0x06),   # âmbar
                      "C": RGBColor(0x6B, 0x72, 0x80)}   # cinza

        _abc_col_x = [Inches(0.2), Inches(4.55), Inches(8.9)]
        _abc_col_w = Inches(4.2)
        _abc_top   = Inches(1.1)
        _fat_total = df_abc["Receita (R$)"].sum()

        for gi, grupo in enumerate(["A", "B", "C"]):
            gdf   = df_abc[df_abc["Curva"] == grupo]
            if gdf.empty:
                continue
            cx    = _abc_col_x[gi]
            cor   = _abc_cores[grupo]
            n_prod   = len(gdf)
            fat_grp  = gdf["Receita (R$)"].sum()
            pct_grp  = fat_grp / _fat_total * 100 if _fat_total else 0

            # Cabeçalho colorido com letra
            add_rect(sl, cx, _abc_top, _abc_col_w, Inches(0.5), cor)
            add_text(sl, f"GRUPO  {grupo}",
                     cx + Inches(0.1), _abc_top + Inches(0.08),
                     _abc_col_w - Inches(0.2), Inches(0.36),
                     font_size=18, bold=True, color=BRANCO)

            # Resumo estatístico
            _stat_top = _abc_top + Inches(0.55)
            add_rect(sl, cx, _stat_top, _abc_col_w, Inches(0.62),
                     RGBColor(0xF3, 0xF4, 0xF6))
            add_text(sl, f"{n_prod} produto(s)  ·  {pct_grp:.1f}% da receita  ·  {brl(fat_grp)}",
                     cx + Inches(0.1), _stat_top + Inches(0.12),
                     _abc_col_w - Inches(0.2), Inches(0.40),
                     font_size=12, color=TEXTO)

            # Top 5 produtos do grupo
            top10 = gdf.head(10)
            _row_top = _stat_top + Inches(0.68)
            _rh_abc  = Inches(0.34)
            hdrs_abc = ["#", "Produto", "Receita", "%"]
            cws_abc  = [Inches(0.38), Inches(2.18), Inches(1.12), Inches(0.52)]
            # Cabeçalho tabela
            _x = cx
            for hdr_a, cw_a in zip(hdrs_abc, cws_abc):
                add_rect(sl, _x, _row_top, cw_a, _rh_abc, cor)
                add_text(sl, hdr_a,
                         _x + Inches(0.03), _row_top + Inches(0.06),
                         cw_a - Inches(0.06), _rh_abc - Inches(0.08),
                         font_size=11, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
                _x += cw_a
            # Linhas de dados
            for ri, (_, row_a) in enumerate(top10.iterrows()):
                _yr  = _row_top + _rh_abc * (ri + 1)
                _bg  = RGBColor(0xEC, 0xFD, 0xF5) if gi == 0 else (
                       RGBColor(0xFF, 0xF9, 0xEB) if gi == 1 else
                       RGBColor(0xF9, 0xFA, 0xFB)) if ri % 2 == 0 else BRANCO
                _x = cx
                _vals_a = [str(int(row_a["Rank"])),
                           str(row_a["Produto"])[:30],
                           brl(row_a["Receita (R$)"]),
                           f"{row_a['% Receita']:.1f}%"]
                for ji, (va, cw_a) in enumerate(zip(_vals_a, cws_abc)):
                    add_rect(sl, _x, _yr, cw_a, _rh_abc, _bg)
                    _al = PP_ALIGN.CENTER if ji in (0, 2, 3) else PP_ALIGN.LEFT
                    add_text(sl, va,
                             _x + Inches(0.03), _yr + Inches(0.06),
                             cw_a - Inches(0.06), _rh_abc - Inches(0.08),
                             font_size=11, color=TEXTO, align=_al)
                    _x += cw_a

        # Legenda rodapé
        add_rect(sl, Inches(0.2), Inches(6.9), W - Inches(0.4), Inches(0.38),
                 RGBColor(0xEB, 0xF5, 0xFF))
        add_text(sl,
                 "A: primeiros 80% da receita (produtos vitais)  "
                 "·  B: 80–95% (produtos importantes)  "
                 "·  C: 95–100% (produtos de menor impacto)",
                 Inches(0.35), Inches(6.92), W - Inches(0.7), Inches(0.32),
                 font_size=11, color=AZUL_ESC)

    #  SLIDE DIA DA SEMANA (se disponível)
    if df_dia_tipo is not None and not df_dia_tipo.empty:
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
        add_text(sl, "VENDAS POR DIA DA SEMANA",
                 Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
                 font_size=22, bold=True, color=BRANCO)

        col_w_dia = Inches(6.1)
        tipos = [("Dia Útil", "", AZUL_MED), ("Final de Semana", "", RGBColor(0x8E, 0x44, 0xAD))]
        for ti, (tipo, icon, cor_tipo) in enumerate(tipos):
            left_d = Inches(0.3 + ti * 6.6)
            sub_d  = df_dia_tipo[df_dia_tipo["tipo_dia"] == tipo].head(8)
            add_rect(sl, left_d, Inches(1.1), col_w_dia, Inches(0.55), cor_tipo)
            add_text(sl, f"{icon} {tipo.upper()}",
                     left_d + Inches(0.1), Inches(1.17),
                     col_w_dia - Inches(0.2), Inches(0.4),
                     font_size=15, bold=True, color=BRANCO)
            for ki, (_, r) in enumerate(sub_d.iterrows()):
                y_r = Inches(1.65 + ki * 0.68)
                add_rect(sl, left_d, y_r, col_w_dia, Inches(0.62),
                         CINZA_CLR if ki % 2 == 0 else BRANCO)
                add_text(sl, f"{ki+1}. {str(r['xProd'])[:52]}",
                         left_d + Inches(0.1), y_r + Inches(0.05),
                         col_w_dia - Inches(1.8), Inches(0.3),
                         font_size=13, bold=True, color=TEXTO)
                add_text(sl, brl(r["receita"]),
                         left_d + Inches(0.1), y_r + Inches(0.33),
                         col_w_dia - Inches(0.2), Inches(0.24),
                         font_size=13, color=LARANJA)
                add_text(sl, fmt_num(r["frequencia"]) + " pedidos",
                         left_d + col_w_dia - Inches(1.7), y_r + Inches(0.18),
                         Inches(1.6), Inches(0.26),
                         font_size=13, color=RGBColor(0x6B, 0x72, 0x80), align=PP_ALIGN.RIGHT)

    #  SLIDE 10: TICKET DRIVERS 
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
    add_text(sl, "PRODUTOS QUE INFLUENCIAM O TICKET MÉDIO",
             Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=BRANCO)

    # Esquerda: elevadores
    add_rect(sl, Inches(0.3), Inches(1.15), Inches(6.1), Inches(0.5),
             RGBColor(0x10, 0xB9, 0x81))
    add_text(sl, "ELEVAM O TICKET",
             Inches(0.4), Inches(1.22), Inches(5.9), Inches(0.38),
             font_size=14, bold=True, color=BRANCO)
    for ki, row in enumerate(df_elev.head(7).iterrows()):
        _, r = row
        y_r = Inches(1.65 + ki * 0.63)
        add_rect(sl, Inches(0.3), y_r, Inches(6.1), Inches(0.58),
                 RGBColor(0xD1, 0xFA, 0xE5) if ki % 2 == 0 else BRANCO)
        add_text(sl, str(r["Produto"])[:40],
                 Inches(0.4), y_r + Inches(0.04), Inches(3.5), Inches(0.28),
                 font_size=13, bold=True, color=TEXTO)
        diff_r = r.get("Diferença R$", 0)
        add_text(sl, f"+{brl(abs(diff_r))} vs média",
                 Inches(0.4), y_r + Inches(0.3), Inches(3.5), Inches(0.24),
                 font_size=12, color=RGBColor(0x10, 0xB9, 0x81))
        add_text(sl, brl(r.get("Ticket Médio c/ Produto", 0)),
                 Inches(3.9), y_r + Inches(0.15), Inches(2.4), Inches(0.28),
                 font_size=14, bold=True, color=AZUL_ESC, align=PP_ALIGN.CENTER)

    # Direita: redutores
    add_rect(sl, Inches(6.8), Inches(1.15), Inches(6.1), Inches(0.5),
             RGBColor(0xEF, 0x44, 0x44))
    add_text(sl, "REDUZEM O TICKET",
             Inches(6.9), Inches(1.22), Inches(5.9), Inches(0.38),
             font_size=14, bold=True, color=BRANCO)
    for ki, row in enumerate(df_redu.head(7).iterrows()):
        _, r = row
        y_r = Inches(1.65 + ki * 0.63)
        add_rect(sl, Inches(6.8), y_r, Inches(6.1), Inches(0.58),
                 RGBColor(0xFE, 0xE2, 0xE2) if ki % 2 == 0 else BRANCO)
        add_text(sl, str(r["Produto"])[:40],
                 Inches(6.9), y_r + Inches(0.04), Inches(3.5), Inches(0.28),
                 font_size=13, bold=True, color=TEXTO)
        add_text(sl, "Ticket médio menor que a média geral",
                 Inches(6.9), y_r + Inches(0.3), Inches(3.5), Inches(0.24),
                 font_size=12, color=RGBColor(0xEF, 0x44, 0x44))

    _tm_nfce_str  = brl(df_elev["Ticket Médio Geral"].iloc[0]) if not df_elev.empty else "—"
    _tm_total_str = brl(kpis["ticket_medio"])
    add_text(sl,
             f"Ticket médio NFC-e (consumidor final): {_tm_nfce_str}",
             Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.38),
             font_size=13, bold=True, color=RGBColor(0x1E, 0x3A, 0x5F))
    add_text(sl,
             (f"Nota: o ticket médio exibido no banner do relatorio ({_tm_total_str}) considera "
              f"NFC-e + NF-e (vendas B2B). Esta analise de drivers usa apenas NFC-e — "
              f"compras do consumidor final — excluindo pedidos corporativos que distorcem o perfil de consumo."),
             Inches(0.3), Inches(6.6), Inches(12.7), Inches(0.72),
             font_size=11, color=RGBColor(0x6B, 0x72, 0x80))

    #  SLIDE 11: SIMULAÇÃO DE RECEITA 
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
    add_text(sl, "SIMULAÇÕES DE CRESCIMENTO DE RECEITA",
             Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
             font_size=22, bold=True, color=BRANCO)

    cores_sim = [RGBColor(0x10, 0xB9, 0x81), RGBColor(0x25, 0x63, 0xEB),
                 RGBColor(0x8E, 0x44, 0xAD)]
    if not df_sim_rec.empty:
        card_w_s = Inches(4.1)
        for si, (_, row) in enumerate(df_sim_rec.iterrows()):
            left_s = Inches(0.2 + si * 4.35)
            add_rect(sl, left_s, Inches(1.2), card_w_s, Inches(4.8), cores_sim[si % 4])
            add_text(sl, str(row["Estratégia"]),
                     left_s + Inches(0.12), Inches(1.35), card_w_s - Inches(0.24), Inches(0.8),
                     font_size=15, bold=True, color=BRANCO)
            add_text(sl, brl(row["Impacto Mensal"]),
                     left_s + Inches(0.12), Inches(2.35), card_w_s - Inches(0.24), Inches(0.65),
                     font_size=22, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
            add_text(sl, "por mês",
                     left_s + Inches(0.12), Inches(2.95), card_w_s - Inches(0.24), Inches(0.3),
                     font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
            add_text(sl, f"Anual: {brl(row['Impacto Anual'])}",
                     left_s + Inches(0.12), Inches(3.35), card_w_s - Inches(0.24), Inches(0.35),
                     font_size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
            add_text(sl, str(row["Complexidade"]),
                     left_s + Inches(0.12), Inches(3.85), card_w_s - Inches(0.24), Inches(0.3),
                     font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
            add_text(sl, str(row["Como"]),
                     left_s + Inches(0.12), Inches(4.3), card_w_s - Inches(0.24), Inches(0.7),
                     font_size=13, color=RGBColor(0xFF, 0xFF, 0xFF))

        melhor = df_sim_rec.nlargest(1, "Impacto Mensal").iloc[0]
        add_rect(sl, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.8), RGBColor(0xFE, 0xF3, 0xC7))
        _melhor_txt = (f"Maior oportunidade: {melhor['Estratégia']}"
                       f"  |  {brl(melhor['Impacto Mensal'])} por mes"
                       f"  |  {brl(melhor['Impacto Anual'])} por ano")
        add_text(sl, _melhor_txt,
                 Inches(0.5), Inches(6.35), Inches(12.5), Inches(0.7),
                 font_size=14, bold=True, color=RGBColor(0x92, 0x40, 0x0E))

    #  SLIDE 12: SIMULAÇÃO DE PREÇOS 
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
    add_text(sl, "SIMULAÇÃO DE AJUSTE DE PREÇOS (+10% / +15% / +20%)",
             Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
             font_size=21, bold=True, color=BRANCO)
    add_text(sl, "Estimativa com queda de 5% no volume — top produtos por receita",
             Inches(0.5), Inches(1.08), Inches(12), Inches(0.35),
             font_size=14, color=RGBColor(0x6B, 0x72, 0x80))

    if not df_sim_preco.empty:
        hdrs_sp  = ["Produto", "Receita Atual", "+10% (−5% vol)", "Δ+10%", "+15%", "Δ+15%", "+20%", "Δ+20%"]
        col_ws_p = [Inches(3.5), Inches(1.5), Inches(1.5), Inches(1.0),
                    Inches(1.2), Inches(1.0), Inches(1.2), Inches(1.0)]
        y_s = Inches(1.5)
        rh  = Inches(0.38)
        x0  = Inches(0.2)
        x   = x0
        for hdr, ww in zip(hdrs_sp, col_ws_p):
            add_rect(sl, x, y_s, ww, rh, AZUL_ESC)
            add_text(sl, hdr, x + Inches(0.04), y_s + Inches(0.06), ww - Inches(0.08), rh - Inches(0.1),
                     font_size=12, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
            x += ww

        for ri, (_, row) in enumerate(df_sim_preco.head(12).iterrows()):
            y_r  = y_s + rh * (ri + 1)
            bg_r = CINZA_CLR if ri % 2 == 0 else BRANCO
            x    = x0
            r10  = row.get("+10% preço (-5% vol)", 0)
            r15  = row.get("+15% preço (-5% vol)", 0)
            r20  = row.get("+20% preço (-5% vol)", 0)
            d10  = row.get("Δ +10%", 0)
            d15  = row.get("Δ +15%", 0)
            d20  = row.get("Δ +20%", 0)
            vals_sp = [str(row["xProd"])[:38], brl(row["receita"]),
                       brl(r10), f"+{brl(d10)}",
                       brl(r15), f"+{brl(d15)}",
                       brl(r20), f"+{brl(d20)}"]
            for vi, (val, ww) in enumerate(zip(vals_sp, col_ws_p)):
                add_rect(sl, x, y_r, ww, rh, bg_r)
                al = PP_ALIGN.LEFT if vi == 0 else PP_ALIGN.CENTER
                col_v = RGBColor(0x10, 0xB9, 0x81) if vi in (3, 5, 7) else TEXTO
                add_text(sl, val, x + Inches(0.04), y_r + Inches(0.06),
                         ww - Inches(0.08), rh - Inches(0.1),
                         font_size=12, color=col_v, align=al)
                x += ww

    #  SLIDE 13: COMBOS PRECIFICADOS 
    if not df_combos.empty:
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
        add_text(sl, "PRECIFICAÇÃO DE COMBOS",
                 Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
                 font_size=22, bold=True, color=BRANCO)
        add_text(sl, "Baseado nos pares mais frequentes — desconto de 5% e 10%",
                 Inches(0.5), Inches(1.05), Inches(12), Inches(0.35),
                 font_size=14, color=RGBColor(0x6B, 0x72, 0x80))

        hdrs_cb  = ["Combo", "Total Individual", "5% desc.", "10% desc.", "Freq."]
        col_ws_c = [Inches(5.5), Inches(2.0), Inches(1.8), Inches(1.8), Inches(1.2)]
        y_c = Inches(1.5)
        rh  = Inches(0.42)
        x   = Inches(0.3)
        for hdr, ww in zip(hdrs_cb, col_ws_c):
            add_rect(sl, x, y_c, ww, rh, AZUL_ESC)
            add_text(sl, hdr, x + Inches(0.05), y_c + Inches(0.07),
                     ww - Inches(0.1), rh - Inches(0.1),
                     font_size=14, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
            x += ww
        for ri, row in enumerate(df_combos.head(10).itertuples()):
            y_r = y_c + rh * (ri + 1)
            bg  = CINZA_CLR if ri % 2 == 0 else BRANCO
            x   = Inches(0.3)
            vals_c = [str(row.Combo)[:65], brl(row._4),
                      brl(row._5), brl(row._6), fmt_num(row.Frequência)]
            for vi, (val, ww) in enumerate(zip(vals_c, col_ws_c)):
                add_rect(sl, x, y_r, ww, rh, bg)
                al = PP_ALIGN.LEFT if vi == 0 else PP_ALIGN.CENTER
                col_v = RGBColor(0x10, 0xB9, 0x81) if vi in (2, 3) else TEXTO
                add_text(sl, val, x + Inches(0.05), y_r + Inches(0.07),
                         ww - Inches(0.1), rh - Inches(0.1),
                         font_size=13, color=col_v, align=al)
                x += ww

        add_rect(sl, Inches(0.3), Inches(6.35), Inches(12.7), Inches(0.75), RGBColor(0xD1, 0xFA, 0xE5))
        add_text(sl, "O combo com 5% de desconto mantém margem saudável e aumenta percepção de valor para o cliente",
                 Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.55),
                 font_size=15, color=RGBColor(0x06, 0x5F, 0x46))

    #  SLIDE 14: METAS 
    if not df_metas.empty:
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
        add_text(sl, "METAS MENSAIS POR PRODUTO",
                 Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
                 font_size=22, bold=True, color=BRANCO)

        hdrs_mt  = ["Produto", "Receita Atual", "Meta +10%", "Meta +20%"]
        col_ws_m = [Inches(5.5), Inches(2.5), Inches(2.5), Inches(2.5)]
        y_m = Inches(1.15)
        rh  = Inches(0.37)
        x   = Inches(0.2)
        for hdr, ww in zip(hdrs_mt, col_ws_m):
            add_rect(sl, x, y_m, ww, rh, AZUL_ESC)
            add_text(sl, hdr, x + Inches(0.05), y_m + Inches(0.05),
                     ww - Inches(0.1), rh - Inches(0.08),
                     font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
            x += ww
        for ri, (_, row) in enumerate(df_metas.head(9).iterrows()):
            y_r = y_m + rh * (ri + 1)
            bg  = CINZA_CLR if ri % 2 == 0 else BRANCO
            x   = Inches(0.2)
            vals_m = [str(row["xProd"])[:55],
                      brl(row["receita"]), brl(row["Meta +10%"]), brl(row["Meta +20%"])]
            for vi, (val, ww) in enumerate(zip(vals_m, col_ws_m)):
                add_rect(sl, x, y_r, ww, rh, bg)
                al = PP_ALIGN.LEFT if vi == 0 else PP_ALIGN.CENTER
                col_v = (RGBColor(0x16, 0xA0, 0x85) if vi == 2
                         else (RGBColor(0x8E, 0x44, 0xAD) if vi == 3 else TEXTO))
                add_text(sl, val, x + Inches(0.05), y_r + Inches(0.06),
                         ww - Inches(0.1), rh - Inches(0.1),
                         font_size=12, color=col_v, align=al)
                x += ww

        fat_top10 = df_metas["receita"].sum()
        add_rect(sl, Inches(0.2), Inches(4.9), Inches(4.0), Inches(1.0), RGBColor(0xEC, 0xF0, 0xF1))
        add_text(sl, "Top 10 Atual",
                 Inches(0.3), Inches(4.97), Inches(3.8), Inches(0.3),
                 font_size=13, color=RGBColor(0x6B, 0x72, 0x80))
        add_text(sl, brl(fat_top10),
                 Inches(0.3), Inches(5.27), Inches(3.8), Inches(0.55),
                 font_size=16, bold=True, color=AZUL_ESC)

        add_rect(sl, Inches(4.5), Inches(4.9), Inches(4.0), Inches(1.0), RGBColor(0xD1, 0xFA, 0xE5))
        add_text(sl, "Meta +10%",
                 Inches(4.6), Inches(4.97), Inches(3.8), Inches(0.3),
                 font_size=13, color=RGBColor(0x6B, 0x72, 0x80))
        add_text(sl, brl(fat_top10 * 1.10),
                 Inches(4.6), Inches(5.27), Inches(3.8), Inches(0.55),
                 font_size=16, bold=True, color=RGBColor(0x10, 0xB9, 0x81))

        add_rect(sl, Inches(8.8), Inches(4.9), Inches(4.0), Inches(1.0), RGBColor(0xDB, 0xEA, 0xFE))
        add_text(sl, "Meta +20%",
                 Inches(8.9), Inches(4.97), Inches(3.8), Inches(0.3),
                 font_size=13, color=RGBColor(0x6B, 0x72, 0x80))
        add_text(sl, brl(fat_top10 * 1.20),
                 Inches(8.9), Inches(5.27), Inches(3.8), Inches(0.55),
                 font_size=16, bold=True, color=AZUL_MED)

    #  SLIDE 15: FLUXO DE VENDAS POR HORÁRIO
    if df_por_hora is not None and not df_por_hora.empty:
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, W, Inches(1.0), AZUL_ESC)
        add_text(sl, "FLUXO DE VENDAS POR HORÁRIO DO DIA",
                 Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
                 font_size=22, bold=True, color=BRANCO)

        CORES_TURNO_H = {"Manhã": "#2563EB", "Tarde": "#E67E22", "Noite": "#1A234E"}

        fig_fh, ax_fh = plt.subplots(figsize=(8.5, 5.0))
        df_fh = df_por_hora.sort_values("hora")
        cores_fh = df_fh["turno"].map(CORES_TURNO_H).fillna("#95A5A6")
        bars_fh = ax_fh.bar(df_fh["hora"], df_fh["transacoes"],
                            color=cores_fh, edgecolor="white", width=0.75)
        for bar, n in zip(bars_fh, df_fh["transacoes"]):
            if n > 0:
                ax_fh.text(bar.get_x() + bar.get_width() / 2,
                           bar.get_height() + df_fh["transacoes"].max() * 0.01,
                           f"{int(n):,}".replace(",", "."),
                           ha="center", fontsize=8, fontweight="bold", color="#1F2937")
        ax_fh.set_xticks(range(0, 24))
        ax_fh.set_xticklabels([f"{h:02d}h" for h in range(0, 24)], fontsize=8, rotation=45)
        ax_fh.set_xlabel("Hora do dia", fontsize=10)
        ax_fh.set_ylabel("Nº de Pedidos", fontsize=10)
        ax_fh.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{int(x):,}".replace(",", ".")))
        ax_fh.spines[["top", "right"]].set_visible(False)
        # Legenda de turnos
        patches_fh = [mpatches.Patch(color=v, label=k) for k, v in CORES_TURNO_H.items()]
        ax_fh.legend(handles=patches_fh, fontsize=9, loc="upper right")
        fig_fh.tight_layout()
        plt_to_pptx_image(fig_fh, sl, Inches(0.3), Inches(1.1), Inches(8.8), Inches(6.2))

        # Painel de turnos à direita
        if df_por_turno is not None and not df_por_turno.empty:
            hora_pico = df_fh.loc[df_fh["transacoes"].idxmax()]
            turno_lider = df_por_turno.loc[df_por_turno["transacoes"].idxmax()]

            add_rect(sl, Inches(9.3), Inches(1.1), Inches(3.7), Inches(6.2),
                     RGBColor(0xF0, 0xF9, 0xFF))
            add_text(sl, "RESUMO POR TURNO",
                     Inches(9.45), Inches(1.25), Inches(3.4), Inches(0.4),
                     font_size=14, bold=True, color=AZUL_ESC)

            for ti2, (_, tr) in enumerate(df_por_turno.iterrows()):
                cor_t = RGBColor(*[int(CORES_TURNO_H.get(tr["turno"], "#2563EB").lstrip("#")[i:i+2], 16) for i in (0, 2, 4)])
                y_t = Inches(1.75 + ti2 * 1.5)
                add_rect(sl, Inches(9.3), y_t, Inches(3.7), Inches(1.35), cor_t)
                pct_t = f"{tr['pct']:.1f}".replace(".", ",") + "%"
                add_text(sl, tr["turno"].upper(),
                         Inches(9.4), y_t + Inches(0.07), Inches(3.5), Inches(0.3),
                         font_size=15, bold=True, color=BRANCO)
                add_text(sl, pct_t,
                         Inches(9.4), y_t + Inches(0.36), Inches(3.5), Inches(0.5),
                         font_size=22, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
                add_text(sl, f"{fmt_num(tr['transacoes'])} pedidos  ·  {brl(tr['receita'])}",
                         Inches(9.4), y_t + Inches(0.88), Inches(3.5), Inches(0.3),
                         font_size=13, color=BRANCO, align=PP_ALIGN.CENTER)

            add_rect(sl, Inches(9.3), Inches(6.25), Inches(3.7), Inches(0.75),
                     RGBColor(0xDB, 0xEA, 0xFE))
            add_text(sl,
                     f"Pico: {int(hora_pico['hora']):02d}h  ·  "
                     f"Turno líder: {turno_lider['turno']}",
                     Inches(9.4), Inches(6.33), Inches(3.5), Inches(0.55),
                     font_size=14, bold=True, color=AZUL_ESC, align=PP_ALIGN.CENTER)

    # Slide "Horários Inexplorados" removido a pedido do usuário

    #  SLIDE NF-e B2B (se carregado) 
    if df_nfe is not None and not df_nfe.empty and "vNF" in df_nfe.columns and "chave" in df_nfe.columns:
        sl = prs.slides.add_slide(blank)
        add_rect(sl, 0, 0, W, Inches(1.0), RGBColor(0x1A, 0x23, 0x4E))
        add_text(sl, "NF-e — VENDAS PARA EMPRESAS (B2B)",
                 Inches(0.5), Inches(0.1), Inches(12), Inches(0.8),
                 font_size=22, bold=True, color=BRANCO)

        notas_b2b = df_nfe.drop_duplicates("chave")
        fat_b2b   = notas_b2b["vNF"].sum()
        n_b2b     = len(notas_b2b)
        tm_b2b    = fat_b2b / n_b2b if n_b2b else 0
        fat_total = kpis["faturamento"]
        fat_nfce  = kpis_nfce["faturamento"]

        cards_b2b = [
            (brl(fat_nfce),  "NFC-e (Consumidor)", AZUL_MED),
            (brl(fat_b2b),   "NF-e (Empresas)",    RGBColor(0x8E, 0x44, 0xAD)),
            (brl(fat_total), "Total Consolidado",   RGBColor(0x10, 0xB9, 0x81)),
            (fmt_num(n_b2b) + " notas", "Notas B2B Emitidas", RGBColor(0xD3, 0x54, 0x00)),
        ]
        for ci, (val, lbl, cor_c) in enumerate(cards_b2b):
            lx = Inches(0.3 + ci * 3.25)
            add_rect(sl, lx, Inches(1.2), Inches(3.0), Inches(1.6), cor_c)
            add_text(sl, val, lx + Inches(0.1), Inches(1.35),
                     Inches(2.8), Inches(0.8),
                     font_size=22, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
            add_text(sl, lbl, lx + Inches(0.1), Inches(2.1),
                     Inches(2.8), Inches(0.5),
                     font_size=14, color=BRANCO, align=PP_ALIGN.CENTER)

        # Top produtos NF-e (nomes limpos — sem Ped./Nro.Item)
        if "xProd" in df_nfe.columns and "vProd" in df_nfe.columns:
            import re as _re_tbl
            def _clean_tbl(s):
                s = _re_tbl.sub(r"\s+Ped\.+\s*:.*$", "", str(s), flags=_re_tbl.IGNORECASE)
                s = _re_tbl.sub(r"\s+Nro\.?\s*Item\s*:.*$", "", s, flags=_re_tbl.IGNORECASE)
                return s.strip()
            _nfe_tbl = df_nfe.copy()
            _nfe_tbl["xProd"] = _nfe_tbl["xProd"].apply(_clean_tbl)
            _dest_b2b = "destinatario" if "destinatario" in _nfe_tbl.columns else None
            _grp_b2b  = ["xProd"] + ([_dest_b2b] if _dest_b2b else [])
            top_b2b = (_nfe_tbl.groupby(_grp_b2b)
                       .agg(receita=("vProd", "sum"), notas=("chave", "nunique"))
                       .sort_values("receita", ascending=False).head(8).reset_index())

            add_text(sl, "TOP PRODUTOS (NF-e)",
                     Inches(0.3), Inches(3.0), Inches(7.5), Inches(0.45),
                     font_size=14, bold=True, color=AZUL_ESC)

            if _dest_b2b:
                hdrs_b2b   = ["#", "Produto", "Empresa", "Notas"]
                col_ws_b2b = [Inches(0.4), Inches(3.3), Inches(3.0), Inches(0.6)]
            else:
                hdrs_b2b   = ["#", "Produto", "Notas"]
                col_ws_b2b = [Inches(0.5), Inches(6.4), Inches(0.8)]

            y_b = Inches(3.5)
            rh  = Inches(0.38)
            x   = Inches(0.3)
            for hdr, ww in zip(hdrs_b2b, col_ws_b2b):
                add_rect(sl, x, y_b, ww, rh, RGBColor(0x1A, 0x23, 0x4E))
                add_text(sl, hdr, x + Inches(0.04), y_b + Inches(0.05),
                         ww - Inches(0.08), rh - Inches(0.08),
                         font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.CENTER)
                x += ww
            for ri, (_, row) in enumerate(top_b2b.iterrows()):
                y_r = y_b + rh * (ri + 1)
                bg  = CINZA_CLR if ri % 2 == 0 else BRANCO
                x   = Inches(0.3)
                if _dest_b2b:
                    vals_b = [str(ri+1), str(row["xProd"])[:40],
                              str(row.get(_dest_b2b, ""))[:38],
                              fmt_num(row["notas"])]
                else:
                    vals_b = [str(ri+1), str(row["xProd"])[:60], fmt_num(row["notas"])]
                for vi, (val, ww) in enumerate(zip(vals_b, col_ws_b2b)):
                    add_rect(sl, x, y_r, ww, rh, bg)
                    al = PP_ALIGN.LEFT if vi in (1, 2) else PP_ALIGN.CENTER
                    add_text(sl, val, x + Inches(0.04), y_r + Inches(0.05),
                             ww - Inches(0.08), rh - Inches(0.08),
                             font_size=11, color=TEXTO, align=al)
                    x += ww

        # Top produtos NF-e à direita (nomes limpos — sem Ped./Nro.Item)
        if df_nfe is not None and not df_nfe.empty:
            import re as _re_nfe
            def _clean_nfe(s):
                s = _re_nfe.sub(r"\s+Ped\.+\s*:.*$", "", str(s), flags=_re_nfe.IGNORECASE)
                s = _re_nfe.sub(r"\s+Nro\.?\s*Item\s*:.*$", "", s, flags=_re_nfe.IGNORECASE)
                return s.strip()
            _nfe_c = df_nfe.copy()
            _nfe_c["xProd"] = _nfe_c["xProd"].apply(_clean_nfe)
            top_nfe = (_nfe_c.groupby("xProd")["vProd"].sum()
                       .nlargest(8).reset_index().sort_values("vProd"))
            add_text(sl, "TOP PRODUTOS (NF-e)",
                     Inches(8.0), Inches(3.0), Inches(5.1), Inches(0.4),
                     font_size=15, bold=True, color=AZUL_ESC)
            fig_cb2, ax_cb2 = plt.subplots(figsize=(4.8, 3.8))
            ax_cb2.barh(top_nfe["xProd"].str[:25], top_nfe["vProd"],
                        color="#2563EB", height=0.6)
            ax_cb2.tick_params(labelsize=7)
            ax_cb2.spines[["top", "right"]].set_visible(False)
            fig_cb2.tight_layout()
            plt_to_pptx_image(fig_cb2, sl, Inches(8.0), Inches(3.45), Inches(5.1), Inches(3.8))

    # ══════════════════════════════════════════════════════════════════
    # SLIDE: SIMPLES NACIONAL — REGRA DOS 80%
    # ══════════════════════════════════════════════════════════════════
    if sn_result is not None and sn_result.get("status") != "SEM_DADOS":
        sl = prs.slides.add_slide(blank)
        _sn_total   = sn_result["total_compras_comercializacao"]
        _sn_pct     = sn_result["pct_faturamento"]
        _sn_status  = sn_result["status"]
        _sn_fat     = kpis["faturamento"]

        # Header
        add_rect(sl, Inches(0), Inches(0), W, Inches(1.05), AZUL_ESC)
        add_text(sl, "SIMPLES NACIONAL — REGRA DOS 80%",
                 Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.55),
                 font_size=26, bold=True, color=BRANCO)
        add_text(sl, f"{cliente}  ·  {periodo}  ·  Compras de Comercialização vs Faturamento",
                 Inches(0.3), Inches(0.68), Inches(12.7), Inches(0.35),
                 font_size=13, color=RGBColor(0xBA, 0xD0, 0xF0))

        # ── KPI Cards (row) ──────────────────────────────────────────
        _cards = [
            ("Faturamento Total",          brl(_sn_fat),   AZUL_MED),
            ("Compras Comercialização",     brl(_sn_total), AZUL_ESC),
            ("% do Faturamento",           f"{_sn_pct:.1f}%".replace(".", ","),
             RGBColor(0x16, 0xa3, 0x4a) if _sn_pct <= 80 else RGBColor(0xDC, 0x26, 0x26)),
            ("Limite Legal",               "80,0%",        RGBColor(0x92, 0x40, 0x0E)),
        ]
        _cw = Inches(13.33 / len(_cards))
        for _ci, (_lbl, _val, _clr) in enumerate(_cards):
            _cx = Inches(_ci * (13.33 / len(_cards)))
            add_rect(sl, _cx, Inches(1.15), _cw, Inches(1.3), CINZA_CLR)
            add_text(sl, _lbl, _cx + Inches(0.1), Inches(1.18), _cw - Inches(0.2), Inches(0.35),
                     font_size=11, color=RGBColor(0x6B, 0x72, 0x80))
            add_text(sl, _val, _cx + Inches(0.1), Inches(1.55), _cw - Inches(0.2), Inches(0.6),
                     font_size=22, bold=True, color=_clr)

        # ── Barra de progresso visual ─────────────────────────────────
        _bar_y = Inches(2.65)
        # Fundo cinza
        add_rect(sl, Inches(0.5), _bar_y, Inches(12.33), Inches(0.55),
                 RGBColor(0xE5, 0xE7, 0xEB))
        # Preenchimento colorido
        _bar_clr = (RGBColor(0x16, 0xa3, 0x4a) if _sn_pct <= 70
                    else (RGBColor(0xF5, 0x9E, 0x0B) if _sn_pct <= 80
                          else RGBColor(0xDC, 0x26, 0x26)))
        _bar_w = Inches(12.33 * min(_sn_pct, 100) / 100)
        add_rect(sl, Inches(0.5), _bar_y, _bar_w, Inches(0.55), _bar_clr)
        add_text(sl, f"{_sn_pct:.1f}%", Inches(0.5), _bar_y + Inches(0.1),
                 _bar_w, Inches(0.35), font_size=13, bold=True, color=BRANCO, align=PP_ALIGN.RIGHT)
        # Linha de limite 80%
        _lim_x = Inches(0.5 + 12.33 * 0.80)
        add_rect(sl, _lim_x - Inches(0.015), _bar_y - Inches(0.08),
                 Inches(0.03), Inches(0.71), AZUL_ESC)
        add_text(sl, "Limite 80%", _lim_x - Inches(0.5), _bar_y + Inches(0.6),
                 Inches(1.0), Inches(0.3), font_size=10, bold=True, color=AZUL_ESC, align=PP_ALIGN.CENTER)

        # ── Status banner ─────────────────────────────────────────────
        _status_msgs = {
            "OK":       ("✓ DENTRO DO LIMITE — compras representam {pct} do faturamento (≤ 80%)", RGBColor(0x16, 0xa3, 0x4a)),
            "ALERTA":   ("⚠ ATENÇÃO — compras em {pct} do faturamento — próximo do limite de 80%", RGBColor(0xD9, 0x77, 0x06)),
            "EXCEDIDO": ("✗ LIMITE EXCEDIDO — compras em {pct} do faturamento — acima de 80%",     RGBColor(0xDC, 0x26, 0x26)),
        }
        _msg_tmpl, _msg_clr = _status_msgs.get(_sn_status, ("—", AZUL_ESC))
        _msg = _msg_tmpl.format(pct=f"{_sn_pct:.1f}%".replace(".", ","))
        add_rect(sl, Inches(0.5), Inches(3.45), Inches(12.33), Inches(0.55), _msg_clr)
        add_text(sl, _msg, Inches(0.6), Inches(3.5), Inches(12.1), Inches(0.45),
                 font_size=14, bold=True, color=BRANCO)

        # ── Tabela: Compras por CFOP ──────────────────────────────────
        df_por_cfop = sn_result.get("df_por_cfop", pd.DataFrame())
        add_text(sl, "COMPRAS POR CFOP",
                 Inches(0.5), Inches(4.15), Inches(5.5), Inches(0.4),
                 font_size=14, bold=True, color=AZUL_ESC)
        if not df_por_cfop.empty:
            _hdr_cfop = ["CFOP", "Total Compras", "Notas", "Itens"]
            _cw_cfop  = [Inches(1.1), Inches(2.4), Inches(1.0), Inches(1.0)]
            _rh = Inches(0.38)
            _y_c = Inches(4.58)
            _x_c = Inches(0.5)
            for _hi, (_hh, _hw) in enumerate(zip(_hdr_cfop, _cw_cfop)):
                add_rect(sl, _x_c, _y_c, _hw, _rh, AZUL_ESC)
                add_text(sl, _hh, _x_c + Inches(0.05), _y_c + Inches(0.07),
                         _hw - Inches(0.1), _rh - Inches(0.1),
                         font_size=11, bold=True, color=BRANCO)
                _x_c += _hw
            for _ri, (_, _row) in enumerate(df_por_cfop.head(6).iterrows()):
                _y_c += _rh
                _x_c = Inches(0.5)
                bg = CINZA_CLR if _ri % 2 == 0 else BRANCO
                _vals_c = [_row["CFOP"], brl(_row["total_compras"]),
                           fmt_num(_row["notas"]), fmt_num(_row["itens"])]
                for _vi, (_vv, _vw) in enumerate(zip(_vals_c, _cw_cfop)):
                    add_rect(sl, _x_c, _y_c, _vw, _rh, bg)
                    _al = PP_ALIGN.LEFT if _vi == 0 else PP_ALIGN.RIGHT
                    add_text(sl, str(_vv), _x_c + Inches(0.05), _y_c + Inches(0.07),
                             _vw - Inches(0.1), _rh - Inches(0.1),
                             font_size=11, color=TEXTO, align=_al)
                    _x_c += _vw

        # ── Tabela: Top Fornecedores ───────────────────────────────────
        df_por_forn = sn_result.get("df_por_fornecedor", pd.DataFrame())
        add_text(sl, "TOP FORNECEDORES",
                 Inches(6.8), Inches(4.15), Inches(6.0), Inches(0.4),
                 font_size=14, bold=True, color=AZUL_ESC)
        if not df_por_forn.empty:
            _hdr_forn = ["Fornecedor", "Total Compras", "Notas"]
            _cw_forn  = [Inches(3.6), Inches(1.8), Inches(0.8)]
            _rh = Inches(0.38)
            _y_f = Inches(4.58)
            _x_f = Inches(6.8)
            for _hi, (_hh, _hw) in enumerate(zip(_hdr_forn, _cw_forn)):
                add_rect(sl, _x_f, _y_f, _hw, _rh, AZUL_ESC)
                add_text(sl, _hh, _x_f + Inches(0.05), _y_f + Inches(0.07),
                         _hw - Inches(0.1), _rh - Inches(0.1),
                         font_size=11, bold=True, color=BRANCO)
                _x_f += _hw
            for _ri, (_, _row) in enumerate(df_por_forn.head(6).iterrows()):
                _y_f += _rh
                _x_f = Inches(6.8)
                bg = CINZA_CLR if _ri % 2 == 0 else BRANCO
                _nome_f = str(_row["emitente"])[:40]
                _vals_f = [_nome_f, brl(_row["total_compras"]), fmt_num(_row["notas"])]
                for _vi, (_vv, _vw) in enumerate(zip(_vals_f, _cw_forn)):
                    add_rect(sl, _x_f, _y_f, _vw, _rh, bg)
                    _al = PP_ALIGN.LEFT if _vi == 0 else PP_ALIGN.RIGHT
                    add_text(sl, str(_vv), _x_f + Inches(0.05), _y_f + Inches(0.07),
                             _vw - Inches(0.1), _rh - Inches(0.1),
                             font_size=11, color=TEXTO, align=_al)
                    _x_f += _vw

        # Nota de rodapé
        add_text(sl,
                 "CFOPs considerados — comercialização: 1102, 2102, 1403, 2403, 1104, 2104, 1117, 2117, 1410, 2410  "
                 "| Excluídos: imobilizado (14xx), uso e consumo (15xx), industrialização (11xx)",
                 Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.35),
                 font_size=9, color=RGBColor(0x9C, 0xA3, 0xAF))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


#
# INTERFACE STREAMLIT
#
def main():
    st.markdown("""
    <style>
    .block-container { padding-top: 1.2rem; }
    .metric-label { font-size: 12px !important; }
    /* Remove bordas arredondadas da logo no painel lateral */
    section[data-testid="stSidebar"] img {
        border-radius: 0 !important;
        -webkit-border-radius: 0 !important;
    }
    </style>
    """, unsafe_allow_html=True)

    #  SIDEBAR
    with st.sidebar:
        from pathlib import Path as _Path
        _logo_sidebar = _Path(__file__).parent / "LOGO S FUNDO 2.png"
        if _logo_sidebar.exists():
            st.image(str(_logo_sidebar), use_container_width=True)
        st.markdown("## Análise de Vendas CP")
        st.divider()

        # Auto-detectar empresa e período dos XMLs já processados
        _auto_empresa = ""
        _auto_periodo = ""
        if "_analise" in st.session_state:
            _R = st.session_state["_analise"]
            _df_tmp = _R.get("df_nfce")
            if _df_tmp is not None and not _df_tmp.empty and "emitente" in _df_tmp.columns:
                _vals = _df_tmp["emitente"].dropna()
                _vals = _vals[_vals != ""]
                if not _vals.empty:
                    _auto_empresa = _vals.mode()[0]
            if _df_tmp is not None and not _df_tmp.empty and "dhEmi" in _df_tmp.columns:
                _MESES_PT_s = {1:"Janeiro",2:"Fevereiro",3:"Março",4:"Abril",5:"Maio",6:"Junho",
                               7:"Julho",8:"Agosto",9:"Setembro",10:"Outubro",11:"Novembro",12:"Dezembro"}
                _meses = sorted(_df_tmp["dhEmi"].dropna().dt.to_period("M").unique())
                if len(_meses) == 1:
                    _auto_periodo = f"{_MESES_PT_s[_meses[0].month]} {_meses[0].year}"
                elif len(_meses) >= 2:
                    _m1, _m2 = _meses[0], _meses[-1]
                    if _m1.year == _m2.year:
                        _auto_periodo = f"{_MESES_PT_s[_m1.month]} - {_MESES_PT_s[_m2.month]} {_m1.year}"
                    else:
                        _auto_periodo = f"{_MESES_PT_s[_m1.month]} {_m1.year} - {_MESES_PT_s[_m2.month]} {_m2.year}"

        cliente = st.text_input("Nome do cliente", value=_auto_empresa, placeholder="Detectado automaticamente")
        periodo = st.text_input("Período", value=_auto_periodo, placeholder="Detectado automaticamente")
        top_n   = st.slider("Itens no Market Basket", 5, 20, 10)
        st.divider()

        # ── Área única de upload ─────────────────────────────
        st.markdown("**📂 Arquivos de Notas Fiscais**")
        st.caption("ZIP, RAR, 7z, XML ou Excel — qualquer formato, qualquer quantidade")
        arquivos_upload = st.file_uploader(
            "Arraste ou selecione os arquivos",
            type=["zip", "rar", "7z", "xml", "xlsx", "xls"],
            accept_multiple_files=True,
            label_visibility="collapsed",
            key=f"uploader_{st.session_state.get('_upload_key', 0)}",
        )

        # Resumo do que foi carregado
        if arquivos_upload:
            _n_comp = sum(1 for f in arquivos_upload if f.name.lower().endswith((".zip",".rar",".7z")))
            _n_xmls = sum(1 for f in arquivos_upload if f.name.lower().endswith(".xml"))
            _n_xls  = sum(1 for f in arquivos_upload if f.name.lower().endswith((".xlsx",".xls")))
            _partes = []
            if _n_comp: _partes.append(f"{_n_comp} compactado(s)")
            if _n_xmls: _partes.append(f"{_n_xmls} XML(s)")
            if _n_xls:  _partes.append(f"{_n_xls} Excel(s)")
            st.success("✅ " + " · ".join(_partes))

        # ── Pasta local (só desktop) ──────────────────────────
        if not _is_cloud():
            st.markdown("**📁 Ou selecione uma pasta**")
            if "pastas_xml" not in st.session_state:
                st.session_state["pastas_xml"] = []
            if st.button("+ Adicionar pasta", use_container_width=True, key="btn_add_pasta"):
                _caminho = _abrir_seletor_pasta()
                if _caminho and _caminho not in st.session_state["pastas_xml"]:
                    st.session_state["pastas_xml"].append(_caminho)
            _pastas_validas = []
            for _idx, _p in enumerate(list(st.session_state["pastas_xml"])):
                _n    = _contar_xmls_pasta(_p)
                _nome = _Path(_p).name
                _cor     = "#D1FAE5" if _n > 0 else "#FEE2E2"
                _txt_cor = "#065F46" if _n > 0 else "#991B1B"
                _col_info, _col_rm = st.columns([5, 1])
                with _col_info:
                    st.markdown(
                        f"<div style='background:{_cor};border-radius:6px;padding:5px 9px;"
                        f"font-size:11px;color:{_txt_cor};margin:2px 0'>"
                        f"<b>{_nome}</b> · {_n} XMLs</div>",
                        unsafe_allow_html=True)
                with _col_rm:
                    if st.button("✕", key=f"rm_pasta_{_idx}"):
                        st.session_state["pastas_xml"].pop(_idx)
                        st.rerun()
                if _n > 0:
                    _pastas_validas.append(_p)
        else:
            if "pastas_xml" not in st.session_state:
                st.session_state["pastas_xml"] = []
            _pastas_validas = []

        # ── Determina fonte ativa ─────────────────────────────
        f_nfce, f_nfe = None, None  # legado Excel não mais exposto na UI
        _tem_dados = bool(arquivos_upload) or bool(_pastas_validas)

        st.divider()
        # ── Simples Nacional ──────────────────────────────────
        is_simples = st.checkbox(
            "🏪 Empresa é Simples Nacional?",
            value=st.session_state.get("_chk_simples", False),
            key="chk_simples",
            help="Ativa a verificação da regra dos 80% — compras de comercialização não podem ultrapassar 80% do faturamento.",
        )
        st.session_state["_chk_simples"] = is_simples
        if is_simples:
            st.caption("📥 XMLs de **entrada** (notas de compra) — upload ou pasta")
            arquivos_entrada = st.file_uploader(
                "XMLs de Entrada",
                type=["xml", "zip"],
                accept_multiple_files=True,
                label_visibility="collapsed",
                key=f"uploader_entrada_{st.session_state.get('_upload_key', 0)}",
            )
            if arquivos_entrada:
                st.success(f"✅ {len(arquivos_entrada)} arquivo(s) de entrada carregado(s)")

            # ── Pasta local de entradas (só desktop) ─────────────────
            if not _is_cloud():
                if st.button("📁 Selecionar pasta de entradas", use_container_width=True, key="btn_pasta_entrada"):
                    _pe = _abrir_seletor_pasta()
                    if _pe:
                        st.session_state["_pasta_entrada"] = _pe
                _pasta_entrada = st.session_state.get("_pasta_entrada", "")
                if _pasta_entrada:
                    _n_ent = _contar_xmls_pasta(_pasta_entrada)
                    _cor_e = "#D1FAE5" if _n_ent > 0 else "#FEE2E2"
                    _tc_e  = "#065F46" if _n_ent > 0 else "#991B1B"
                    _col_ei, _col_er = st.columns([5, 1])
                    with _col_ei:
                        st.markdown(
                            f"<div style='background:{_cor_e};border-radius:6px;padding:5px 9px;"
                            f"font-size:11px;color:{_tc_e};margin:2px 0'>"
                            f"<b>{_Path(_pasta_entrada).name}</b> · {_n_ent} XMLs de entrada</div>",
                            unsafe_allow_html=True)
                    with _col_er:
                        if st.button("✕", key="rm_pasta_entrada"):
                            del st.session_state["_pasta_entrada"]
                            st.rerun()
            else:
                _pasta_entrada = ""
        else:
            arquivos_entrada = []
            _pasta_entrada = ""

        st.divider()
        btn_analisar = st.button(
            "▶ Analisar" if _tem_dados else "Carregue os arquivos acima",
            use_container_width=True,
            disabled=not _tem_dados,
            type="primary",
        )

        # ── Botão Nova Análise ────────────────────────────────
        if "_analise" in st.session_state or _tem_dados:
            if st.button("🔄 Nova Análise", use_container_width=True, key="btn_nova_analise"):
                for _k in ["_analise", "_analise_fp", "pastas_xml"]:
                    if _k in st.session_state:
                        del st.session_state[_k]
                # Incrementa chave do uploader para forçar limpeza dos arquivos
                st.session_state["_upload_key"] = st.session_state.get("_upload_key", 0) + 1
                st.rerun()

    # ── Fingerprint da fonte de dados ──
    _fp_entrada = tuple(sorted((f.name, f.size) for f in arquivos_entrada)) if arquivos_entrada else ()
    _fp_pe = _pasta_entrada if _pasta_entrada else ""
    if arquivos_upload:
        _fp = ("uploads", tuple(sorted((f.name, f.size) for f in arquivos_upload)), tuple(sorted(_pastas_validas)), top_n, is_simples, _fp_entrada, _fp_pe)
    elif _pastas_validas:
        _fp = ("pastas", tuple(sorted(_pastas_validas)), top_n, is_simples, _fp_entrada, _fp_pe)
    elif f_nfce:
        _fp = ("excel", f_nfce.name, getattr(f_nfce, "size", 0), top_n, is_simples, _fp_entrada, _fp_pe)
    else:
        _fp = None

    _tem_cache = (
        _fp is not None and
        "_analise" in st.session_state and
        st.session_state.get("_analise_fp") == _fp
    )

    #  SEM DADOS / AGUARDANDO BOTÃO
    if not _tem_dados or (not _tem_cache and not btn_analisar):
        st.markdown("""
        <div style="background:linear-gradient(135deg,#1e3a5f,#2563eb);
                    padding:36px 40px;border-radius:16px;color:white;margin-bottom:20px">
          <h1 style="margin:0;font-size:32px;font-weight:700;letter-spacing:-0.5px">
            Análise de Vendas CP
          </h1>
          <p style="margin:10px 0 0;opacity:.8;font-size:15px">
            Carregue os arquivos na barra lateral e clique em <b>▶ Analisar</b>.
          </p>
        </div>
        """, unsafe_allow_html=True)

        if _tem_dados and not btn_analisar:
            if arquivos_upload:
                _nomes = ", ".join(f.name for f in arquivos_upload[:3])
                if len(arquivos_upload) > 3:
                    _nomes += f" +{len(arquivos_upload)-3} mais"
                st.info(f"📎 **{_nomes}** carregado(s). Clique em **▶ Analisar** para iniciar.")
            elif _pastas_validas:
                st.info(f"📁 **{len(_pastas_validas)} pasta(s)** selecionada(s). Clique em **▶ Analisar** para iniciar.")
        else:
            col1, col2, col3 = st.columns(3)
            col1.info("**Formatos aceitos**\n\nZIP · RAR · 7z · XML · Excel")
            col2.info("**Market Basket**\n\nPares e combos de 3 produtos mais frequentes")
            col3.info("**Exportação**\n\nPowerPoint completo + Excel com todas as abas")
        return

    if _tem_cache:
        # ── Restaura resultados do cache de sessão ──
        _R           = st.session_state["_analise"]
        df_nfce      = _R["df_nfce"]
        df_nfe       = _R["df_nfe"]
        df           = df_nfce
        df_all       = _R["df_all"]
        cli_label    = _R["cli_label"]
        per_label    = _R["per_label"]
        cnpj_label   = _R.get("cnpj_label", "")
        tem_nfe      = _R["tem_nfe"]
        fonte_label  = _R["fonte_label"]
        kpis         = _R["kpis"]
        kpis_nfce    = _R["kpis_nfce"]
        df_pares     = _R["df_pares"]
        df_trios     = _R["df_trios"]
        df_cesta     = _R["df_cesta"]
        df_bcg       = _R["df_bcg"]
        df_abc       = _R.get("df_abc", pd.DataFrame())
        df_remocao   = _R["df_remocao"]
        df_turno     = _R["df_turno"]
        df_solo      = _R["df_solo"]
        df_anti      = _R["df_anti"]
        df_dia_tipo  = _R["df_dia_tipo"]
        df_dia_semana= _R["df_dia_semana"]
        df_elev      = _R["df_elev"]
        df_redu      = _R["df_redu"]
        df_sim_preco = _R["df_sim_preco"]
        df_sim_rec   = _R["df_sim_rec"]
        df_combos    = _R["df_combos"]
        df_metas     = _R["df_metas"]
        df_horas     = _R["df_horas"]
        df_por_hora  = _R["df_por_hora"]
        df_por_turno = _R["df_por_turno"]
        df_entradas  = _R.get("df_entradas", pd.DataFrame())
        sn_result    = _R.get("sn_result", None)

    else:
        # ── CARREGAR + CALCULAR com tela de progresso ──
        import time as _time
        _n_xml  = 0
        _n_skip = 0
        _t0     = _time.time()

        def _fmt_elapsed(t0):
            s = int(_time.time() - t0)
            return f"{s//60}m {s%60:02d}s" if s >= 60 else f"{s}s"

        def _render_prog(pct: int, etapa: str, t0: float, _box_txt, _box_bar):
            s       = int(_time.time() - t0)
            elapsed = f"{s//60}m {s%60:02d}s" if s >= 60 else (f"{s}s" if s > 0 else "iniciando...")
            _box_txt.markdown(
                f"""
<div style="padding:18px 24px;background:#f0f4ff;border-radius:12px;
            border-left:5px solid #4C6EF5;margin-bottom:8px;">
  <div style="font-size:13px;color:#6b7280;margin-bottom:4px;">
    ⏱ Tempo decorrido: <strong>{elapsed}</strong>
    <span style="margin-left:12px;font-size:11px;color:#9ca3af;">
      (o relógio avança a cada etapa concluída)
    </span>
  </div>
  <div style="font-size:17px;font-weight:600;color:#1e3a5f;">
    {etapa}
  </div>
  <div style="font-size:13px;color:#6b7280;margin-top:4px;">
    {pct}% concluído
  </div>
</div>
""",
                unsafe_allow_html=True,
            )
            _box_bar.progress(pct)

        st.markdown("<div style='height:12px'></div>", unsafe_allow_html=True)
        st.markdown("### ⚙️ Gerando análises...")
        _box_txt = st.empty()
        _box_bar = st.empty()

        # ── Carregar dados ──
        _render_prog(2, "📂 Vasculhando arquivos e lendo XMLs... (pode demorar na primeira vez)", _t0, _box_txt, _box_bar)

        # Separa Excels dos demais (XML/ZIP/RAR/7z)
        _arqs_xml  = [f for f in (arquivos_upload or []) if not f.name.lower().endswith((".xlsx",".xls"))]
        _arqs_xls  = [f for f in (arquivos_upload or []) if f.name.lower().endswith((".xlsx",".xls"))]

        df_nfce, df_nfe, _n_xml, _n_skip = pd.DataFrame(), pd.DataFrame(), 0, 0

        # Processa XMLs/ZIPs/RARs/7z + pastas
        if _arqs_xml or _pastas_validas:
            _arqs_tuple = tuple((f.name, f.read()) for f in _arqs_xml)
            df_nfce, df_nfe, _n_xml, _n_skip = processar_fontes_universal(
                _arqs_tuple, tuple(_pastas_validas)
            )

        # Processa Excels e combina
        for _xf in _arqs_xls:
            _xb = _xf.read()
            try:
                _df_x = carregar_nfce(_xb)
                if not _df_x.empty:
                    df_nfce = pd.concat([df_nfce, _df_x], ignore_index=True) if not df_nfce.empty else _df_x
            except Exception:
                try:
                    _df_x = carregar_nfe(_xb)
                    if not _df_x.empty:
                        df_nfe = pd.concat([df_nfe, _df_x], ignore_index=True) if not df_nfe.empty else _df_x
                except Exception:
                    _n_skip += 1

        if df_nfce.empty and df_nfe.empty:
            st.error("Nenhum dado encontrado nos arquivos enviados. Verifique o formato.")
            return

        if df_nfce.empty:
            st.error("Nenhum dado de NFC-e encontrado. Verifique o arquivo ou pasta.")
            return

        # ── Validação: apenas uma empresa por análise ─────────
        if "emitente" in df_nfce.columns:
            _em_check = df_nfce["emitente"].dropna()
            _em_check = _em_check[_em_check != ""]
            _em_unicas = sorted(_em_check.unique())
            if len(_em_unicas) > 1:
                st.error(
                    f"⚠️ **{len(_em_unicas)} empresas diferentes** foram detectadas nos arquivos enviados. "
                    "Cada análise deve conter notas fiscais de **uma única empresa**. "
                    "Remova os arquivos da empresa incorreta e tente novamente.\n\n"
                    "**Empresas encontradas:**\n" +
                    "\n".join(f"- {_e}" for _e in _em_unicas)
                )
                return

        df     = df_nfce.copy()
        df_all = pd.concat([df_nfce, df_nfe], ignore_index=True) if not df_nfe.empty else df_nfce.copy()

        _MESES_PT = {1:"Janeiro",2:"Fevereiro",3:"Março",4:"Abril",5:"Maio",6:"Junho",
                     7:"Julho",8:"Agosto",9:"Setembro",10:"Outubro",11:"Novembro",12:"Dezembro"}
        # Auto-detectar nome e CNPJ da empresa a partir das notas
        _auto_cli = ""
        cnpj_label = ""
        if "emitente" in df.columns:
            _em_vals = df["emitente"].dropna()
            _em_vals = _em_vals[_em_vals != ""]
            if not _em_vals.empty:
                _auto_cli = _em_vals.mode()[0]
        if "cnpj_emit" in df.columns:
            _cn_vals = df["cnpj_emit"].dropna()
            _cn_vals = _cn_vals[_cn_vals != ""]
            if not _cn_vals.empty:
                _raw_cnpj = _cn_vals.mode()[0]
                # Formata: 00.000.000/0000-00
                _d = "".join(c for c in str(_raw_cnpj) if c.isdigit())
                if len(_d) == 14:
                    cnpj_label = f"{_d[:2]}.{_d[2:5]}.{_d[5:8]}/{_d[8:12]}-{_d[12:]}"
                else:
                    cnpj_label = _raw_cnpj
        cli_label = cliente or _auto_cli or "Cliente"

        if periodo:
            per_label = periodo
        elif "dhEmi" in df.columns and df["dhEmi"].notna().any():
            _meses = sorted(df["dhEmi"].dropna().dt.to_period("M").unique())
            if len(_meses) == 1:
                per_label = f"{_MESES_PT[_meses[0].month]} {_meses[0].year}"
            elif len(_meses) >= 2:
                _m1, _m2 = _meses[0], _meses[-1]
                if _m1.year == _m2.year:
                    per_label = f"{_MESES_PT[_m1.month]} - {_MESES_PT[_m2.month]} {_m1.year}"
                else:
                    per_label = f"{_MESES_PT[_m1.month]} {_m1.year} - {_MESES_PT[_m2.month]} {_m2.year}"
        else:
            per_label = "Período"
        tem_nfe     = not df_nfe.empty
        fonte_label = "NFC-e + NF-e" if tem_nfe else "NFC-e"

        # ── Análises ──
        _render_prog(12, "📊 Calculando faturamento e KPIs...", _t0, _box_txt, _box_bar)
        kpis      = calc_kpis(df_all)
        kpis_nfce = calc_kpis(df)

        _render_prog(20, "🛒 Market Basket — pares de produtos... (pode demorar)", _t0, _box_txt, _box_bar)
        df_pares = calc_basket_pares(df, top_n)

        _render_prog(42, "🛍️ Market Basket — combos de 3 produtos... (pode demorar)", _t0, _box_txt, _box_bar)
        df_trios = calc_basket_trios(df, top_n)

        _render_prog(55, "🧺 Distribuição da cesta de compras...", _t0, _box_txt, _box_bar)
        df_cesta = calc_cesta(df)

        _render_prog(60, "🏷️ Classificando produtos (BCG, remoção, solo)...", _t0, _box_txt, _box_bar)
        df_bcg     = calc_bcg(df_all)
        df_abc     = calc_curva_abc(df_all)
        df_remocao = calc_remocao(df_all)
        df_solo    = calc_solo_produtos(df)
        df_anti    = calc_anti_pares(df)

        _render_prog(67, "🌅 Análise por turno e dia da semana...", _t0, _box_txt, _box_bar)
        df_turno = calc_turno(df)
        df_dia_tipo, df_dia_semana = calc_por_dia_semana(df)

        _render_prog(80, "💡 Identificando ticket drivers...", _t0, _box_txt, _box_bar)
        df_elev, df_redu = calc_ticket_drivers(df)

        _render_prog(87, "🧮 Gerando simulações de preço e receita...", _t0, _box_txt, _box_bar)
        df_sim_preco = calc_simulacao_precos(df_all)
        df_sim_rec   = calc_simulacao_receita(kpis, df_all)
        df_combos    = calc_combo_pricing(df_pares, df)
        df_metas     = calc_metas(df_all)

        _render_prog(93, "⏰ Analisando fluxo por horário...", _t0, _box_txt, _box_bar)
        df_horas = calc_horas_oportunidade(df)
        df_por_hora, df_por_turno = calc_vendas_horario(df)

        # ── Simples Nacional ──────────────────────────────────
        df_entradas = pd.DataFrame()
        sn_result   = None
        _tem_entradas = is_simples and (bool(arquivos_entrada) or bool(_pasta_entrada))
        if _tem_entradas:
            _render_prog(97, "📋 Verificando regra dos 80% — Simples Nacional...", _t0, _box_txt, _box_bar)
            # Re-seek dos uploads (streams já lidos no fingerprint)
            for _ae in arquivos_entrada:
                try:
                    _ae.seek(0)
                except Exception:
                    pass
            df_entradas = parse_entradas_xml(arquivos_entrada)

            # Pasta local de entradas
            if _pasta_entrada:
                from pathlib import Path as _PEPath
                import xml.etree.ElementTree as _ET_PE
                _pe_files = sorted(
                    set(_PEPath(_pasta_entrada).rglob("*.xml")) |
                    set(_PEPath(_pasta_entrada).rglob("*.XML"))
                )
                # Cria objetos file-like a partir dos bytes dos XMLs da pasta
                import io as _io_pe

                class _BytesWrapper:
                    def __init__(self, data, name):
                        self._buf = _io_pe.BytesIO(data)
                        self.name = name
                    def read(self): return self._buf.read()
                    def seek(self, p): self._buf.seek(p)

                _pe_wrappers = []
                for _pf in _pe_files:
                    try:
                        _pe_wrappers.append(_BytesWrapper(_pf.read_bytes(), _pf.name))
                    except Exception:
                        pass

                if _pe_wrappers:
                    _df_pe = parse_entradas_xml(_pe_wrappers)
                    if not _df_pe.empty:
                        df_entradas = pd.concat([df_entradas, _df_pe], ignore_index=True) if not df_entradas.empty else _df_pe

            sn_result = calc_simples_nacional(df_entradas, kpis["faturamento"])

        _render_prog(100, "✅ Análise concluída!", _t0, _box_txt, _box_bar)

        _elapsed_total = _fmt_elapsed(_t0)
        _box_txt.markdown(
            f"""
<div style="padding:18px 24px;background:#ecfdf5;border-radius:12px;
            border-left:5px solid #10b981;margin-bottom:8px;">
  <div style="font-size:17px;font-weight:700;color:#065f46;">
    ✅ Análise concluída em {_elapsed_total}!
  </div>
  <div style="font-size:13px;color:#047857;margin-top:4px;">
    Role a página para ver os resultados abaixo.
  </div>
</div>
""",
            unsafe_allow_html=True,
        )
        _box_bar.progress(100)

        # Feedback XML — card visual
        if _n_xml > 0:
            n_nfce_notas = df_nfce["chave"].nunique() if not df_nfce.empty else 0
            n_nfe_notas  = df_nfe["chave"].nunique()  if not df_nfe.empty  else 0
            fat_nfce     = df_nfce.drop_duplicates("chave")["vNF"].sum() if not df_nfce.empty else 0
            fat_nfe      = df_nfe.drop_duplicates("chave")["vNF"].sum()  if not df_nfe.empty  else 0
            n_notas_total = n_nfce_notas + n_nfe_notas

            # Build extra columns conditionally
            _col_nfe = (
                f"<td style='padding:0 20px;border-left:2px solid #86efac;'>"
                f"<div style='font-size:11px;color:#6b7280;font-weight:600;letter-spacing:.5px'>NF-e (B2B)</div>"
                f"<div style='font-size:18px;font-weight:800;color:#4f46e5'>{brl(fat_nfe)}</div>"
                f"<div style='font-size:11px;color:#9ca3af'>{fmt_num(n_nfe_notas)} notas</div>"
                f"</td>"
            ) if n_nfe_notas > 0 else ""

            _col_skip = (
                f"<td style='padding:0 20px;border-left:2px solid #86efac;'>"
                f"<div style='font-size:11px;color:#6b7280;font-weight:600;letter-spacing:.5px'>DESCARTADAS</div>"
                f"<div style='font-size:18px;font-weight:800;color:#d97706'>{fmt_num(_n_skip)}</div>"
                f"<div style='font-size:11px;color:#9ca3af'>canceladas / denegadas</div>"
                f"</td>"
            ) if _n_skip > 0 else ""

            _html_card = (
"<div style='background:#f0fdf4;border:1px solid #86efac;border-radius:10px;padding:14px 20px;margin-bottom:4px'>"
"<table style='border-collapse:collapse;width:auto'><tr>"
f"<td style='padding:0 20px 0 0'>"
f"<div style='font-size:11px;color:#6b7280;font-weight:600;letter-spacing:.5px'>NOTAS NA ANÁLISE</div>"
f"<div style='font-size:22px;font-weight:800;color:#15803d'>{fmt_num(n_notas_total)}</div>"
f"<div style='font-size:11px;color:#9ca3af'>{fmt_num(_n_xml)} XMLs lidos</div></td>"
f"<td style='padding:0 20px;border-left:2px solid #86efac'>"
f"<div style='font-size:11px;color:#6b7280;font-weight:600;letter-spacing:.5px'>NFC-e (CONSUMIDOR)</div>"
f"<div style='font-size:22px;font-weight:800;color:#15803d'>{brl(fat_nfce)}</div>"
f"<div style='font-size:11px;color:#9ca3af'>{fmt_num(n_nfce_notas)} notas autorizadas</div></td>"
f"{_col_nfe}{_col_skip}"
"</tr></table></div>"
            )
            st.markdown(_html_card, unsafe_allow_html=True)

        # ── Salva tudo no cache de sessão ──
        st.session_state["_analise_fp"] = _fp
        st.session_state["_analise"] = {
            "df_nfce": df_nfce,      "df_nfe": df_nfe,        "df_all": df_all,
            "cli_label": cli_label,  "per_label": per_label,  "cnpj_label": cnpj_label,
            "tem_nfe": tem_nfe,      "fonte_label": fonte_label,
            "kpis": kpis,            "kpis_nfce": kpis_nfce,
            "df_pares": df_pares,    "df_trios": df_trios,
            "df_cesta": df_cesta,    "df_bcg": df_bcg,   "df_abc": df_abc,
            "df_remocao": df_remocao,"df_turno": df_turno,    "df_solo": df_solo,
            "df_anti": df_anti,      "df_dia_tipo": df_dia_tipo, "df_dia_semana": df_dia_semana,
            "df_elev": df_elev,      "df_redu": df_redu,
            "df_sim_preco": df_sim_preco, "df_sim_rec": df_sim_rec,
            "df_combos": df_combos,  "df_metas": df_metas,
            "df_horas": df_horas,    "df_por_hora": df_por_hora, "df_por_turno": df_por_turno,
            "df_entradas": df_entradas, "sn_result": sn_result,
        }
        # Força re-render para o painel lateral enxergar o cache
        # (sidebar é renderizado antes da análise rodar)
        st.rerun()

    _cnpj_line = (f'<p style="margin:2px 0 0;opacity:.6;font-size:12px;letter-spacing:.4px">CNPJ {cnpj_label}</p>'
                  if cnpj_label else "")
    st.markdown(
        '<div style="background:linear-gradient(135deg,#1e3a5f,#2563eb);'
        'padding:20px 28px;border-radius:10px;color:white;margin-bottom:16px">'
        f'<h2 style="margin:0">{cli_label} — Análise Estratégica de Vendas</h2>'
        + _cnpj_line +
        f'<p style="margin:6px 0 0;opacity:.75;font-size:14px">{per_label} &nbsp;·&nbsp; {fonte_label}</p>'
        '</div>',
        unsafe_allow_html=True)

    #  KPIs PRINCIPAIS
    if tem_nfe and "vNF" in df_nfe.columns and "chave" in df_nfe.columns:
        # Com NF-e: mostra consolidado + breakdown NFC-e / NF-e
        notas_b2b = df_nfe.drop_duplicates("chave")
        fat_b2b   = notas_b2b["vNF"].sum()
        n_b2b     = len(notas_b2b)
        fat_nfce  = kpis_nfce["faturamento"]
        pct_nfe   = fat_b2b  / kpis["faturamento"] * 100 if kpis["faturamento"] else 0
        pct_nfce  = fat_nfce / kpis["faturamento"] * 100 if kpis["faturamento"] else 0

        c1, c2, c3, c4, c5 = st.columns(5)
        c1.metric("Faturamento Consolidado (NFC-e + NF-e)", brl(kpis["faturamento"]))
        c2.metric("Pedidos",       fmt_num(kpis["n_pedidos"]))
        c3.metric("Ticket Médio",  brl(kpis["ticket_medio"]))
        c4.metric("Itens por Pedido (Média)",           f"{kpis['ipc']:.2f} itens")
        c5.metric("Total de Itens",fmt_num(kpis["total_itens"]))

        st.markdown(
            f"""<div style="background:#F0F9FF;border:1px solid #BAE6FD;border-radius:8px;
            padding:10px 16px;margin-top:6px;display:flex;gap:32px;align-items:center;flex-wrap:wrap">
            <span style="font-size:13px;color:#0369A1;font-weight:600">📊 Composição do Faturamento</span>
            <span style="font-size:13px;color:#374151">
              <b>NFC-e (Consumidor Final):</b> {brl(fat_nfce)} &nbsp;·&nbsp; {pct_nfce:.1f}%
            </span>
            <span style="font-size:13px;color:#374151">
              <b>NF-e (B2B / Empresas):</b> {brl(fat_b2b)} &nbsp;·&nbsp; {pct_nfe:.1f}%
            </span>
            <span style="font-size:12px;color:#6B7280">
              ⚠️ O sistema fiscal pode exibir apenas NFC-e — compare {brl(fat_nfce)} com o relatório de NFC-e.
            </span>
            </div>""",
            unsafe_allow_html=True,
        )
    else:
        # Só NFC-e: simples e direto
        c1, c2, c3, c4, c5 = st.columns(5)
        c1.metric("Faturamento NFC-e", brl(kpis["faturamento"]))
        c2.metric("Pedidos",           fmt_num(kpis["n_pedidos"]))
        c3.metric("Ticket Médio",      brl(kpis["ticket_medio"]))
        c4.metric("Itens por Pedido (Média)",               f"{kpis['ipc']:.2f} itens")
        c5.metric("Total de Itens",    fmt_num(kpis["total_itens"]))

    st.divider()

    #  ABAS 
    abas = [
        "Market Basket",
        "Cesta",
        "Candidatos a Remoção",
        "Curva ABC",
        "Ticket Drivers",
        "Simulações",
        "Metas",
    ]
    if df_turno is not None or df_dia_tipo is not None or df_por_hora is not None:
        abas.insert(6, "Temporal")
    if not df_nfe.empty:
        abas.append("NF-e (B2B)")
    if sn_result is not None:
        abas.append("Simples Nacional")

    tabs = st.tabs(abas)
    tab_idx = {name: i for i, name in enumerate(abas)}

    #  MARKET BASKET
    with tabs[tab_idx["Market Basket"]]:
        col_p, col_t3 = st.columns(2)

        with col_p:
            st.subheader("Top Pares de Produtos")
            st.dataframe(df_pares, use_container_width=True, hide_index=True, height=380)
            if not df_pares.empty:
                top1 = df_pares.iloc[0]
                st.success(f"Par mais frequente: **{top1['Produto A']}** + **{top1['Produto B']}** "
                           f"— {fmt_num(top1['Frequência'])} pedidos juntos")

        with col_t3:
            st.subheader("Combos de 3 Produtos")
            st.dataframe(df_trios, use_container_width=True, hide_index=True, height=380)
            if not df_trios.empty:
                top3 = df_trios.iloc[0]
                st.success(f"Combo mais frequente: **{top3['Produto A']}** + **{top3['Produto B']}** "
                           f"+ **{top3['Produto C']}** — {fmt_num(top3['Frequência'])} vezes")

    #  TURNO 
    if "Turno" in tab_idx:
        with tabs[tab_idx["Turno"]]:
            st.subheader("Produtos Mais Vendidos por Turno")
            turnos = ["Manhã", "Tarde", "Noite"]
            cols_t = st.columns(3)
            for col_t, turno in zip(cols_t, turnos):
                sub = df_turno[df_turno["turno"] == turno].head(8)
                icon = "" if turno == "Manhã" else ("" if turno == "Tarde" else "")
                col_t.markdown(f"**{icon} {turno}**")
                if sub.empty:
                    col_t.info("Sem dados")
                else:
                    tbl_t = sub[["xProd", "frequencia", "receita"]].copy()
                    tbl_t["receita"] = tbl_t["receita"].apply(brl)
                    tbl_t.columns = ["Produto", "Freq.", "Receita"]
                    col_t.dataframe(tbl_t, use_container_width=True,
                                    hide_index=True, height=320)

    #  CESTA 
    with tabs[tab_idx["Cesta"]]:
        subtab_cesta = st.tabs(["Distribuição", "Produtos Solo"])

        with subtab_cesta[0]:
            st.subheader("Distribuição do Tamanho da Cesta")
            col_g, col_i = st.columns([3, 2])
            with col_g:
                st.plotly_chart(fig_cesta(df_cesta), use_container_width=True)
            with col_i:
                pct_1_2 = df_cesta[df_cesta["Itens/Pedido"] <= 2]["Nº Pedidos"].sum() / df_cesta["Nº Pedidos"].sum() * 100
                st.metric("Pedidos com 1-2 itens", f"{pct_1_2:.1f}".replace(".", ",") + "%")
                st.metric("Cada +1 item equivale a", brl(kpis["ticket_medio"]))
                st.dataframe(df_cesta.head(15), use_container_width=True,
                             hide_index=True, height=300)
            pct_1_2_br = f"{pct_1_2:.1f}".replace(".", ",") + "%"
            st.warning(f" {pct_1_2_br} dos pedidos têm apenas 1 ou 2 itens — "
                       "converter para 3+ itens é o maior alavancador de faturamento.")

        with subtab_cesta[1]:
            st.subheader("Produtos Âncora — Comprados como Único Item")
            st.caption("Produtos que saem sozinhos com maior frequência — "
                       "clientes que vêm só por isso. Oportunidade de upsell direto.")
            if df_solo.empty:
                st.info("Sem dados.")
            else:
                tbl_s = df_solo.copy()
                tbl_s["receita"] = tbl_s["receita"].apply(brl)
                tbl_s.columns = ["Produto", "Pedidos Solo", "Receita Solo"]
                st.dataframe(tbl_s, use_container_width=True, hide_index=True)
                top_solo = df_solo.iloc[0]
                st.info(f" **{top_solo['xProd']}** é o produto mais comprado sozinho "
                        f"({fmt_num(top_solo['frequencia'])} pedidos). "
                        "Sugira um complemento natural no momento da compra.")

    #  CANDIDATOS A REMOÇÃO 
    with tabs[tab_idx["Candidatos a Remoção"]]:
        st.subheader("Candidatos a Remoção do Cardápio")
        st.caption("Produtos com menor receita e baixa frequência — avalie custo de produção antes de remover")

        col_rem, col_info = st.columns([3, 2])
        with col_rem:
            rem = df_remocao.copy()
            rem["receita"] = rem["receita"].apply(brl)
            rem.columns = ["Produto", "Frequência", "Receita"]
            st.dataframe(rem, use_container_width=True, hide_index=True, height=500)
        with col_info:
            st.markdown("""
            <div style="background:#FFF8E1;border-left:4px solid #F59E0B;padding:16px;border-radius:6px;margin-top:8px">
            <b> Como avaliar</b><br><br>
            Antes de remover um produto do cardápio, considere:<br><br>
            • <b>Custo de produção</b> — o produto pode ser simples de fazer<br>
            • <b>Perfil do cliente</b> — pode ser item de nicho fiel<br>
            • <b>Sazonalidade</b> — pode ter melhor desempenho em outro período<br>
            • <b>Complementaridade</b> — pode estar puxando venda de outro item<br><br>
            Remova apenas o que não cobre sequer seu custo variável.
            </div>
            """, unsafe_allow_html=True)

            top_rem = df_remocao.iloc[0] if not df_remocao.empty else None
            if top_rem is not None:
                st.info(f" **{top_rem['xProd']}** — produto com menor receita da lista. "
                        f"Apenas {fmt_num(top_rem['frequencia'])} vendas no período.")

    #  CURVA ABC
    with tabs[tab_idx["Curva ABC"]]:
        st.subheader("Curva ABC — Relevância dos Produtos")
        st.caption("Classifica todos os produtos pela participação acumulada na receita total.")

        if df_abc is not None and not df_abc.empty:
            _fat_abc = df_abc["Receita (R$)"].sum()
            _ga = df_abc[df_abc["Curva"] == "A"]
            _gb = df_abc[df_abc["Curva"] == "B"]
            _gc = df_abc[df_abc["Curva"] == "C"]

            ca, cb, cc = st.columns(3)
            with ca:
                st.markdown(
                    f"<div style='background:#D1FAE5;border-left:5px solid #059669;"
                    f"padding:14px 16px;border-radius:6px'>"
                    f"<div style='font-size:22px;font-weight:700;color:#065F46'>GRUPO A</div>"
                    f"<div style='font-size:13px;color:#065F46;margin-top:4px'>"
                    f"{len(_ga)} produtos · {_ga['Receita (R$)'].sum()/_fat_abc*100:.1f}% da receita</div>"
                    f"<div style='font-size:18px;font-weight:600;color:#059669;margin-top:6px'>"
                    f"{brl(_ga['Receita (R$)'].sum())}</div>"
                    f"<div style='font-size:11px;color:#6B7280;margin-top:2px'>Produtos vitais (0–80%)</div>"
                    f"</div>", unsafe_allow_html=True)
            with cb:
                st.markdown(
                    f"<div style='background:#FEF3C7;border-left:5px solid #D97706;"
                    f"padding:14px 16px;border-radius:6px'>"
                    f"<div style='font-size:22px;font-weight:700;color:#92400E'>GRUPO B</div>"
                    f"<div style='font-size:13px;color:#92400E;margin-top:4px'>"
                    f"{len(_gb)} produtos · {_gb['Receita (R$)'].sum()/_fat_abc*100:.1f}% da receita</div>"
                    f"<div style='font-size:18px;font-weight:600;color:#D97706;margin-top:6px'>"
                    f"{brl(_gb['Receita (R$)'].sum())}</div>"
                    f"<div style='font-size:11px;color:#6B7280;margin-top:2px'>Produtos importantes (80–95%)</div>"
                    f"</div>", unsafe_allow_html=True)
            with cc:
                st.markdown(
                    f"<div style='background:#F3F4F6;border-left:5px solid #6B7280;"
                    f"padding:14px 16px;border-radius:6px'>"
                    f"<div style='font-size:22px;font-weight:700;color:#374151'>GRUPO C</div>"
                    f"<div style='font-size:13px;color:#374151;margin-top:4px'>"
                    f"{len(_gc)} produtos · {_gc['Receita (R$)'].sum()/_fat_abc*100:.1f}% da receita</div>"
                    f"<div style='font-size:18px;font-weight:600;color:#6B7280;margin-top:6px'>"
                    f"{brl(_gc['Receita (R$)'].sum())}</div>"
                    f"<div style='font-size:11px;color:#6B7280;margin-top:2px'>Longa cauda (95–100%)</div>"
                    f"</div>", unsafe_allow_html=True)

            st.markdown("---")

            _filtro_abc = st.radio("Filtrar grupo:", ["Todos", "A", "B", "C"],
                                   horizontal=True, key="radio_abc")
            _df_show = df_abc if _filtro_abc == "Todos" else df_abc[df_abc["Curva"] == _filtro_abc]

            _df_display = _df_show.copy()
            _df_display["Receita (R$)"] = _df_display["Receita (R$)"].apply(brl)
            _df_display["% Receita"]    = _df_display["% Receita"].round(2).astype(str) + "%"
            _df_display["% Acumulado"] = _df_display["% Acumulado"].round(2).astype(str) + "%"

            st.dataframe(_df_display, use_container_width=True, hide_index=True, height=420)
        else:
            st.info("Rode a análise para ver a Curva ABC.")

    #  TICKET DRIVERS
    with tabs[tab_idx["Ticket Drivers"]]:
        st.subheader("Produtos que Influenciam o Valor do Pedido")

        col_e, col_r = st.columns(2)

        # Ticket médio geral como referência no cabeçalho
        if not df_elev.empty:
            _tm_geral = df_elev["Ticket Médio Geral"].iloc[0]
            st.caption(f"Ticket médio geral do período: **{brl(_tm_geral)}** — base de comparação para todos os produtos abaixo")

        with col_e:
            st.markdown("####  Elevam o Ticket Médio")
            st.caption("Quando presentes, o pedido tende a ser maior")
            if df_elev.empty:
                st.info("Dados insuficientes.")
            else:
                tbl = df_elev[["Produto", "Nº Pedidos",
                               "Ticket Médio c/ Produto", "Diferença R$", "Diferença %"]].copy()
                tbl["Ticket Médio c/ Produto"] = tbl["Ticket Médio c/ Produto"].apply(brl)
                tbl["Diferença R$"] = tbl["Diferença R$"].apply(lambda v: f"+{brl(v)}")
                tbl["Diferença %"]  = tbl["Diferença %"].apply(lambda v: f"+{v:.1f}%".replace(".", ","))
                st.dataframe(tbl, use_container_width=True, hide_index=True, height=400)
                top = df_elev.iloc[0]
                st.success(f"**{top['Produto']}** — pedidos com esse produto têm ticket "
                           f"{brl(top['Diferença R$'])} acima da média. Deixe-o visível e sugira sempre.")

        with col_r:
            st.markdown("####  Reduzem o Ticket Médio")
            st.caption("Quando presentes, o pedido tende a ser menor — avaliar perfil")
            if df_redu.empty:
                st.info("Dados insuficientes.")
            else:
                tbl2 = df_redu[["Produto", "Nº Pedidos",
                                "Ticket Médio c/ Produto", "Diferença R$", "Diferença %"]].copy()
                tbl2["Ticket Médio c/ Produto"] = tbl2["Ticket Médio c/ Produto"].apply(brl)
                tbl2["Diferença R$"] = tbl2["Diferença R$"].apply(lambda v: brl(v) if v < 0 else f"+{brl(v)}")
                tbl2["Diferença %"]  = tbl2["Diferença %"].apply(lambda v: f"{v:.1f}%".replace(".", ","))
                # Filtra só os que realmente reduzem (diferença negativa)
                tbl2_show = tbl2[df_redu["Diferença R$"] < 0]
                if tbl2_show.empty:
                    st.info("Nenhum produto com impacto negativo significativo no ticket.")
                else:
                    st.dataframe(tbl2_show, use_container_width=True, hide_index=True, height=400)
                    bot = df_redu[df_redu["Diferença R$"] < 0].iloc[0]
                    st.info(f"ℹ **{bot['Produto']}** — pedidos com esse produto têm ticket menor. "
                            "Pode ser perfil de cliente diferente ou produto de compra rápida/solo.")

    #  TEMPORAL 
    if "Temporal" in tab_idx:
        with tabs[tab_idx["Temporal"]]:
            st.subheader("Análise Temporal de Vendas")

            subtabs = st.tabs(["Fluxo por Horário", "Por Turno", "Dia da Semana", "Horários com Potencial"])

            # ── SUBTAB 1: FLUXO POR HORÁRIO ──────────────────
            with subtabs[0]:
                if df_por_hora is None:
                    st.warning("Dados de horário não disponíveis — o arquivo precisa ter coluna dhEmi.")
                else:
                    st.markdown("#### Volume de Vendas por Hora do Dia")
                    st.caption("Cada barra representa o total de pedidos registrados naquele horário ao longo do mês")

                    CORES_TURNO = {"Manhã": "#2563EB", "Tarde": "#E67E22", "Noite": "#1A234E"}

                    fig_hora = px.bar(
                        df_por_hora.sort_values("hora"),
                        x="hora", y="transacoes",
                        color="turno",
                        color_discrete_map=CORES_TURNO,
                        text="transacoes",
                        labels={"hora": "Hora do dia", "transacoes": "Nº de Pedidos", "turno": "Turno"},
                        category_orders={"turno": ["Manhã", "Tarde", "Noite"]},
                    )
                    fig_hora.update_traces(texttemplate="%{text}", textposition="outside")
                    fig_hora.update_xaxes(
                        tickvals=list(range(0, 24)),
                        ticktext=[f"{h:02d}h" for h in range(0, 24)],
                    )
                    fig_hora.update_layout(
                        height=420,
                        bargap=0.15,
                        legend_title_text="Turno",
                        xaxis_title="Hora do dia",
                        yaxis_title="Nº de Pedidos",
                    )
                    st.plotly_chart(fig_hora, use_container_width=True)

                    # Métricas-resumo
                    hora_pico  = df_por_hora.loc[df_por_hora["transacoes"].idxmax()]
                    hora_fraca = df_por_hora.loc[df_por_hora["transacoes"].idxmin()]
                    turno_lider = df_por_turno.loc[df_por_turno["transacoes"].idxmax()] if df_por_turno is not None else None

                    c1, c2, c3, c4 = st.columns(4)
                    c1.metric("Hora de pico",
                              f"{int(hora_pico['hora']):02d}h",
                              f"{fmt_num(hora_pico['transacoes'])} pedidos")
                    c2.metric("Hora mais fraca",
                              f"{int(hora_fraca['hora']):02d}h",
                              f"{fmt_num(hora_fraca['transacoes'])} pedidos")
                    if turno_lider is not None:
                        pct_br = f"{turno_lider['pct']:.1f}".replace(".", ",") + "%"
                        c3.metric("Turno principal", turno_lider["turno"], f"{pct_br} dos pedidos")
                        c4.metric("Ticket médio no turno", brl(turno_lider["ticket_medio"]))

                    st.divider()

                    # Cards por turno
                    st.markdown("**Resumo por turno**")
                    if df_por_turno is not None and not df_por_turno.empty:
                        desc_turno = {
                            "Manhã":  "Das 05h às 11h — café, pães e salgados dominam",
                            "Tarde":  "Das 12h às 17h — lanche e refeição rápida",
                            "Noite":  "Das 18h em diante — finalização do dia",
                        }
                        cols_turno = st.columns(len(df_por_turno))
                        for col_c, (_, row) in zip(cols_turno, df_por_turno.iterrows()):
                            cor = CORES_TURNO.get(row["turno"], "#555")
                            pct_s = f"{row['pct']:.1f}".replace(".", ",") + "%"
                            col_c.markdown(
                                f"""<div style="background:{cor};border-radius:10px;padding:16px;color:white;text-align:center">
                                <div style="font-size:18px;font-weight:700">{row['turno']}</div>
                                <div style="font-size:28px;font-weight:900;margin:6px 0">{pct_s}</div>
                                <div style="font-size:13px;opacity:.85">dos pedidos</div>
                                <hr style="border-color:rgba(255,255,255,.3);margin:8px 0">
                                <div style="font-size:13px">{fmt_num(row['transacoes'])} pedidos</div>
                                <div style="font-size:13px">{brl(row['receita'])} receita</div>
                                <div style="font-size:12px;opacity:.8;margin-top:6px">{desc_turno.get(row['turno'],'')}</div>
                                </div>""",
                                unsafe_allow_html=True,
                            )

                    # Interpretação automática
                    st.divider()
                    if turno_lider is not None:
                        pct_lider = f"{turno_lider['pct']:.1f}".replace(".", ",") + "%"
                        st.info(
                            f"**Leitura para o gestor:** O turno da **{turno_lider['turno']}** concentra "
                            f"{pct_lider} de todos os pedidos do mês. A hora de maior movimento é as "
                            f"**{int(hora_pico['hora']):02d}h** com {fmt_num(hora_pico['transacoes'])} pedidos. "
                            f"O horário das **{int(hora_fraca['hora']):02d}h** é o mais fraco — avalie se vale "
                            f"manter operação plena ou usar esse período para promoções direcionadas."
                        )

            # ── SUBTAB 2: POR TURNO (produtos) ───────────────
            with subtabs[1]:
                st.markdown("#### Produtos Mais Vendidos por Turno")
                if df_turno is None:
                    st.warning("Dados de horário não disponíveis no arquivo.")
                else:
                    turnos = ["Manhã", "Tarde", "Noite"]
                    cols_t = st.columns(3)
                    for col_t, turno in zip(cols_t, turnos):
                        sub = df_turno[df_turno["turno"] == turno].head(8)
                        col_t.markdown(f"**{turno}**")
                        if sub.empty:
                            col_t.info("Sem dados")
                        else:
                            tbl_t = sub[["xProd", "frequencia", "receita"]].copy()
                            tbl_t["receita"] = tbl_t["receita"].apply(brl)
                            tbl_t.columns = ["Produto", "Freq.", "Receita"]
                            col_t.dataframe(tbl_t, use_container_width=True, hide_index=True, height=340)

            # ── SUBTAB 3: DIA DA SEMANA ───────────────────────
            with subtabs[2]:
                st.markdown("#### Produtos por Dia da Semana")
                if df_dia_tipo is None:
                    st.warning("Dados de data não disponíveis no arquivo.")
                else:
                    col_u, col_f = st.columns(2)
                    for col_v, tipo in [(col_u, "Dia Útil"), (col_f, "Final de Semana")]:
                        sub = df_dia_tipo[df_dia_tipo["tipo_dia"] == tipo].head(10)
                        col_v.markdown(f"**{tipo}**")
                        if sub.empty:
                            col_v.info("Sem dados")
                        else:
                            tbl_d = sub[["xProd", "frequencia", "receita"]].copy()
                            tbl_d["receita"] = tbl_d["receita"].apply(brl)
                            tbl_d.columns = ["Produto", "Freq.", "Receita"]
                            col_v.dataframe(tbl_d, use_container_width=True, hide_index=True, height=400)

            # ── SUBTAB 4: HORÁRIOS COM POTENCIAL ─────────────
            with subtabs[3]:
                st.markdown("#### Horários com Potencial Inexplorado")
                st.caption("Horários com volume abaixo da média — oportunidade de ação comercial")
                if df_horas is None:
                    st.warning("Dados de horário não disponíveis.")
                else:
                    fig_h = px.bar(
                        df_horas, x="hora", y="notas", color="Status",
                        color_discrete_map={
                            "Ativo": "#27AE60",
                            "Abaixo da média": "#F39C12",
                            "Potencial inexplorado": "#E74C3C",
                        },
                        labels={"hora": "Hora", "notas": "Nº Pedidos", "Status": ""},
                    )
                    fig_h.update_xaxes(tickvals=list(range(0,24)),
                                       ticktext=[f"{h:02d}h" for h in range(0,24)])
                    fig_h.update_layout(height=400)
                    st.plotly_chart(fig_h, use_container_width=True)


    #  SIMULAÇÕES 
    with tabs[tab_idx["Simulações"]]:
        st.subheader("Simulações de Crescimento")

        subtab_sim = st.tabs(["Crescimento de Receita", "Simulação de Preços", "Combos"])

        with subtab_sim[0]:
            st.markdown("#### Simulação de Crescimento de Receita")
            sim = df_sim_rec.copy()
            sim["Impacto Mensal"] = sim["Impacto Mensal"].apply(brl)
            sim["Impacto Anual"]  = sim["Impacto Anual"].apply(brl)
            st.dataframe(sim, use_container_width=True, hide_index=True)

            melhor = df_sim_rec.nlargest(1, "Impacto Mensal").iloc[0]
            st.success(f"Maior oportunidade: **{melhor['Estratégia']}**"
                       f" — {brl(melhor['Impacto Mensal'])} por mês"
                       f" — {brl(melhor['Impacto Anual'])} por ano")

        with subtab_sim[1]:
            st.markdown("#### Simulação de Ajuste de Preços")
            st.caption("Simula reajuste de 10%, 15% e 20% com queda estimada de 5% no volume")
            sim_p = df_sim_preco.copy()
            for col in ["receita", "+10% preço (-5% vol)", "+15% preço (-5% vol)", "+20% preço (-5% vol)"]:
                sim_p[col] = sim_p[col].apply(brl)
            for col in ["Δ +10%", "Δ +15%", "Δ +20%"]:
                sim_p[col] = df_sim_preco[col].apply(lambda v: f"+{brl(v)}" if v > 0 else brl(v))
            st.dataframe(
                sim_p[["xProd", "preco_medio", "receita",
                        "+10% preço (-5% vol)", "Δ +10%",
                        "+15% preço (-5% vol)", "Δ +15%",
                        "+20% preço (-5% vol)", "Δ +20%"]].rename(columns={
                    "xProd": "Produto",
                    "preco_medio": "Preço Médio Atual", "receita": "Receita Atual",
                }),
                use_container_width=True, hide_index=True, height=500,
            )

        with subtab_sim[2]:
            st.markdown("#### Precificação de Combos")
            if df_combos.empty:
                st.info("Sem dados suficientes de preço para calcular combos.")
            else:
                tbl_c = df_combos.copy()
                for col in ["Preço A", "Preço B", "Total Individual", "Combo c/ 5% desc.", "Combo c/ 10% desc."]:
                    tbl_c[col] = tbl_c[col].apply(brl)
                st.dataframe(tbl_c, use_container_width=True, hide_index=True)
                st.info("O combo com 5% de desconto mantém margem saudável e aumenta percepção de valor.")

    #  METAS 
    with tabs[tab_idx["Metas"]]:
        st.subheader("Plano de Metas Mensais por Produto")
        st.caption("Baseado nos top produtos por receita — metas de +10% e +20%")

        tbl_m = df_metas.copy()
        for col in ["receita", "Meta +10%", "Meta +20%"]:
            tbl_m[col] = tbl_m[col].apply(brl)
        tbl_m["volume"] = tbl_m["volume"].apply(lambda v: f"{v:.1f}")

        st.dataframe(
            tbl_m.rename(columns={
                "xProd": "Produto",
                "volume": "Volume Atual", "receita": "Receita Atual",
                "frequencia": "Freq. Pedidos",
            }),
            use_container_width=True, hide_index=True, height=450,
        )

        fat_total = df_metas["receita"].sum() if "receita" in df_metas.columns else 0
        # recalcular com valores originais
        fat_top10 = df.groupby("xProd")["vProd"].sum().nlargest(10).sum()
        st.metric("Receita atual (Top 10 produtos)", brl(fat_top10))
        c1, c2 = st.columns(2)
        c1.metric("Meta +10%", brl(fat_top10 * 1.10), f"+{brl(fat_top10 * 0.10)}")
        c2.metric("Meta +20%", brl(fat_top10 * 1.20), f"+{brl(fat_top10 * 0.20)}")

    #  CESTA — adiciona solo produtos 
    # (já existe a aba Cesta, vamos adicionar solo dentro dela via subtab)

    #  NF-e B2B 
    if "NF-e (B2B)" in tab_idx and not df_nfe.empty:
        with tabs[tab_idx["NF-e (B2B)"]]:
            st.subheader("NF-e — Vendas para Empresas (B2B)")
            st.caption("Notas Fiscais de saída emitidas para CNPJ — distribuidor, fornecedor, revenda")

            if "vNF" in df_nfe.columns and "chave" in df_nfe.columns:
                notas_b2b = df_nfe.drop_duplicates("chave")
                fat_b2b   = notas_b2b["vNF"].sum()
                n_b2b     = len(notas_b2b)
                c1, c2, c3, c4 = st.columns(4)
                c1.metric("Faturamento B2B",   brl(fat_b2b))
                c2.metric("Notas Emitidas",    fmt_num(n_b2b))
                c3.metric("Ticket Médio B2B",  brl(fat_b2b / n_b2b if n_b2b else 0))
                c4.metric("% do Total",        (f"{fat_b2b/kpis['faturamento']*100:.1f}".replace(".", ",") + "%") if kpis["faturamento"] else "—")

            if "xProd" in df_nfe.columns and "vProd" in df_nfe.columns:
                import re as _re
                def _limpar_xprod(s: str) -> str:
                    s = _re.sub(r"\s+Ped\.+\s*:.*$", "", s, flags=_re.IGNORECASE)
                    s = _re.sub(r"\s+Nro\.?\s*Item\s*:.*$", "", s, flags=_re.IGNORECASE)
                    return s.strip()

                _df_nfe_view = df_nfe.copy()
                _df_nfe_view["xProd_clean"] = _df_nfe_view["xProd"].astype(str).apply(_limpar_xprod)
                _dest_col = "destinatario" if "destinatario" in _df_nfe_view.columns else None

                _grp_cols = ["xProd_clean"] + ([_dest_col] if _dest_col else [])
                top_b2b = (
                    _df_nfe_view.groupby(_grp_cols)
                    .agg(receita=("vProd", "sum"), notas=("chave", "nunique"))
                    .sort_values("receita", ascending=False).head(15)
                    .reset_index()
                )
                top_b2b["Receita"] = top_b2b["receita"].apply(brl)

                _rename = {"xProd_clean": "Produto", "notas": "Notas"}
                if _dest_col:
                    _rename[_dest_col] = "Empresa"
                _cols_show = ["Produto"] + (["Empresa"] if _dest_col else []) + ["Notas", "Receita"]
                top_b2b = top_b2b.rename(columns=_rename)

                col_b1, col_b2 = st.columns([3, 2])
                with col_b1:
                    st.markdown("#### Top Produtos (NF-e)")
                    st.dataframe(top_b2b[_cols_show],
                                 use_container_width=True, hide_index=True, height=450)

                with col_b2:
                    st.markdown("#### Receita por Produto (NF-e)")
                    import plotly.express as _px
                    _chart_df = top_b2b.head(10).iloc[::-1].copy()
                    _chart_df["_receita_num"] = _df_nfe_view.groupby(
                        _df_nfe_view["xProd"].astype(str).apply(_limpar_xprod))["vProd"].sum().nlargest(10).iloc[::-1].values
                    _fig_nfe = _px.bar(
                        _chart_df, x="_receita_num", y="Produto",
                        orientation="h",
                        labels={"_receita_num": "Receita (R$)", "Produto": ""},
                        color_discrete_sequence=["#2563EB"],
                    )
                    _fig_nfe.update_layout(
                        margin=dict(l=0, r=10, t=10, b=10),
                        xaxis_tickformat=",.0f",
                        plot_bgcolor="white",
                        height=420,
                    )
                    _fig_nfe.update_xaxes(showgrid=True, gridcolor="#f0f0f0")
                    _fig_nfe.update_yaxes(tickfont=dict(size=11))
                    st.plotly_chart(_fig_nfe, use_container_width=True)

    #  SIMPLES NACIONAL
    if "Simples Nacional" in tab_idx and sn_result is not None:
        with tabs[tab_idx["Simples Nacional"]]:
            st.subheader("Simples Nacional — Verificação dos 80%")
            st.caption(
                "A legislação do Simples Nacional exige que as compras para comercialização "
                "não ultrapassem **80% do faturamento total**. Aqui você confere se a empresa está dentro do limite."
            )

            _sn_total   = sn_result["total_compras_comercializacao"]
            _sn_pct     = sn_result["pct_faturamento"]
            _sn_status  = sn_result["status"]
            _sn_fat     = kpis["faturamento"]

            # ── KPI Cards ────────────────────────────────────────────
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("Faturamento Total",           brl(_sn_fat))
            c2.metric("Compras de Comercialização",  brl(_sn_total))
            c3.metric("% do Faturamento",
                      f"{_sn_pct:.1f}%".replace(".", ","),
                      delta=f"Limite: 80%",
                      delta_color="off")
            _status_emoji = {"OK": "✅ OK", "ALERTA": "⚠️ Atenção", "EXCEDIDO": "🔴 Excedido", "SEM_DADOS": "—"}
            c4.metric("Status",  _status_emoji.get(_sn_status, _sn_status))

            # ── Barra visual de progresso ─────────────────────────────
            _bar_color = "#16a34a" if _sn_pct <= 70 else ("#f59e0b" if _sn_pct <= 80 else "#dc2626")
            _pct_capped = min(_sn_pct, 100)
            st.markdown(
                f"""
<div style="margin:12px 0 20px;background:#f1f5f9;border-radius:10px;height:28px;
            position:relative;overflow:hidden;">
  <div style="background:{_bar_color};width:{_pct_capped:.1f}%;height:100%;
              border-radius:10px;display:flex;align-items:center;justify-content:flex-end;
              padding-right:10px;font-size:13px;font-weight:700;color:white;">
    {_sn_pct:.1f}%
  </div>
  <div style="position:absolute;top:0;left:80%;width:2px;height:100%;background:#1e3a5f;opacity:0.5;"></div>
  <div style="position:absolute;top:0;left:80%;transform:translateX(-50%);margin-top:30px;
              font-size:11px;color:#6b7280;white-space:nowrap;">limite 80%</div>
</div>
""",
                unsafe_allow_html=True,
            )

            if _sn_status == "EXCEDIDO":
                st.error(
                    f"🔴 **ALERTA FISCAL**: as compras de comercialização representam **{_sn_pct:.1f}%** do faturamento — "
                    f"acima do limite de 80% exigido pelo Simples Nacional. "
                    f"Isso pode indicar subfaturamento ou excesso de estoques. Consulte o contador."
                )
            elif _sn_status == "ALERTA":
                st.warning(
                    f"⚠️ **Atenção**: as compras estão em **{_sn_pct:.1f}%** do faturamento — "
                    f"próximo do limite de 80%. Acompanhe de perto para não ultrapassar."
                )
            else:
                st.success(
                    f"✅ **Dentro do limite**: as compras de comercialização representam **{_sn_pct:.1f}%** "
                    f"do faturamento — abaixo dos 80% exigidos."
                )

            st.divider()

            # ── Breakdown por CFOP ────────────────────────────────────
            df_por_cfop = sn_result["df_por_cfop"]
            df_por_forn = sn_result["df_por_fornecedor"]

            col_cfop, col_forn = st.columns(2)

            with col_cfop:
                st.markdown("#### Compras por CFOP")
                if not df_por_cfop.empty:
                    _cfop_show = df_por_cfop.copy()
                    _cfop_show["total_compras"] = _cfop_show["total_compras"].apply(brl)
                    _cfop_show = _cfop_show.rename(columns={
                        "CFOP": "CFOP", "total_compras": "Total Compras",
                        "notas": "Notas", "itens": "Itens",
                    })
                    st.dataframe(_cfop_show, use_container_width=True, hide_index=True, height=320)
                else:
                    st.info("Nenhum CFOP de comercialização encontrado nas entradas.")

            with col_forn:
                st.markdown("#### Top Fornecedores")
                if not df_por_forn.empty:
                    _forn_show = df_por_forn.copy()
                    _forn_show["total_compras"] = _forn_show["total_compras"].apply(brl)
                    _forn_show = _forn_show.rename(columns={
                        "emitente": "Fornecedor", "total_compras": "Total Compras", "notas": "Notas",
                    })
                    st.dataframe(_forn_show, use_container_width=True, hide_index=True, height=320)
                else:
                    st.info("Sem dados de fornecedores.")

            # ── CFOPs considerados ────────────────────────────────────
            with st.expander("ℹ️ CFOPs considerados como compras de comercialização"):
                st.markdown(
                    "Os seguintes CFOPs são contabilizados como **compras de comercialização** "
                    "(revenda de mercadorias):\n\n"
                    + "\n".join(f"- **{c}**" for c in sorted(CFOPS_COMERCIALIZACAO))
                    + "\n\n_Ficam de fora: imobilizado (14xx), uso e consumo (15xx), industrialização (11xx)._"
                )

    #  EXPORTAÇÃO
    # CSS de impressão injetado na página principal (não no iframe)
    st.markdown("""
<style>
@media print {
  /* Oculta elementos de navegação e controle */
  section[data-testid="stSidebar"],
  [data-testid="stHeader"],
  [data-testid="stToolbar"],
  [data-testid="stDecoration"],
  [data-testid="stStatusWidget"],
  .stDeployButton,
  footer,
  #exportar-relatorio,
  #exportar-relatorio ~ div { display: none !important; }

  /* Remove margens do Streamlit para aproveitar a folha */
  .main .block-container {
    padding: 0 !important;
    max-width: 100% !important;
  }

  /* Força impressão colorida */
  * {
    -webkit-print-color-adjust: exact !important;
    print-color-adjust: exact !important;
  }

  /* Evita quebra de página no meio de gráficos e tabelas */
  .stPlotlyChart, .stDataFrame, [data-testid="stTable"] {
    page-break-inside: avoid;
  }
}
</style>
""", unsafe_allow_html=True)

    st.divider()
    st.subheader("Exportar Relatório")

    _nome_base = f"Analise_de_Vendas_{cli_label.replace(' ', '_')}_{per_label.replace(' ', '_')}"

    # ── Cache PPTX e Excel no session_state para evitar re-geração a cada clique ──
    # Usa o fingerprint da análise como chave: se o dado mudou, regenera; caso
    # contrário reutiliza o bytes já gerado, evitando o rerun "pesado" que
    # causa o reinício da tela.
    _fp_atual = st.session_state.get("_analise_fp", "")

    if st.session_state.get("_export_fp") != _fp_atual:
        # Dados novos — gera os arquivos e guarda no cache
        _pptx_bytes = None
        try:
            from pptx import Presentation  # noqa: F401
            import matplotlib  # noqa: F401
            _pptx_bytes = exportar_pptx(
                kpis, df_pares, df_trios,
                df_cesta, df_turno, df_bcg,
                df_elev, df_redu, df_sim_rec, df_sim_preco,
                df_combos, df_metas, df_horas,
                df_solo, df_remocao, df_dia_tipo,
                df_nfe, kpis_nfce,
                fonte_label, cli_label, per_label,
                df_abc=df_abc,
                df_por_hora=df_por_hora,
                df_por_turno=df_por_turno,
                sn_result=sn_result,
            )
        except ImportError:
            pass

        _xlsx_bytes = exportar_excel(kpis, df_pares, df_trios,
                                     df_cesta, df_bcg, df_abc, df_remocao,
                                     df_elev, df_redu, df_sim_preco, df_sim_rec,
                                     df_combos, df_metas,
                                     cli_label, per_label,
                                     sn_result=sn_result)

        st.session_state["_export_pptx"]  = _pptx_bytes
        st.session_state["_export_xlsx"]  = _xlsx_bytes
        st.session_state["_export_fp"]    = _fp_atual
    else:
        # Mesmos dados — reutiliza cache (clique no botão não regenera nada)
        _pptx_bytes = st.session_state.get("_export_pptx")
        _xlsx_bytes = st.session_state.get("_export_xlsx")

    col_xl, col_pp, col_pdf = st.columns(3)

    with col_xl:
        st.download_button(
            label="Baixar Excel (.xlsx)",
            data=_xlsx_bytes,
            file_name=f"{_nome_base}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
        )

    with col_pp:
        if _pptx_bytes:
            st.download_button(
                label="Baixar PowerPoint (.pptx)",
                data=_pptx_bytes,
                file_name=f"{_nome_base}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        else:
            st.warning("Instale: `pip install python-pptx matplotlib`")

    with col_pdf:
        if _pptx_bytes:
            _pdf_bytes = pptx_para_pdf(_pptx_bytes)
            if _pdf_bytes:
                st.download_button(
                    label="Baixar PDF (.pdf)",
                    data=_pdf_bytes,
                    file_name=f"{_nome_base}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                )
            else:
                # LibreOffice não disponível — fallback: botão de impressão
                import streamlit.components.v1 as _components
                _components.html(
                    """<style>
  body{margin:0;padding:0}
  .bp{display:block;width:100%;padding:5px 16px;background:white;color:#4f46e5;
      border:1px solid #d1d5db;border-radius:6px;font-size:14px;font-weight:500;
      cursor:pointer;font-family:sans-serif;text-align:center;box-sizing:border-box}
  .bp:hover{background:#f5f3ff;border-color:#4f46e5}
</style>
<button class="bp" onclick="window.parent.print()">🖨️ Imprimir / Salvar PDF</button>""",
                    height=40,
                )
        else:
            st.info("Gere o PPTX primeiro para habilitar o PDF.")


if __name__ == "__main__":
    main()
